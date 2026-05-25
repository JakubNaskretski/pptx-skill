"""v5 redesign — structural-skeleton digest (self-contained).

This module is intentionally a single file so v5 is removable as one
delete + a tiny hook-removal in cli.py if the redesign doesn't pan
out. Until phase F flips the build flag, v5 outputs are purely
additive: v4 slide.yaml / theme.yaml continue to be written under
workspace/decks/<deck>/ untouched, and v5 writes alongside under
workspace/themes/<deck>/ and workspace/skeletons/<deck>_<NN>/.

See REDESIGN.md (root) for the architecture; phase callouts in
function docstrings reference sub-phases B1-B5 + C1.
"""

from __future__ import annotations

import hashlib
import shutil
from pathlib import Path

import yaml
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER


# ---------------------------------------------------------------------------
# Local helpers — duplicated from cli.py to keep this module
# import-cycle-free and easy to delete as a unit.
# ---------------------------------------------------------------------------


def _write_yaml(path: Path, data: dict) -> None:
    path.write_text(
        yaml.safe_dump(data, sort_keys=False, allow_unicode=True, width=100),
        encoding="utf-8",
    )


def _aspect_ratio(w: int, h: int) -> str:
    if w <= 0 or h <= 0:
        return "free"
    r = w / h
    candidates = {"16:9": 16 / 9, "4:3": 4 / 3, "1:1": 1.0, "3:4": 3 / 4, "9:16": 9 / 16}
    best, best_err = "free", 0.10
    for name, target in candidates.items():
        err = abs(r - target) / target
        if err < best_err:
            best, best_err = name, err
    return best


def _resolve_palette_v5(theme_v4: dict) -> dict:
    """Map v4's raw clrScheme slots → v5 semantic palette roles.

    Used by both digest_theme (B1) and digest_skeleton (B2 — for
    color_role resolution on slot styles). Centralised so the two
    stay aligned.
    """
    palette_v4 = theme_v4.get("palette") or {}
    aliases = theme_v4.get("aliases") or {}

    def _r(role: str) -> str:
        slot = aliases.get(role, "")
        return palette_v4.get(slot, "") if slot else ""

    return {
        "primary": _r("primary"),
        "accent": _r("accent"),
        "text_default": _r("text"),
        "background": _r("background"),
    }


# ---------------------------------------------------------------------------
# B1 — Theme extraction
# ---------------------------------------------------------------------------


def digest_theme(
    prs,
    original_path: Path,
    deck_stem: str,
    theme_v4: dict,
    themes_root: Path,
) -> dict:
    """Write workspace/themes/<deck>/theme.yaml + master.pptx.

    Reads from the v4 theme dict (palette + fonts + aliases from
    cli.extract_deck_theme) and emits the v5 schema: semantic palette
    roles, fonts, master_pptx + preview references, plus the
    auto-classified decorations array from B3 (master shapes tagged
    as top_bar / bottom_bar / corner_logo / page_number / section_panel
    / background_image via geometric heuristics).

    master.pptx is currently a copy of original.pptx. Proper master-only
    extraction (drop slides, keep masters + layouts + theme + media)
    is a future refinement — the phase-E build engine can strip slides
    at build time, so this is correctness-equivalent.
    """
    theme_dir = themes_root / deck_stem
    theme_dir.mkdir(parents=True, exist_ok=True)
    _extract_master_pptx(original_path, theme_dir / "master.pptx")

    palette_v5 = {k: v for k, v in _resolve_palette_v5(theme_v4).items() if v}
    fonts = {k: v for k, v in (theme_v4.get("fonts") or {}).items() if v}
    decorations = _classify_decorations(prs, palette_v5)

    out = {
        "id": deck_stem,
        "palette": palette_v5,
        "fonts": fonts,
        "master_pptx": "master.pptx",
        "preview": "preview.png",
        "decorations": decorations,
    }
    _write_yaml(theme_dir / "theme.yaml", out)
    return out


def _extract_master_pptx(src: Path, dst: Path) -> None:
    """Copy the original deck as master.pptx.

    Placeholder for proper master-only extraction; the build engine in
    phase E can handle slide stripping at build time so this is
    correctness-equivalent, just larger than the minimal artifact.
    """
    shutil.copyfile(src, dst)


# ---------------------------------------------------------------------------
# B3 — Decoration auto-classification
# ---------------------------------------------------------------------------


def _classify_decorations(prs, palette_v5: dict) -> list[dict]:
    """Walk slide masters; classify each shape per geometric heuristics
    into top_bar / bottom_bar / corner_logo / page_number / section_panel
    / background_image. Returns a list of decoration entries with kind
    + fractional geometry + optional color_role.

    Decorations are informational only in v5 (they ride along with the
    master at build time, not addressable from plans). Captured for
    future v5.x decoration-mixer; user verifies in the Flask UI.
    """
    slide_w = prs.slide_width or 9144000
    slide_h = prs.slide_height or 6858000

    out: list[dict] = []
    for master in prs.slide_masters:
        for shape in list(master.shapes):
            deco = _classify_one_decoration(shape, slide_w, slide_h, palette_v5)
            if deco is not None:
                out.append(deco)
    return out


def _classify_one_decoration(shape, slide_w: int, slide_h: int, palette_v5: dict) -> dict | None:
    geom = _fractional_geometry(shape, slide_w, slide_h)
    x, y, w, h = geom["x"], geom["y"], geom["w"], geom["h"]
    if w <= 0 or h <= 0:
        return None

    kind = _infer_decoration_kind(shape, x, y, w, h)
    if kind is None:
        return None

    entry: dict = {"kind": kind, "geometry": geom}
    color_role = _shape_fill_role(shape, palette_v5)
    if color_role:
        entry["color_role"] = color_role
    return entry


def _infer_decoration_kind(shape, x: float, y: float, w: float, h: float) -> str | None:
    # Slide-number placeholder wins regardless of geometry — masters
    # typically place it in a corner but custom templates may differ.
    if _is_page_number(shape):
        return "page_number"

    # Full-width thin bars: top_bar / bottom_bar
    if w > 0.9 and h < 0.02:
        if y < 0.03:
            return "top_bar"
        if y + h > 0.97:
            return "bottom_bar"

    # Small corner artifact: logo if Picture, decoration otherwise
    if w < 0.10 and h < 0.10:
        near_left = x < 0.05
        near_right = (x + w) > 0.95
        near_top = y < 0.05
        near_bottom = (y + h) > 0.95
        if (near_left or near_right) and (near_top or near_bottom):
            if _shape_is_picture(shape):
                return "corner_logo"
            return "corner_decoration"

    # Background image: covers nearly whole slide
    if w > 0.9 and h > 0.9 and _shape_is_picture(shape):
        return "background_image"

    # Section panel: large solid fill (>30% × >30%) likely a divider
    # background. Lower priority than the more specific kinds above.
    if w > 0.30 and h > 0.30 and _shape_has_solid_fill(shape):
        return "section_panel"

    return None


def _is_page_number(shape) -> bool:
    try:
        if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.SLIDE_NUMBER:
            return True
    except (AttributeError, ValueError):
        pass
    # PPT stores the page-number field as <a:fld type="slidenum">; the
    # surface text on extracted shape.text often comes through as ‹#›
    # (per locale). Check both shapes commonly seen in real decks.
    try:
        text = (shape.text_frame.text or "")
        if "‹#›" in text or "<#>" in text or "#" == text.strip():
            return True
    except (AttributeError, ValueError):
        pass
    return False


def _shape_has_solid_fill(shape) -> bool:
    try:
        rgb = shape.fill.fore_color.rgb
        return rgb is not None
    except (AttributeError, ValueError, TypeError):
        return False


def _shape_fill_role(shape, palette_v5: dict) -> str | None:
    try:
        rgb = shape.fill.fore_color.rgb
        if rgb is None:
            return None
        return _resolve_color_role(f"#{str(rgb).upper()}", palette_v5)
    except (AttributeError, ValueError, TypeError):
        return None


# ---------------------------------------------------------------------------
# B2 — Slot inventory + skeleton digest
# ---------------------------------------------------------------------------


# Placeholder types that produce textual slots (heading / paragraph /
# bullets). Mirrors cli.PLACEHOLDER_TEXTUAL but locally defined so v5
# can drift without coupling.
_PLACEHOLDER_TEXTUAL = frozenset({
    PP_PLACEHOLDER.TITLE,
    PP_PLACEHOLDER.CENTER_TITLE,
    PP_PLACEHOLDER.SUBTITLE,
    PP_PLACEHOLDER.BODY,
    PP_PLACEHOLDER.OBJECT,
})

# Placeholder types treated as footer-kind slots (page number, date,
# footer text). Build engine binds page-number placeholders to host
# master auto-numbering; other footers are plain text.
_PLACEHOLDER_FOOTER = frozenset({
    PP_PLACEHOLDER.FOOTER,
    PP_PLACEHOLDER.DATE,
    PP_PLACEHOLDER.SLIDE_NUMBER,
})


def compute_repeated_picture_info(prs, threshold: float = 0.5) -> dict[str, float]:
    """Return {sha: median_area_fraction} for pictures appearing on
    >= `threshold` fraction of slides — i.e. deck-wide brand decoration.

    The median area is the *typical* rendered size of the picture across
    the deck. Callers compare a specific instance's area against this:
    if a brand mark is rendered noticeably larger on one slide (e.g. a
    cover/title slide that features the logo at hero size), that
    instance is treated as a content slot rather than decoration.

    SHA1 of the image binary collapses visually-identical pictures
    inserted as separate shape instances. A picture appearing multiple
    times on the SAME slide counts once toward the deck-wide repeat
    fraction but all instances contribute their areas to the median.

    Call once per deck at the start of ingest; pass the result into
    digest_skeleton via `repeated_picture_info`.
    """
    sha_to_areas: dict[str, list[float]] = {}
    sha_to_slide_count: dict[str, int] = {}
    n_slides = len(prs.slides)
    if n_slides == 0:
        return {}

    sw = prs.slide_width or 9144000
    sh = prs.slide_height or 6858000
    slide_area = sw * sh

    for slide in prs.slides:
        seen_on_slide: set[str] = set()
        for shape in slide.shapes:
            if getattr(shape, "is_placeholder", False):
                continue
            if not _shape_is_picture(shape):
                continue
            sha = _picture_sha(shape)
            if sha is None:
                continue
            area = ((shape.width or 0) * (shape.height or 0)) / slide_area if slide_area else 0
            sha_to_areas.setdefault(sha, []).append(area)
            if sha not in seen_on_slide:
                seen_on_slide.add(sha)
                sha_to_slide_count[sha] = sha_to_slide_count.get(sha, 0) + 1

    out: dict[str, float] = {}
    for sha, slide_count in sha_to_slide_count.items():
        if slide_count / n_slides < threshold:
            continue
        areas = sorted(sha_to_areas[sha])
        m = len(areas)
        median = areas[m // 2] if m % 2 == 1 else (areas[m // 2 - 1] + areas[m // 2]) / 2
        out[sha] = median
    return out


# Backwards-compat name during this session — alias for callers that
# may have imported the old name. Remove once cli.py is updated.
def compute_repeated_picture_shas(prs, threshold: float = 0.5) -> set[str]:
    return set(compute_repeated_picture_info(prs, threshold).keys())


def _picture_sha(shape) -> str | None:
    try:
        return hashlib.sha1(shape.image.blob).hexdigest()
    except (AttributeError, ValueError):
        return None


_FEATURED_SIZE_MULTIPLIER = 2.0  # >= 2x median area = treat as featured slot
                                 # even if SHA is repeated deck-wide


def digest_skeleton(
    slide,
    slide_w: int,
    slide_h: int,
    deck_stem: str,
    slide_number: int,
    theme_v4: dict,
    skeletons_root: Path,
    v4_preview_path: Path | None = None,
    repeated_picture_info: dict[str, float] | None = None,
) -> dict:
    """Write workspace/skeletons/<deck>_<NN>/skeleton.yaml.

    Detects slots with full v5 schema (kind, fractional geometry, style
    block with font_role / color_role / size_pt / alignment / bold,
    per-kind constraints with max_chars / max_items / max_rows /
    max_cols / required). Status is `pending`; categories list empty
    (B5 populates), background_image null (B4 may set).

    `repeated_picture_info`: {sha: median_area_fraction} for pictures
    appearing on >= 50% of slides — caller obtains via
    compute_repeated_picture_info(prs) once per deck. A picture whose
    SHA is in this map is skipped UNLESS its area on this specific
    slide is >= 2x the median (then it's "featured" content for this
    slide — e.g. the brand logo at hero size on the cover slide).
    Passing None disables the filter.

    If v4_preview_path is supplied and exists, copies the file into
    the skeleton dir so the skeleton dir is self-contained for the
    Flask UI in C1.
    """
    skeleton_id = f"{deck_stem}_{slide_number:02d}"
    skeleton_dir = skeletons_root / skeleton_id
    skeleton_dir.mkdir(parents=True, exist_ok=True)

    palette_v5 = _resolve_palette_v5(theme_v4)
    theme_fonts = theme_v4.get("fonts") or {}
    palette_v4 = theme_v4.get("palette") or {}

    fresh_slots = _extract_slots(
        slide, slide_w, slide_h, palette_v4, palette_v5, theme_fonts,
        repeated_picture_info=repeated_picture_info or {},
    )

    existing = _read_existing_skeleton(skeleton_dir / "skeleton.yaml")
    # C-actions: if a slot the user promoted exists in skeleton.yaml
    # but isn't in fresh_slots (heuristic doesn't claim it), AND its
    # shape_id still exists in the source slide, preserve it verbatim.
    slots = _merge_user_promoted_slots(fresh_slots, existing, slide)
    consumed_shape_ids = {s["shape_id"] for s in slots if "shape_id" in s}

    # Capture everything else as unmapped_shapes so the user can
    # promote anything our heuristic missed via the /v5 UI.
    unmapped_shapes = _extract_unmapped_shapes(
        slide, slide_w, slide_h, consumed_shape_ids, repeated_picture_info or {},
    )

    overlap_info = _detect_overlap_candidates(slide, slide_w, slide_h, slots)

    # Preserve existing descriptive fields on re-ingest so user
    # categorisation isn't blown away. Status is preserved too — if
    # the user has already rejected, re-ingest doesn't auto-recategorise.
    preserved_status = existing.get("status") if existing else None
    preserved_categories = existing.get("categories") if existing else None
    preserved_overlap_decision = existing.get("overlap_decision") if existing else None

    # B5: propose categories from slot inventory only if user hasn't
    # set their own. User-set categories survive re-ingest verbatim.
    categories = (
        preserved_categories
        if preserved_categories
        else _propose_categories(slots)
    )

    out = {
        "id": skeleton_id,
        "source_deck": deck_stem,
        "source_slide_index": slide_number,
        "status": preserved_status or "pending",
        "categories": categories,
        "preview": "preview.png",
        "background_image": None,
        "slots": slots,
        "unmapped_shapes": unmapped_shapes,
    }
    if preserved_overlap_decision:
        out["overlap_decision"] = preserved_overlap_decision
    if overlap_info:
        out["digest_warnings"] = overlap_info["warnings"]
        out["overlap_candidates"] = overlap_info["candidates"]
    _write_yaml(skeleton_dir / "skeleton.yaml", out)

    if v4_preview_path and v4_preview_path.exists():
        shutil.copyfile(v4_preview_path, skeleton_dir / "preview.png")

    return out


def _merge_user_promoted_slots(fresh_slots: list[dict], existing: dict | None, slide) -> list[dict]:
    """Re-ingest preservation for user-curated slots.

    Two cases:
    1. **user-promoted**: an existing slot whose shape_id is in the
       source slide but not in fresh_slots (heuristic skipped it; the
       user promoted it via C-actions). Preserve verbatim.
    2. **user-edited**: an existing slot whose shape_id matches a
       fresh slot AND carries `user_edited: true` (kind was changed
       in the UI). Replace the fresh slot with the user's verbatim
       version — heuristic's kind/constraints lose to user intent.

    Without preservation, a re-ingest would silently revert both
    kinds of user curation.
    """
    if not existing:
        return list(fresh_slots)
    existing_slots = existing.get("slots") or []
    fresh_by_shape_id = {s["shape_id"]: s for s in fresh_slots if "shape_id" in s}
    source_shape_ids = {sh.shape_id for sh in slide.shapes}

    out: list[dict] = []
    consumed_shape_ids: set[int] = set()
    for fresh in fresh_slots:
        sid = fresh.get("shape_id")
        existing_match = next(
            (e for e in existing_slots if e.get("shape_id") == sid),
            None,
        )
        if existing_match and existing_match.get("user_edited"):
            # User overrode the heuristic's kind — keep the user version
            out.append(existing_match)
        else:
            out.append(fresh)
        if sid is not None:
            consumed_shape_ids.add(sid)

    # Append user-promoted slots (in existing but not in fresh)
    for e in existing_slots:
        sid = e.get("shape_id")
        if sid is None or sid in consumed_shape_ids:
            continue
        if sid not in source_shape_ids:
            continue
        out.append(e)

    return out


def _read_existing_skeleton(path: Path) -> dict | None:
    if not path.exists():
        return None
    try:
        return yaml.safe_load(path.read_text(encoding="utf-8")) or {}
    except Exception:
        return None


# ---------------------------------------------------------------------------
# Slot extraction
# ---------------------------------------------------------------------------


def _extract_slots(
    slide, slide_w: int, slide_h: int,
    palette_v4: dict, palette_v5: dict, theme_fonts: dict,
    repeated_picture_info: dict[str, float],
) -> list[dict]:
    slots: list[dict] = []
    used_ids: set[str] = set()
    consumed_shape_ids: set[int] = set()

    # First pass — placeholders. Ordered by reading-order convention
    # (title first, then body / hero / footer).
    for shape in list(slide.placeholders):
        slot = _slot_from_placeholder(
            shape, slide_w, slide_h, palette_v4, palette_v5, theme_fonts, used_ids,
        )
        if slot is not None:
            slot["shape_id"] = shape.shape_id
            slots.append(slot)
            consumed_shape_ids.add(shape.shape_id)

    # Second pass — free shapes (non-placeholder pictures, tables,
    # charts) that the agent might want as slots.
    for shape in list(slide.shapes):
        if shape.shape_id in consumed_shape_ids:
            continue
        if getattr(shape, "is_placeholder", False):
            continue
        slot = _slot_from_free_shape(
            shape, slide_w, slide_h, palette_v4, palette_v5, theme_fonts, used_ids,
            repeated_picture_info=repeated_picture_info,
        )
        if slot is not None:
            slot["shape_id"] = shape.shape_id
            slots.append(slot)

    return slots


# ---------------------------------------------------------------------------
# C-actions: unmapped-shapes capture
# ---------------------------------------------------------------------------


def _extract_unmapped_shapes(
    slide, slide_w: int, slide_h: int,
    consumed_shape_ids: set[int],
    repeated_picture_info: dict[str, float],
) -> list[dict]:
    """Walk every shape on the slide; emit an entry for each one the
    heuristic skipped (i.e. shape_id NOT in consumed_shape_ids). Lets
    the C-actions UI surface what was filtered so the user can promote
    real content the heuristic missed.

    Each entry carries shape_id (for re-ingest matching and the
    promote-shape endpoint), kind_hint (what the source shape is),
    fractional geometry, source_excerpt (first 80 chars of text if
    any), and skipped_reason (human-readable explanation).
    """
    out: list[dict] = []
    for shape in slide.shapes:
        if shape.shape_id in consumed_shape_ids:
            continue
        entry: dict = {
            "shape_id": shape.shape_id,
            "kind_hint": _shape_kind_hint(shape),
            "geometry": _fractional_geometry(shape, slide_w, slide_h),
            "skipped_reason": _infer_skipped_reason(shape, slide_w, slide_h, repeated_picture_info),
        }
        excerpt = _shape_text_excerpt(shape)
        if excerpt:
            entry["source_excerpt"] = excerpt
        out.append(entry)
    return out


def _shape_kind_hint(shape) -> str:
    """Human-friendly tag for shape.shape_type — used in the UI to
    say 'this is what was here'.
    """
    try:
        st = shape.shape_type
    except (AttributeError, ValueError):
        return "unknown"
    mapping = {
        MSO_SHAPE_TYPE.AUTO_SHAPE: "auto_shape",
        MSO_SHAPE_TYPE.PICTURE: "picture",
        MSO_SHAPE_TYPE.TEXT_BOX: "text_box",
        MSO_SHAPE_TYPE.FREEFORM: "freeform",
        MSO_SHAPE_TYPE.GROUP: "group",
        MSO_SHAPE_TYPE.LINE: "line",
        MSO_SHAPE_TYPE.CHART: "chart",
        MSO_SHAPE_TYPE.TABLE: "table",
        MSO_SHAPE_TYPE.PLACEHOLDER: "placeholder",
    }
    return mapping.get(st, "other")


def _shape_text_excerpt(shape) -> str:
    try:
        if not getattr(shape, "has_text_frame", False):
            return ""
        text = (shape.text_frame.text or "").strip()
        return _truncate(text, 80) if text else ""
    except (AttributeError, ValueError):
        return ""


def _infer_skipped_reason(shape, slide_w: int, slide_h: int,
                          repeated_picture_info: dict[str, float]) -> str:
    """Best-effort retroactive explanation for why a shape was skipped.
    Mirrors the filters in _slot_from_placeholder + _slot_from_free_shape.
    """
    if getattr(shape, "is_placeholder", False):
        try:
            ph_type = shape.placeholder_format.type
        except (AttributeError, ValueError):
            ph_type = None
        if ph_type is not None and ph_type not in _PLACEHOLDER_TEXTUAL \
                and ph_type not in _PLACEHOLDER_FOOTER \
                and ph_type != PP_PLACEHOLDER.PICTURE:
            return f"placeholder kind {_safe_enum_short_name(ph_type)} not yet supported"

    if _shape_is_picture(shape):
        area = _area_fraction(shape, slide_w, slide_h)
        if area < 0.005:
            return "sub-pixel area"
        if repeated_picture_info:
            sha = _picture_sha(shape)
            if sha and sha in repeated_picture_info:
                median = repeated_picture_info[sha]
                if median > 0 and area / median < _FEATURED_SIZE_MULTIPLIER:
                    return f"repeated brand mark ({len(repeated_picture_info)} decks-wide pictures detected)"
        return "decoration"

    try:
        st = shape.shape_type
    except (AttributeError, ValueError):
        return "unknown shape type"

    if st in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM):
        if _shape_has_meaningful_text(shape):
            return "auto-shape with text (would be captured if free text box)"
        return "conscious-drop: freeform / auto-shape"
    if st == MSO_SHAPE_TYPE.GROUP:
        return "conscious-drop: group"
    if st == MSO_SHAPE_TYPE.LINE:
        return "conscious-drop: line / connector"
    if st == MSO_SHAPE_TYPE.TEXT_BOX and not _shape_has_meaningful_text(shape):
        return "empty text box"
    return "skipped (no matching slot rule)"


def _slot_from_placeholder(
    shape, slide_w: int, slide_h: int,
    palette_v4: dict, palette_v5: dict, theme_fonts: dict,
    used_ids: set[str],
) -> dict | None:
    ph = shape.placeholder_format
    ph_type = ph.type

    # Tables in OBJECT placeholders — check before textual fallback.
    if getattr(shape, "has_table", False):
        return _slot_table(shape, slide_w, slide_h, used_ids)

    # Charts in OBJECT placeholders.
    if getattr(shape, "has_chart", False):
        return _slot_chart(shape, slide_w, slide_h, used_ids)

    if ph_type == PP_PLACEHOLDER.PICTURE:
        return _slot_image(shape, slide_w, slide_h, used_ids, slot_id_base="hero")

    if ph_type in _PLACEHOLDER_FOOTER:
        return _slot_footer(
            shape, slide_w, slide_h, palette_v4, palette_v5, theme_fonts, used_ids, ph_type,
        )

    if ph_type in _PLACEHOLDER_TEXTUAL:
        return _slot_text_or_bullets(
            shape, slide_w, slide_h, palette_v4, palette_v5, theme_fonts, used_ids, ph_type,
        )

    # Unsupported placeholder kind (e.g. media, header). Skip — the
    # build engine doesn't know how to fill these.
    return None


def _slot_from_free_shape(
    shape, slide_w: int, slide_h: int,
    palette_v4: dict, palette_v5: dict, theme_fonts: dict,
    used_ids: set[str],
    repeated_picture_info: dict[str, float] | None = None,
) -> dict | None:
    if getattr(shape, "has_table", False):
        return _slot_table(shape, slide_w, slide_h, used_ids)
    if getattr(shape, "has_chart", False):
        return _slot_chart(shape, slide_w, slide_h, used_ids)
    if _shape_is_picture(shape):
        # Deck-wide brand marks are usually decoration — skip them.
        # BUT if this instance is rendered much larger than typical
        # (e.g. a logo featured at hero size on a cover slide), the
        # author has used the asset as featured content; capture it.
        if repeated_picture_info:
            sha = _picture_sha(shape)
            if sha and sha in repeated_picture_info:
                median_area = repeated_picture_info[sha]
                this_area = _area_fraction(shape, slide_w, slide_h)
                if median_area > 0 and this_area / median_area < _FEATURED_SIZE_MULTIPLIER:
                    return None  # typical size → decoration
                # else fall through — featured instance, treat as slot
        if _is_image_slot_worthy(shape, slide_w, slide_h):
            return _slot_image(shape, slide_w, slide_h, used_ids, slot_id_base="hero")
    # Free text boxes (TEXT_BOX shape type, or any non-placeholder with
    # a non-empty text_frame). Designer-dropped captions, body
    # paragraphs, annotations.
    if _shape_has_meaningful_text(shape):
        return _slot_from_free_text(
            shape, slide_w, slide_h, palette_v4, palette_v5, theme_fonts, used_ids,
        )
    # Freeforms / auto-shapes / connectors / groups — conscious drop
    # per REDESIGN.md. B4 may rescue some as frozen background.
    return None


def _shape_has_meaningful_text(shape) -> bool:
    """True if shape has a text_frame with non-empty text. Skips
    shapes that have a frame but no content (decorative auto-shapes
    that happen to expose has_text_frame).
    """
    try:
        if not getattr(shape, "has_text_frame", False):
            return False
        return bool((shape.text_frame.text or "").strip())
    except (AttributeError, ValueError):
        return False


def _slot_from_free_text(
    shape, slide_w: int, slide_h: int,
    palette_v4: dict, palette_v5: dict, theme_fonts: dict,
    used_ids: set[str],
) -> dict:
    """Build a slot from a non-placeholder text shape (TEXT_BOX or
    auto-shape with text). Mirrors the placeholder-text logic but
    infers kind from content shape since there's no PH type to
    consult.
    """
    tf = shape.text_frame
    text = tf.text or ""
    paras = list(getattr(tf, "paragraphs", []) or [])
    non_empty = [p for p in paras if (p.text or "").strip()]

    # Pull first-run size to inform kind inference: large + short =
    # heading; otherwise paragraph or bullets per paragraph count.
    size_pt = 0.0
    try:
        runs = list(getattr(paras[0], "runs", []) or [])
        if runs and runs[0].font.size is not None:
            size_pt = float(runs[0].font.size.pt)
    except (AttributeError, ValueError, IndexError):
        pass

    if len(non_empty) > 1:
        kind = "bullets"
        base_id = "body"
    elif len(text) <= 80 and size_pt >= 24:
        kind = "heading"
        base_id = "heading"
    else:
        kind = "paragraph"
        base_id = "body"

    # Style — pass ph_type=None; _resolve_font_role falls back to
    # minor for non-title placeholders, which is right for free text.
    style = _extract_style(shape, None, palette_v4, palette_v5, theme_fonts)
    geometry = _fractional_geometry(shape, slide_w, slide_h)
    required = _is_required(shape, slide_w, slide_h, True)

    if kind == "bullets":
        item_lens = [len((p.text or "").strip()) for p in non_empty]
        longest = max(item_lens) if item_lens else 0
        constraints = {
            "max_items": max(1, len(non_empty)),
            "max_chars_per_item": max(20, int(longest * 1.5)) if longest else 80,
            "required": required,
        }
        items = [(p.text or "").strip() for p in non_empty[:3]]
        excerpt = _truncate(" / ".join(f"• {it}" for it in items), 80)
    else:
        constraints = {
            "max_chars": max(20, int(len(text) * 1.5)) if text else 60,
            "max_lines": max(1, len(non_empty) or 1),
            "required": required,
        }
        excerpt = _truncate(text.strip(), 80) if text else ""

    return {
        "id": _unique(base_id, used_ids),
        "kind": kind,
        "geometry": geometry,
        "style": style,
        "constraints": constraints,
        "source_excerpt": excerpt,
    }


# ---------------------------------------------------------------------------
# B5 — Auto-classifier (category proposal)
#
# Heuristic rules over the slot inventory only; no NLP, no content
# semantics beyond simple substring matches. The user confirms in the
# Flask UI (C2, future) so wrong proposals are easy to fix. Multiple
# categories per slide are valid — a data slide with a quote tag is
# both. Fallback to "content" if nothing else fires.
# ---------------------------------------------------------------------------


# Categories enum mirrors REDESIGN.md.
_CATEGORIES = (
    "opening", "section_divider", "content", "comparison",
    "data", "metric", "quote", "closing",
)


def _propose_categories(slots: list[dict]) -> list[str]:
    if not slots:
        return ["content"]

    kinds = [s.get("kind") for s in slots]
    excerpts = [s.get("source_excerpt", "") or "" for s in slots]

    proposed: list[str] = []

    # data — has a structured data slot
    if "table" in kinds or "chart" in kinds:
        proposed.append("data")

    # opening — small slot count, dominant heading, no body content
    heading_slots = [s for s in slots if s.get("kind") == "heading"]
    body_slots = [s for s in slots if s.get("kind") in ("bullets", "paragraph")]
    if (
        len(slots) <= 3
        and heading_slots
        and not body_slots
        and any((h.get("geometry") or {}).get("h", 0) > 0.10 for h in heading_slots)
    ):
        proposed.append("opening")

    # comparison — 2 slots of same kind, mirrored horizontally
    if _has_side_by_side_pair(slots):
        proposed.append("comparison")

    # quote — any source excerpt contains quote glyphs
    if any(_looks_like_quote(t) for t in excerpts):
        proposed.append("quote")

    # closing — any source excerpt matches farewell / Q&A patterns
    if any(_looks_like_closing(t) for t in excerpts):
        proposed.append("closing")

    # Fallback so the user always sees at least one suggestion.
    if not proposed:
        proposed.append("content")

    return proposed


def _has_side_by_side_pair(slots: list[dict]) -> bool:
    """Two slots of the same kind that mirror horizontally — same y/h,
    different x, similar w. Catches 2-column compare/contrast slides.
    """
    eligible = [s for s in slots if s.get("kind") in ("bullets", "paragraph", "image")]
    for i in range(len(eligible)):
        for j in range(i + 1, len(eligible)):
            a, b = eligible[i], eligible[j]
            if a.get("kind") != b.get("kind"):
                continue
            ga = a.get("geometry") or {}
            gb = b.get("geometry") or {}
            if (
                abs(ga.get("y", 0) - gb.get("y", 0)) < 0.05
                and abs(ga.get("h", 0) - gb.get("h", 0)) < 0.10
                and abs(ga.get("w", 0) - gb.get("w", 0)) < 0.10
                and abs(ga.get("x", 0) - gb.get("x", 0)) > 0.15
            ):
                return True
    return False


# Quote-character pairs. "Quote" fires only if a paired span is
# substantial (>=30 chars) or makes up half the excerpt — bare quotes
# on a single noun (company name, "Suplemencik") shouldn't fire.
_QUOTE_PAIRS = (
    ('"', '"'),
    ("“", "”"),
    ("„", "”"),   # German / Polish low-9 + high-9
    ("„", '"'),   # Polish low-9 + straight (common in extracted text)
    ("‟", "”"),
    ("«", "»"),
)
_CLOSING_PATTERNS = (
    # English
    "thank you", "thanks", "questions?", "any questions",
    "q&a", "q & a", "next steps", "contact us", "get in touch",
    # Polish (test deck Naskrętski is Polish-language; add more as
    # other-language decks land)
    "dziękuję", "dziekuje", "pytania", "kontakt",
)


def _looks_like_quote(text: str) -> bool:
    if not text or len(text) < 20:
        return False
    longest_span = 0
    for open_c, close_c in _QUOTE_PAIRS:
        start = text.find(open_c)
        if start < 0:
            continue
        end = text.find(close_c, start + len(open_c))
        if end < 0:
            continue
        span = end - start - len(open_c)
        if span > longest_span:
            longest_span = span
    return longest_span >= 30 or (
        longest_span >= 15 and longest_span / max(1, len(text)) >= 0.5
    )


def _looks_like_closing(text: str) -> bool:
    if not text:
        return False
    lower = text.lower()
    return any(p in lower for p in _CLOSING_PATTERNS)


# ---------------------------------------------------------------------------
# B4-detect — Overlap detection (no rendering yet)
#
# EXPERIMENTAL per REDESIGN.md. Detects only — does not modify the
# slot inventory or render a background.png. Flags slides where a
# large picture sits under text shapes, indicating a structural
# illustration (chain diagram, process flow) that the agent shouldn't
# blindly swap. User reviews flagged slides in C1 and decides via
# C2 controls (future) whether to freeze-as-background, override
# (treat as image slot), or reject.
#
# Kept in the same module as the rest of v5 ingest so removal is one
# delete, not a refactor. To disable: comment out the call site in
# digest_skeleton; nothing else depends on overlap_candidates yet.
# ---------------------------------------------------------------------------


def _detect_overlap_candidates(slide, slide_w: int, slide_h: int, slots: list[dict]) -> dict:
    """Find pictures whose bbox overlaps text slots — candidate frozen
    backgrounds. Returns {} when nothing flagged so the field is omitted
    from skeleton.yaml on clean slides.

    Heuristic: a picture is a candidate if (a) it covers ≥15% of slide
    area AND (b) ≥50% of at least one text slot's bbox area sits inside
    the picture's bbox. The 50%/15% thresholds are from REDESIGN.md
    and intentionally conservative — false positives are worse than
    missed detections for the user-review flow.
    """
    if not slots:
        return {}

    text_kinds = {"heading", "paragraph", "bullets", "footer"}
    text_slot_bboxes = []
    for s in slots:
        if s.get("kind") not in text_kinds:
            continue
        g = s.get("geometry") or {}
        x1, y1 = g.get("x", 0.0), g.get("y", 0.0)
        x2, y2 = x1 + g.get("w", 0.0), y1 + g.get("h", 0.0)
        text_slot_bboxes.append((s["id"], x1, y1, x2, y2))

    if not text_slot_bboxes:
        return {}

    candidates: list[dict] = []
    for shape in list(slide.shapes):
        if not _shape_is_picture(shape):
            continue
        if getattr(shape, "is_placeholder", False):
            continue
        if _area_fraction(shape, slide_w, slide_h) < 0.15:
            continue

        g = _fractional_geometry(shape, slide_w, slide_h)
        px1, py1 = g["x"], g["y"]
        px2, py2 = px1 + g["w"], py1 + g["h"]

        overlapping = []
        for slot_id, tx1, ty1, tx2, ty2 in text_slot_bboxes:
            ix1, iy1 = max(px1, tx1), max(py1, ty1)
            ix2, iy2 = min(px2, tx2), min(py2, ty2)
            if ix2 <= ix1 or iy2 <= iy1:
                continue
            intersect_area = (ix2 - ix1) * (iy2 - iy1)
            text_area = (tx2 - tx1) * (ty2 - ty1)
            if text_area > 0 and intersect_area / text_area >= 0.50:
                overlapping.append(slot_id)

        if overlapping:
            candidates.append({
                "picture_shape_id": shape.shape_id,
                "picture_geometry": g,
                "overlapping_slot_ids": overlapping,
            })

    if not candidates:
        return {}
    return {"warnings": ["overlap_detected"], "candidates": candidates}


# ---------------------------------------------------------------------------
# Per-kind slot builders
# ---------------------------------------------------------------------------


def _slot_text_or_bullets(
    shape, slide_w: int, slide_h: int,
    palette_v4: dict, palette_v5: dict, theme_fonts: dict,
    used_ids: set[str], ph_type,
) -> dict | None:
    tf = getattr(shape, "text_frame", None)
    if tf is None:
        return None

    paras = list(getattr(tf, "paragraphs", []) or [])
    text = tf.text or ""

    non_empty = [p for p in paras if (p.text or "").strip()]
    # SUBTITLE is heading-like when short (tagline, deck subtitle) but
    # paragraph-like when long (intro / about-this-deck body). Threshold
    # at 80 chars — below is heading, above is paragraph.
    is_heading = (
        ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
        or (ph_type == PP_PLACEHOLDER.SUBTITLE and len(text) <= 80)
    )
    is_bulleted = (
        not is_heading
        and ph_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.SUBTITLE)
        and len(non_empty) > 1
    )

    base_id = _placeholder_base_id(ph_type)
    slot_id = _unique(base_id, used_ids)
    style = _extract_style(shape, ph_type, palette_v4, palette_v5, theme_fonts)
    geometry = _fractional_geometry(shape, slide_w, slide_h)
    required = _is_required(shape, slide_w, slide_h, is_heading or ph_type == PP_PLACEHOLDER.BODY)

    if is_bulleted:
        item_lens = [len((p.text or "").strip()) for p in non_empty]
        longest = max(item_lens) if item_lens else 0
        constraints = {
            "max_items": max(1, len(non_empty)),
            "max_chars_per_item": max(20, int(longest * 1.5)) if longest else 80,
            "required": required,
        }
        kind = "bullets"
        items = [(p.text or "").strip() for p in non_empty[:3]]
        excerpt = _truncate(" / ".join(f"• {it}" for it in items), 80)
    else:
        constraints = {
            "max_chars": max(20, int(len(text) * 1.5)) if text else 60,
            "max_lines": max(1, len(non_empty) or 1),
            "required": required,
        }
        kind = "heading" if is_heading else "paragraph"
        excerpt = _truncate(text.strip(), 80) if text else ""

    return {
        "id": slot_id,
        "kind": kind,
        "geometry": geometry,
        "style": style,
        "constraints": constraints,
        "source_excerpt": excerpt,
    }


def _slot_footer(
    shape, slide_w: int, slide_h: int,
    palette_v4: dict, palette_v5: dict, theme_fonts: dict,
    used_ids: set[str], ph_type,
) -> dict:
    base = {
        PP_PLACEHOLDER.FOOTER: "footer",
        PP_PLACEHOLDER.DATE: "date",
        PP_PLACEHOLDER.SLIDE_NUMBER: "page_number",
    }.get(ph_type, "footer")
    slot_id = _unique(base, used_ids)
    text = ""
    try:
        text = (shape.text_frame.text or "")
    except (AttributeError, ValueError):
        pass
    excerpt = (
        "auto page number" if ph_type == PP_PLACEHOLDER.SLIDE_NUMBER
        else _truncate(text.strip(), 80) if text else ""
    )
    return {
        "id": slot_id,
        "kind": "footer",
        "geometry": _fractional_geometry(shape, slide_w, slide_h),
        "style": _extract_style(shape, ph_type, palette_v4, palette_v5, theme_fonts),
        "constraints": {
            "max_chars": max(20, int(len(text) * 1.5)) if text else 40,
            "max_lines": 1,
            "required": False,
            "auto_from_host": (ph_type == PP_PLACEHOLDER.SLIDE_NUMBER),
        },
        "source_excerpt": excerpt,
    }


def _slot_image(
    shape, slide_w: int, slide_h: int, used_ids: set[str], slot_id_base: str = "hero",
) -> dict:
    aspect = _aspect_ratio(shape.width or 0, shape.height or 0)
    return {
        "id": _unique(slot_id_base, used_ids),
        "kind": "image",
        "geometry": _fractional_geometry(shape, slide_w, slide_h),
        "constraints": {
            "aspect": aspect,
            "required": _is_required(shape, slide_w, slide_h, True),
            "auto_fit": "cover",
        },
        "source_excerpt": f"image ({aspect})",
    }


def _slot_table(shape, slide_w: int, slide_h: int, used_ids: set[str]) -> dict:
    rows = cols = 0
    has_header = False
    try:
        rows = len(list(shape.table.rows))
        cols = len(list(shape.table.columns))
        has_header = bool(getattr(shape.table, "first_row", False))
    except (AttributeError, ValueError):
        pass
    excerpt = f"{max(1, rows)}×{max(1, cols)} table"
    if has_header:
        excerpt += " (header row)"
    return {
        "id": _unique("data_table", used_ids),
        "kind": "table",
        "geometry": _fractional_geometry(shape, slide_w, slide_h),
        "constraints": {
            "max_rows": max(1, rows),
            "max_cols": max(1, cols),
            "has_header": has_header,
            "required": _is_required(shape, slide_w, slide_h, True),
        },
        "source_excerpt": excerpt,
    }


def _slot_chart(shape, slide_w: int, slide_h: int, used_ids: set[str]) -> dict:
    chart_type = "unknown"
    n_series = 0
    n_categories = 0
    try:
        ct = shape.chart.chart_type
        chart_type = _safe_enum_short_name(ct)
        plots = list(shape.chart.plots)
        if plots:
            n_series = len(list(plots[0].series))
            n_categories = len(list(plots[0].categories))
    except (AttributeError, ValueError):
        pass
    return {
        "id": _unique("data_chart", used_ids),
        "kind": "chart",
        "geometry": _fractional_geometry(shape, slide_w, slide_h),
        "constraints": {
            "chart_type": chart_type,
            "max_series": max(1, n_series),
            "max_categories": max(1, n_categories),
            "required": _is_required(shape, slide_w, slide_h, True),
        },
        "source_excerpt": f"{chart_type} chart, {max(1, n_series)} series × {max(1, n_categories)} cats",
    }


def _truncate(s: str, n: int) -> str:
    if not s:
        return ""
    s = " ".join(s.split())  # collapse newlines/extra spaces for excerpt display
    if len(s) <= n:
        return s
    return s[: n - 1].rstrip() + "…"


# ---------------------------------------------------------------------------
# Style + geometry + role resolution
# ---------------------------------------------------------------------------


def _extract_style(
    shape, ph_type, palette_v4: dict, palette_v5: dict, theme_fonts: dict,
) -> dict:
    """First-run style snapshot for the slot.

    Resolves font.name → font_role: major | minor | explicit (with the
    raw typeface preserved on explicit). Resolves font.color → color
    hex + color_role mapped to the v5 semantic palette. Missing fields
    omitted rather than nulled so the YAML stays readable.
    """
    out: dict = {}
    try:
        tf = shape.text_frame
    except (AttributeError, ValueError):
        return out
    paras = list(getattr(tf, "paragraphs", []) or [])
    if not paras:
        return out
    runs = list(getattr(paras[0], "runs", []) or [])
    para0 = paras[0]
    font = runs[0].font if runs else None

    # Font role
    raw_font_name = None
    if font is not None:
        try:
            raw_font_name = font.name
        except (AttributeError, ValueError):
            raw_font_name = None
    role = _resolve_font_role(raw_font_name, theme_fonts, ph_type)
    if role == "explicit":
        out["font_role"] = "explicit"
        out["typeface"] = raw_font_name
    else:
        out["font_role"] = role

    # Size + weight
    if font is not None:
        try:
            if font.size is not None:
                out["size_pt"] = float(font.size.pt)
        except (AttributeError, ValueError):
            pass
        try:
            if font.bold is not None:
                out["bold"] = bool(font.bold)
        except (AttributeError, ValueError):
            pass
        try:
            if font.italic is not None:
                out["italic"] = bool(font.italic)
        except (AttributeError, ValueError):
            pass

    # Alignment
    try:
        align = para0.alignment
        if align is not None:
            out["alignment"] = _safe_enum_short_name(align)
    except (AttributeError, ValueError):
        pass

    # Color
    if font is not None:
        color_hex = _resolve_run_color_hex(font, palette_v4)
        if color_hex:
            color_role = _resolve_color_role(color_hex, palette_v5)
            if color_role:
                out["color_role"] = color_role
            else:
                out["color"] = color_hex

    return out


def _resolve_font_role(font_name: str | None, theme_fonts: dict, ph_type) -> str:
    """Map a typeface against the theme's major/minor; fall back to
    a sensible default by placeholder kind when the font isn't set.

    - Explicit match against theme.fonts.major → "major"
    - Explicit match against theme.fonts.minor → "minor"
    - No font set (inherits from theme) → default by placeholder kind:
        titles inherit major, everything else inherits minor.
    - Explicit but matches neither → "explicit" (caller preserves the
      raw typeface so build doesn't drift).
    """
    major = (theme_fonts.get("major") or "").strip().lower()
    minor = (theme_fonts.get("minor") or "").strip().lower()
    name = (font_name or "").strip().lower()

    if not name:
        if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE,
                       PP_PLACEHOLDER.SUBTITLE):
            return "major"
        return "minor"
    if major and name == major:
        return "major"
    if minor and name == minor:
        return "minor"
    return "explicit"


def _resolve_color_role(color_hex: str, palette_v5: dict) -> str | None:
    """Match a hex against the semantic palette with a small tolerance.

    Returns the role name (primary / accent / text_default / background)
    if matched, else None — caller falls through to explicit hex.
    """
    target = (color_hex or "").upper().lstrip("#")
    if not target or len(target) != 6:
        return None
    for role, role_hex in palette_v5.items():
        if not role_hex:
            continue
        role_target = role_hex.upper().lstrip("#")
        if _hex_close(target, role_target, tol=5):
            return role
    return None


def _hex_close(a: str, b: str, tol: int = 5) -> bool:
    try:
        return all(abs(int(a[i:i+2], 16) - int(b[i:i+2], 16)) <= tol for i in (0, 2, 4))
    except (ValueError, IndexError):
        return False


def _resolve_run_color_hex(font, palette_v4: dict) -> str:
    """Resolve a run.font.color to a #RRGGBB hex via either direct rgb
    or theme-color lookup against the v4 palette (raw clrScheme).
    """
    try:
        color = font.color
    except (AttributeError, ValueError):
        return ""
    if color is None:
        return ""
    try:
        rgb = color.rgb
        if rgb is not None:
            return f"#{str(rgb).upper()}"
    except (AttributeError, ValueError):
        pass
    try:
        theme_color = color.theme_color
        if theme_color is not None:
            slot = _theme_color_to_slot(theme_color)
            if slot and slot in palette_v4:
                return palette_v4[slot]
    except (AttributeError, ValueError):
        pass
    return ""


def _theme_color_to_slot(theme_color) -> str:
    """Map a python-pptx MSO_THEME_COLOR enum value to a clrScheme slot
    name (matches the slot names v4 stores in palette).
    """
    name = _safe_enum_short_name(theme_color).lower()
    # MSO_THEME_COLOR names: TEXT_1 / TEXT_2 / BACKGROUND_1 / BACKGROUND_2
    # / ACCENT_1..6 / HYPERLINK / FOLLOWED_HYPERLINK. Map to clrScheme.
    mapping = {
        "text_1": "dk1", "text_2": "dk2",
        "background_1": "lt1", "background_2": "lt2",
        "accent_1": "accent1", "accent_2": "accent2", "accent_3": "accent3",
        "accent_4": "accent4", "accent_5": "accent5", "accent_6": "accent6",
        "hyperlink": "hlink", "followed_hyperlink": "folHlink",
    }
    return mapping.get(name, "")


def _fractional_geometry(shape, slide_w: int, slide_h: int) -> dict:
    w = shape.width or 0
    h = shape.height or 0
    left = shape.left or 0
    top = shape.top or 0
    if slide_w <= 0 or slide_h <= 0:
        return {"x": 0.0, "y": 0.0, "w": 0.0, "h": 0.0}
    return {
        "x": round(left / slide_w, 4),
        "y": round(top / slide_h, 4),
        "w": round(w / slide_w, 4),
        "h": round(h / slide_h, 4),
    }


def _area_fraction(shape, slide_w: int, slide_h: int) -> float:
    if slide_w <= 0 or slide_h <= 0:
        return 0.0
    w = shape.width or 0
    h = shape.height or 0
    return (w * h) / (slide_w * slide_h)


def _is_required(shape, slide_w: int, slide_h: int, hint_required: bool) -> bool:
    """Heuristic: slots covering > 10% of the slide are required (titles,
    hero, body); smaller are optional (footers, captions). ``hint_required``
    biases toward True for slot kinds that are typically load-bearing
    (titles, body, hero, table, chart) even when slightly smaller.
    """
    area = _area_fraction(shape, slide_w, slide_h)
    if hint_required and area > 0.05:
        return True
    return area > 0.10


def _shape_is_picture(shape) -> bool:
    try:
        return shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    except (AttributeError, ValueError):
        return False


def _is_image_slot_worthy(shape, slide_w: int, slide_h: int) -> bool:
    """Per-shape filter; very permissive after the per-deck repeat
    detection took over the "is this a brand mark" job.

    Only excludes sub-0.5% area pictures — i.e. tiny inline icons used
    as bullet markers or list separators where promoting them to slots
    would be noise. Aspect filter dropped (banners at extreme aspect
    are legitimate slots when they're not deck-wide brand decoration).

    The real decoration filter lives at compute_repeated_picture_shas;
    THIS check is just a sub-pixel safety floor.
    """
    return _area_fraction(shape, slide_w, slide_h) >= 0.005


# ---------------------------------------------------------------------------
# Id minting + small utilities
# ---------------------------------------------------------------------------


def _placeholder_base_id(ph_type) -> str:
    if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
        return "title"
    if ph_type == PP_PLACEHOLDER.SUBTITLE:
        return "subtitle"
    if ph_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
        return "body"
    if ph_type == PP_PLACEHOLDER.PICTURE:
        return "hero"
    return "field"


def _unique(base: str, used: set[str]) -> str:
    if base not in used:
        used.add(base)
        return base
    i = 2
    while f"{base}_{i}" in used:
        i += 1
    out = f"{base}_{i}"
    used.add(out)
    return out


def _safe_enum_short_name(value) -> str:
    """Lower-snake-case the short name of a python-pptx enum. Handles
    None / non-enum gracefully by falling back to str(value).
    """
    if value is None:
        return ""
    name = getattr(value, "name", None)
    if not name:
        return str(value).lower()
    return name.lower()
