"""pptx-skill authoring CLI.

Turns example .pptx files into a portable template library. See PLAN.md
for the design. Seven commands: ingest, status, next, prompt, validate,
preview, build.

Workspace state lives under authoring/workspace/ and is gitignored.
Hand-edited YAML sidecars (slide_NN.yaml, <sha1>.yaml) are the unit of
description. `validate` auto-promotes complete pending sidecars to done.
`build` emits dist/skill.zip — the consumer artifact.
"""

from __future__ import annotations

import hashlib
import json
import os
import posixpath
import re
import shutil
import subprocess
import sys
import zipfile
from pathlib import Path
from typing import Iterable
from urllib.parse import unquote
from xml.etree import ElementTree as ET

import click
import yaml
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

import ingest_v5  # v5 redesign — self-contained, removable as a unit


HERE = Path(__file__).resolve().parent
WORKSPACE = HERE / "workspace"
DIST = HERE / "dist"
SCHEMAS = HERE / "schemas"
PROMPTS = HERE / "prompts"
CONSUMER = HERE.parent / "consumer"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def ensure_workspace() -> None:
    (WORKSPACE / "decks").mkdir(parents=True, exist_ok=True)
    (WORKSPACE / "assets").mkdir(parents=True, exist_ok=True)


def read_yaml(path: Path) -> dict:
    text = path.read_text(encoding="utf-8")
    data = yaml.safe_load(text) or {}
    if not isinstance(data, dict):
        raise click.ClickException(f"{path}: expected mapping at top level")
    return data


def write_yaml(path: Path, data: dict) -> None:
    path.write_text(
        yaml.safe_dump(data, sort_keys=False, allow_unicode=True, width=100),
        encoding="utf-8",
    )


def aspect_ratio(width: int, height: int) -> str:
    """Pick a human aspect ratio token: 16:9, 4:3, 1:1, 3:4, 9:16, free."""
    if width <= 0 or height <= 0:
        return "free"
    r = width / height
    candidates = {
        "16:9": 16 / 9,
        "4:3": 4 / 3,
        "1:1": 1.0,
        "3:4": 3 / 4,
        "9:16": 9 / 16,
    }
    best, best_err = "free", 0.10  # tolerance
    for name, target in candidates.items():
        err = abs(r - target) / target
        if err < best_err:
            best, best_err = name, err
    return best


def position_quadrant(
    left: int, top: int, width: int, height: int, slide_w: int, slide_h: int
) -> str:
    """Coarse position: left/center/right + top/middle/bottom.

    Uses the shape's visual center, not its top-left, so a wide shape
    spanning the slide doesn't get mislabelled by its anchor corner.
    """
    cx = left + width // 2
    cy = top + height // 2
    if cx < slide_w / 3:
        h = "left"
    elif cx < 2 * slide_w / 3:
        h = "center"
    else:
        h = "right"
    if cy < slide_h / 3:
        v = "top"
    elif cy < 2 * slide_h / 3:
        v = "middle"
    else:
        v = "bottom"
    return f"{v}-{h}"


# ---------------------------------------------------------------------------
# Ingest
# ---------------------------------------------------------------------------


PLACEHOLDER_TEXTUAL = {
    PP_PLACEHOLDER.TITLE,
    PP_PLACEHOLDER.SUBTITLE,
    PP_PLACEHOLDER.BODY,
    PP_PLACEHOLDER.CENTER_TITLE,
    PP_PLACEHOLDER.OBJECT,
}


def _shape_is_picture(shape) -> bool:
    return getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE


# v4: theme-color name → clrScheme slot. MSO_THEME_COLOR has aliases
# (BACKGROUND_1 / LIGHT_1 / DARK_1 are sometimes interchangeable with
# TEXT_1 etc. in different OOXML dialects); the map below normalises
# all known names to the 12 canonical slots stored in theme.yaml.
_THEME_COLOR_TO_SLOT = {
    "accent_1": "accent1", "accent_2": "accent2",
    "accent_3": "accent3", "accent_4": "accent4",
    "accent_5": "accent5", "accent_6": "accent6",
    "background_1": "lt1", "background_2": "lt2",
    "text_1": "dk1", "text_2": "dk2",
    "dark_1": "dk1", "light_1": "lt1",
    "dark_2": "dk2", "light_2": "lt2",
    "hyperlink": "hlink", "followed_hyperlink": "folHlink",
}


def _resolve_run_color(run, theme_palette: dict) -> tuple[str, str]:
    """Inspect a run's font.color. Returns (hex, role).

    `hex` is "#RRGGBB" if resolvable, else empty.
    `role` is the clrScheme slot name (e.g. "accent1") if the run uses
    a theme color, else empty.
    """
    try:
        color = run.font.color
        ctype = color.type
    except (AttributeError, ValueError):
        return "", ""
    if ctype is None:
        return "", ""

    # Direct RGB.
    try:
        rgb = color.rgb
        if rgb is not None:
            return f"#{str(rgb).upper()}", ""
    except (AttributeError, ValueError):
        pass

    # Theme-color reference — resolve to slot, then to hex via theme.
    try:
        tc = color.theme_color
        if tc is not None:
            name = getattr(tc, "name", str(tc)).lower()
            slot = _THEME_COLOR_TO_SLOT.get(name, "")
            if slot:
                return theme_palette.get(slot, ""), slot
    except (AttributeError, ValueError):
        pass

    return "", ""


def _extract_slot_style(shape, theme_palette: dict) -> dict:
    """Snapshot the slot's first-run font properties for v4.

    Inspects only the first paragraph's first run — the "anchor" style
    that compose currently inherits when filling the slot. Returns a
    dict with whichever fields could be resolved; missing fields are
    omitted (not nulled).
    """
    try:
        tf = shape.text_frame
    except (AttributeError, ValueError):
        return {}
    paras = list(getattr(tf, "paragraphs", []) or [])
    if not paras:
        return {}
    runs = list(getattr(paras[0], "runs", []) or [])
    if not runs:
        return {}
    run = runs[0]
    font = run.font

    out: dict = {}
    try:
        if font.name:
            out["font"] = font.name
    except (AttributeError, ValueError):
        pass
    try:
        if font.size is not None:
            out["size_pt"] = float(font.size.pt)
    except (AttributeError, ValueError):
        pass
    try:
        if font.bold is not None:
            out["bold"] = bool(font.bold)
    except (AttributeError, ValueError):
        pass
    try:
        if font.italic is not None:
            out["italic"] = bool(font.italic)
    except (AttributeError, ValueError):
        pass
    color_hex, color_role = _resolve_run_color(run, theme_palette)
    if color_hex:
        out["color"] = color_hex
    if color_role:
        out["color_role"] = color_role
    return out


def _slot_id_for_placeholder(ph_type, used: set) -> str:
    if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
        base = "title"
    elif ph_type == PP_PLACEHOLDER.SUBTITLE:
        base = "subtitle"
    elif ph_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
        base = "body"
    elif ph_type == PP_PLACEHOLDER.PICTURE:
        base = "hero"
    else:
        base = "field"
    return _unique_id(base, used)


def _unique_id(base: str, used: set) -> str:
    if base not in used:
        used.add(base)
        return base
    i = 2
    while f"{base}_{i}" in used:
        i += 1
    out = f"{base}_{i}"
    used.add(out)
    return out


def detect_slots(
    slide, slide_w: int, slide_h: int, theme_palette: dict | None = None
) -> tuple[list[dict], dict]:
    """Detect slot definitions on a slide.

    Returns (slots, shape_renames) where shape_renames maps the original
    shape element id to the slot id we want set as shape.name (so the
    consumer can find each slot by shape name at compose time).

    `theme_palette` is the deck's clrScheme → hex map (from theme.yaml).
    Used to resolve per-slot theme colour references into hex; pass an
    empty dict to skip the v4 style snapshot.
    """
    theme_palette = theme_palette or {}
    slots: list[dict] = []
    used_ids: set = set()
    renames: dict = {}  # shape_id (int) -> slot_id

    slide_area = slide_w * slide_h

    # First pass: placeholders.
    for shape in list(slide.placeholders):
        ph = shape.placeholder_format
        ph_type = ph.type
        if ph_type == PP_PLACEHOLDER.PICTURE:
            slot_id = _unique_id("hero", used_ids)
            slots.append(
                {
                    "id": slot_id,
                    "kind": "image",
                    "aspect": aspect_ratio(shape.width or 0, shape.height or 0),
                }
            )
            renames[shape.shape_id] = slot_id
            continue
        # Tables in a Content/OBJECT placeholder show up before the
        # textual check, so an OBJECT-with-table doesn't get mistyped
        # as a text slot (FINDINGS A2.7).
        if getattr(shape, "has_table", False):
            try:
                rows = len(list(shape.table.rows))
                cols = len(list(shape.table.columns))
            except (AttributeError, ValueError):
                rows = cols = 0
            slot_id = _unique_id("data", used_ids)
            slots.append(
                {
                    "id": slot_id,
                    "kind": "table",
                    "rows": rows,
                    "cols": cols,
                }
            )
            renames[shape.shape_id] = slot_id
            continue
        if ph_type not in PLACEHOLDER_TEXTUAL:
            # Decorative / unsupported — leave frozen.
            continue
        # Textual placeholder.
        tf = getattr(shape, "text_frame", None)
        if tf is None:
            continue
        text = tf.text or ""
        paragraphs = [p for p in tf.paragraphs]
        is_bulleted = (
            ph_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT)
            and len(paragraphs) > 1
            and any((p.text or "").strip() for p in paragraphs)
        )
        slot_id = _slot_id_for_placeholder(ph_type, used_ids)
        style = _extract_slot_style(shape, theme_palette)
        if is_bulleted:
            non_empty = [p for p in paragraphs if (p.text or "").strip()]
            slot: dict = {
                "id": slot_id,
                "kind": "bullets",
                "max_items": max(1, len(non_empty)),
            }
        else:
            slot = {
                "id": slot_id,
                "kind": "text",
                "max_chars": max(20, int(len(text) * 1.5) or 60),
            }
        if style:
            slot["style"] = style
        slots.append(slot)
        renames[shape.shape_id] = slot_id

    # Second pass: free pictures > 20% slide area.
    for shape in list(slide.shapes):
        if not _shape_is_picture(shape):
            continue
        if getattr(shape, "is_placeholder", False):
            continue
        w = shape.width or 0
        h = shape.height or 0
        if slide_area <= 0:
            continue
        if (w * h) / slide_area <= 0.20:
            continue  # frozen background / logo / decoration
        slot_id = _unique_id("hero", used_ids)
        slots.append(
            {
                "id": slot_id,
                "kind": "image",
                "aspect": aspect_ratio(w, h),
            }
        )
        renames[shape.shape_id] = slot_id

    return slots, renames


def _shape_kind_label(shape) -> str | None:
    """Short kind label for a non-placeholder shape, or None if uninteresting.

    Used by ``infer_layout`` to label visually-meaningful content that
    isn't a placeholder slot. Returns one of: ``image``, ``vector``,
    ``table``, ``chart``, ``smartart``, ``callout``, ``freeform``,
    ``text``. Returns ``None`` for shape types we don't classify so the
    caller can skip them (decorations, connectors, etc.).
    """
    if _shape_is_picture(shape):
        return "image"
    try:
        if getattr(shape, "has_table", False):
            return "table"
        if getattr(shape, "has_chart", False):
            return "chart"
    except (AttributeError, ValueError):
        pass
    if _shape_is_smartart(shape):
        return "smartart"
    stype = getattr(shape, "shape_type", None)
    if stype == MSO_SHAPE_TYPE.AUTO_SHAPE:
        return "callout"
    if stype == MSO_SHAPE_TYPE.FREEFORM:
        return "freeform"
    if stype == MSO_SHAPE_TYPE.TEXT_BOX:
        return "text"
    return None


def _shape_geometry(shape, slide_w: int, slide_h: int) -> dict:
    """Per-shape fractional geometry + coarse region label.

    Returns ``{x, y, w, h, region}`` where x/y/w/h are slide-relative
    fractions (0.0-1.0, rounded to 3 decimals to keep YAML terse) and
    ``region`` is the position-quadrant string used by ``infer_layout``.
    Empty dict if slide dimensions are zero.
    """
    if not slide_w or not slide_h:
        return {}
    left = shape.left or 0
    top = shape.top or 0
    width = shape.width or 0
    height = shape.height or 0
    return {
        "x": round(left / slide_w, 3),
        "y": round(top / slide_h, 3),
        "w": round(width / slide_w, 3),
        "h": round(height / slide_h, 3),
        "region": position_quadrant(left, top, width, height, slide_w, slide_h),
    }


def infer_layout(slide, slots: list[dict], renames: dict, slide_w: int, slide_h: int) -> str:
    """One-line spatial summary of every visually-meaningful shape.

    Walks the slide (descending into groups) and labels each shape by
    either its slot id (if it became a template slot during
    ``detect_slots``) or by ``_shape_kind_label`` (otherwise). Skips
    decorations and hairlines via ``_atom_too_small``.

    Output is a comma-separated list of ``label@region`` tokens, e.g.
    ``"title@top-center, image@middle-right, callout@bottom-left"``.
    This is what the compose-time agent reads to understand a slide's
    *anatomy* — distinct from its descriptive ``intent``.
    """
    parts: list[str] = []
    seen_labels: dict[str, int] = {}  # disambiguate repeats: image, image#2
    for shape in _iter_shapes_recursive(slide):
        if shape.shape_id in renames:
            label = renames[shape.shape_id]
        else:
            if _atom_too_small(shape, slide_w, slide_h):
                continue
            kind = _shape_kind_label(shape)
            if not kind:
                continue
            label = kind
        # If we see the same kind twice on one slide (e.g. two pictures),
        # disambiguate as image, image#2, image#3 — gives the compose
        # agent a stable way to refer to the n-th instance positionally.
        count = seen_labels.get(label, 0) + 1
        seen_labels[label] = count
        tag = label if count == 1 else f"{label}#{count}"
        left = shape.left or 0
        top = shape.top or 0
        width = shape.width or 0
        height = shape.height or 0
        region = position_quadrant(left, top, width, height, slide_w, slide_h)
        parts.append(f"{tag}@{region}")
    return ", ".join(parts) if parts else "freeform"


def write_slide_fragment(src_path: Path, slide_idx: int, out_path: Path, renames: dict) -> None:
    """Save a single-slide fragment by reopening the deck and dropping others.

    The fragment keeps its slide masters/layouts so it can be re-composed.
    `renames` maps shape_id -> name we want to set so the consumer can
    locate slots at compose time.

    After python-pptx writes the package, we garbage-collect orphan
    parts (other slides' media, layouts no slide references, etc.) so
    each fragment carries only what its slide actually needs.
    """
    prs = Presentation(str(src_path))

    # Apply renames first (on the current slide, before we drop the others).
    if 0 <= slide_idx < len(prs.slides):
        slide = prs.slides[slide_idx]
        for shape in list(slide.shapes):
            if shape.shape_id in renames:
                try:
                    shape.name = renames[shape.shape_id]
                except Exception:
                    pass

    # Drop every other slide from sldIdLst + the package relationship.
    # KeyError is the only python-pptx-documented failure for drop_rel
    # (rId already gone); anything else is real and should surface.
    sld_id_list = prs.slides._sldIdLst
    sld_ids = list(sld_id_list)
    for i, sld_id in enumerate(sld_ids):
        if i == slide_idx:
            continue
        rId = sld_id.rId
        sld_id_list.remove(sld_id)
        try:
            prs.part.drop_rel(rId)
        except KeyError:
            pass

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))

    # Prune orphans introduced by python-pptx's "leave parts behind" model.
    prune_unreachable_parts(out_path)


# --- Fragment GC: keep only parts reachable from /_rels/.rels ---------------
#
# python-pptx's save() prunes unused layout *rels* when it writes a
# fragment (e.g. all the layouts other slides used) but leaves orphan
# `<p:sldLayoutId>` entries inside the master XML. Those orphans then
# point to rIds that no longer exist — and python-pptx (and PowerPoint)
# blow up when iterating the master's layouts. It also leaves any other
# orphan parts (notes themes, dropped images, etc.) inside the zip.
#
# This pass:
#  1. For each master, removes `<p:sldLayoutId>` entries whose r:id is
#     no longer in the master's rels file.
#  2. Walks /_rels/.rels through every part's rels, collecting the
#     closure of reachable parts.
#  3. Filters `[Content_Types].xml` Override entries to match.
#  4. Rewrites the zip with only reachable parts.

_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
_CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"


def _rels_path_for(part_path: str) -> str:
    dir_part, _, name = part_path.rpartition("/")
    return f"{dir_part}/_rels/{name}.rels" if dir_part else f"_rels/{name}.rels"


def _resolve_rel(rels_path: str, target: str) -> str:
    target = unquote(target)
    if rels_path == "_rels/.rels":
        base = ""
    else:
        base = posixpath.dirname(posixpath.dirname(rels_path))
    if target.startswith("/"):
        return posixpath.normpath(target.lstrip("/"))
    return posixpath.normpath(posixpath.join(base, target)) if base else target


def _parse_rels(xml_bytes: bytes) -> list[tuple[str, str]]:
    """Return list of (target, type) for non-External relationships."""
    try:
        root = ET.fromstring(xml_bytes)
    except ET.ParseError:
        return []
    out: list[tuple[str, str]] = []
    for rel in root.findall(f"{{{_REL_NS}}}Relationship"):
        if rel.get("TargetMode") == "External":
            continue
        out.append((rel.get("Target", ""), rel.get("Type", "")))
    return out


_PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
_DML_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
# `r:` inside document XML (slideMaster.xml, presentation.xml, etc.)
# binds to officeDocument/relationships — distinct from the package
# /relationships namespace used by .rels files themselves.
_DOC_REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _valid_rids(rels_xml: bytes) -> set[str]:
    """Return the Id set declared by a rels file."""
    try:
        root = ET.fromstring(rels_xml)
    except ET.ParseError:
        return set()
    return {rel.get("Id", "") for rel in root.findall(f"{{{_REL_NS}}}Relationship")}


def _layouts_used_by_slides(zf: zipfile.ZipFile, names: list[str]) -> set[str]:
    """Layouts referenced via rels by any slide still in the package."""
    used: set[str] = set()
    names_set = set(names)
    for name in names:
        if not (name.startswith("ppt/slides/slide") and name.endswith(".xml")):
            continue
        rels_path = _rels_path_for(name)
        if rels_path not in names_set:
            continue
        for target, rtype in _parse_rels(zf.read(rels_path)):
            if rtype.endswith("/slideLayout"):
                used.add(_resolve_rel(rels_path, target))
    return used


def _prune_master_rels(rels_xml: bytes, rels_path: str, kept_layouts: set[str]) -> bytes:
    """Drop slideLayout rels not in kept_layouts. Returns rels_xml unchanged
    if nothing to drop (identity-comparable)."""
    try:
        root = ET.fromstring(rels_xml)
    except ET.ParseError:
        return rels_xml
    changed = False
    for rel in list(root.findall(f"{{{_REL_NS}}}Relationship")):
        if not rel.get("Type", "").endswith("/slideLayout"):
            continue
        if _resolve_rel(rels_path, rel.get("Target", "")) not in kept_layouts:
            root.remove(rel)
            changed = True
    if not changed:
        return rels_xml
    ET.register_namespace("", _REL_NS)
    return b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n' + \
        ET.tostring(root, encoding="utf-8")


def _repair_master_xml(master_xml: bytes, valid_rids: set[str]) -> bytes:
    """Drop <p:sldLayoutId> entries whose r:id is no longer in the master's
    rels. python-pptx trims unused layouts from the rels file when it
    saves a fragment, but leaves orphan sldLayoutId entries in the master
    XML — which then break loading because each sldLayoutId is resolved
    via its rId."""
    try:
        root = ET.fromstring(master_xml)
    except ET.ParseError:
        return master_xml
    rid_attr = f"{{{_DOC_REL_NS}}}id"
    changed = False
    for lst in root.findall(f"{{{_PML_NS}}}sldLayoutIdLst"):
        for entry in list(lst.findall(f"{{{_PML_NS}}}sldLayoutId")):
            if entry.get(rid_attr) not in valid_rids:
                lst.remove(entry)
                changed = True
    if not changed:
        return master_xml
    ET.register_namespace("p", _PML_NS)
    ET.register_namespace("a", _DML_NS)
    ET.register_namespace("r", _DOC_REL_NS)
    return b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n' + \
        ET.tostring(root, encoding="utf-8")


def _master_xml_path_for_rels(rels_path: str) -> str:
    # 'ppt/slideMasters/_rels/slideMaster1.xml.rels'
    # -> 'ppt/slideMasters/slideMaster1.xml'
    parent = posixpath.dirname(posixpath.dirname(rels_path))
    name = posixpath.basename(rels_path)
    if name.endswith(".rels"):
        name = name[: -len(".rels")]
    return f"{parent}/{name}"


def prune_unreachable_parts(pptx_path: Path) -> tuple[int, int]:
    """Rewrite a .pptx in place, keeping only parts reachable from the
    package root rels. Returns (kept, removed) part counts."""
    with zipfile.ZipFile(pptx_path, "r") as zf:
        all_names = list(zf.namelist())
        name_set = set(all_names)

        # Active pruning + orphan repair for each master.
        # 1. Drop slideLayout rels not used by any surviving slide.
        # 2. Then drop <p:sldLayoutId> entries in the master XML pointing
        #    to rIds that no longer exist (covers both our prune AND any
        #    orphans python-pptx left behind from its own save logic).
        kept_layouts = _layouts_used_by_slides(zf, all_names)
        rewritten: dict[str, bytes] = {}
        for name in all_names:
            if "/slideMasters/_rels/" not in name or not name.endswith(".rels"):
                continue
            master_xml_path = _master_xml_path_for_rels(name)
            if master_xml_path not in name_set:
                continue

            original_rels = zf.read(name)
            new_rels = _prune_master_rels(original_rels, name, kept_layouts)
            if new_rels is not original_rels:
                rewritten[name] = new_rels

            valid = _valid_rids(new_rels)
            original_master = zf.read(master_xml_path)
            repaired = _repair_master_xml(original_master, valid)
            if repaired is not original_master:
                rewritten[master_xml_path] = repaired

        reachable: set[str] = set()
        for keep in ("[Content_Types].xml", "_rels/.rels"):
            if keep in name_set:
                reachable.add(keep)

        queue: list[str] = ["_rels/.rels"]
        seen_rels: set[str] = set()
        while queue:
            rels_path = queue.pop(0)
            if rels_path in seen_rels or rels_path not in name_set:
                continue
            seen_rels.add(rels_path)
            reachable.add(rels_path)
            xml = rewritten.get(rels_path, zf.read(rels_path))
            for target, _ in _parse_rels(xml):
                resolved = _resolve_rel(rels_path, target)
                if resolved in name_set and resolved not in reachable:
                    reachable.add(resolved)
                    part_rels = _rels_path_for(resolved)
                    if part_rels in name_set and part_rels not in seen_rels:
                        queue.append(part_rels)

        try:
            ct_root = ET.fromstring(zf.read("[Content_Types].xml"))
            for override in list(ct_root.findall(f"{{{_CT_NS}}}Override")):
                pn = override.get("PartName", "").lstrip("/")
                if pn not in reachable:
                    ct_root.remove(override)
            ET.register_namespace("", _CT_NS)
            new_ct = b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n' + \
                ET.tostring(ct_root, encoding="utf-8")
        except ET.ParseError:
            new_ct = zf.read("[Content_Types].xml")

        tmp_path = pptx_path.with_suffix(pptx_path.suffix + ".pruning")
        with zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as zf_out:
            for name in all_names:
                if name not in reachable:
                    continue
                if name == "[Content_Types].xml":
                    zf_out.writestr(name, new_ct)
                elif name in rewritten:
                    zf_out.writestr(name, rewritten[name])
                else:
                    zf_out.writestr(name, zf.read(name))

    tmp_path.replace(pptx_path)
    return len(reachable), len(name_set) - len(reachable)


_SLIDE_DESCRIPTIVE_DEFAULTS = {
    "intent": "",
    "feel": "",
    "suitable_for": [],
    "status": "pending",
    "notes": "",
    "interpretation": "",
}


def write_slide_yaml_stub(
    out_path: Path,
    slide_id: str,
    layout: str,
    slots: list[dict],
    deck_stem: str,
    slide_number: int,
    *,
    theme_colors: dict | None = None,
    fonts: dict | None = None,
    inventory: list | None = None,
) -> None:
    """Write a slide sidecar in canonical-order YAML.

    Structural fields (id, layout, theme_colors, fonts, slots, inventory,
    sources) are always refreshed. Descriptive fields (intent, feel,
    suitable_for, status, notes) are preserved verbatim if the file
    already exists — regardless of status. New stubs default
    descriptive fields to empty + status=pending.
    """
    existing: dict = {}
    if out_path.exists():
        try:
            existing = read_yaml(out_path)
        except Exception:
            existing = {}

    descriptive = {
        k: existing.get(k, default) for k, default in _SLIDE_DESCRIPTIVE_DEFAULTS.items()
    }

    out = {
        "intent": descriptive["intent"],
        "feel": descriptive["feel"],
        "suitable_for": descriptive["suitable_for"],
        "status": descriptive["status"],
        "notes": descriptive["notes"],
        "interpretation": descriptive["interpretation"],
        "id": slide_id,
        "layout": layout,
        "theme_colors": theme_colors or {},
        "fonts": fonts or {},
        "slots": slots,
        "inventory": inventory or [],
        "sources": [{"deck": deck_stem, "slide": slide_number}],
    }
    write_yaml(out_path, out)


def _current_deck_stems() -> set[str]:
    """Names of decks currently present under workspace/decks/.

    Used to prune asset.sources entries that point at decks the user
    has since deleted (e.g. an ingest mistake cleaned up by `rm -rf`).
    Full workspace lifecycle reconciliation is out of v4 scope; this
    keeps cross-deck source lists from accreting noise just by ingest
    touching the asset again.
    """
    decks_dir = WORKSPACE / "decks"
    if not decks_dir.exists():
        return set()
    return {p.name for p in decks_dir.iterdir() if p.is_dir()}


def _prune_dead_sources(sources: list, valid_decks: set[str]) -> list:
    return [s for s in sources if (s.get("deck") or "") in valid_decks]


_ASSET_DESCRIPTIVE_DEFAULTS = {
    "kind": "",
    "subject": "",
    "depicts": "",
    "feel": "",
    "composition": "",
    "colors": [],
    "scope": [],
    "suitable_for": [],
    "status": "pending",
    "notes": "",
    "interpretation": "",
}


def write_asset_yaml_stub(
    out_path: Path,
    asset_id: str,
    sha1: str,
    deck_stem: str,
    slide_number: int,
    *,
    kind: str | None = None,
    colors_hex: list | None = None,
    recolor_targets: list | None = None,
    table: dict | None = None,
    chart: dict | None = None,
    shape: dict | None = None,
    smartart: dict | None = None,
) -> None:
    """Write an asset sidecar in canonical-order YAML.

    Structural fields (id, sha1, colors_hex, recolor_targets, sources)
    are refreshed when the ingest layer supplies new values. Descriptive
    fields (kind/subject/depicts/feel/composition/colors/scope/
    suitable_for/status/notes) are preserved verbatim across re-ingest.

    The `kind` kwarg seeds the field only on first ingest — never
    overwrites a hand-edited value. Pass it when ingest knows the kind
    (e.g. "vector" for an SVG sibling).
    """
    existing: dict = {}
    if out_path.exists():
        try:
            existing = read_yaml(out_path)
        except Exception:
            existing = {}

    descriptive = {
        k: existing.get(k, default) for k, default in _ASSET_DESCRIPTIVE_DEFAULTS.items()
    }
    if kind is not None and not descriptive["kind"]:
        descriptive["kind"] = kind

    # Source list: append the current touch, then prune entries pointing
    # at decks no longer in workspace/decks/ (handles the "I ran ingest
    # wrong then rm -rf'd the bad dir" case).
    sources = list(existing.get("sources") or [])
    if not any(
        s.get("deck") == deck_stem and s.get("slide") == slide_number for s in sources
    ):
        sources.append({"deck": deck_stem, "slide": slide_number})
    sources = _prune_dead_sources(sources, _current_deck_stems())

    # colors_hex: preserve previous extraction if the new pass returned
    # nothing (e.g. PIL absent, or unsupported format) so we don't clobber.
    if colors_hex:
        out_colors_hex = colors_hex
    else:
        out_colors_hex = list(existing.get("colors_hex") or [])

    out = {
        "kind": descriptive["kind"],
        "subject": descriptive["subject"],
        "depicts": descriptive["depicts"],
        "feel": descriptive["feel"],
        "composition": descriptive["composition"],
        "colors": descriptive["colors"],
        "colors_hex": out_colors_hex,
        "scope": descriptive["scope"],
        "suitable_for": descriptive["suitable_for"],
        "status": descriptive["status"],
        "notes": descriptive["notes"],
        "interpretation": descriptive["interpretation"],
        "id": asset_id,
        "sha1": sha1,
        "sources": sources,
    }

    if recolor_targets is not None:
        out["recolor_targets"] = recolor_targets
    elif "recolor_targets" in existing:
        out["recolor_targets"] = existing["recolor_targets"]

    # Kind-specific blocks: prefer fresh input from ingest, else
    # preserve whatever was already there (human edits + prior ingest
    # data both survive when this re-ingest doesn't supply the block).
    for key, value in (
        ("table", table),
        ("chart", chart),
        ("shape", shape),
        ("smartart", smartart),
    ):
        if value is not None:
            out[key] = value
        elif key in existing:
            out[key] = existing[key]

    write_yaml(out_path, out)


# v4: dominant-color extraction (raster) + SVG sibling extraction.
#
# Raster: PIL quantize to a small palette, sort by frequency, dedupe
# near-identicals (per-channel tolerance). Returns hex strings.
#
# SVG sibling: PPTX stores SVG `Picture` shapes with BOTH a raster
# fallback (<a:blip>) and the vector source (<asvg:svgBlip>) inside
# the same <p:pic>. python-pptx's shape.image returns the raster;
# this digs the SVG part out via the rel id.

_ASVG_NS = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"

_SVG_COLOR_RE = re.compile(rb'(?:fill|stroke)\s*=\s*["\']?(#[0-9a-fA-F]{6})')


def _color_close(a: str, b: str, tol: int = 16) -> bool:
    """True if two #RRGGBB strings are within `tol` per channel."""
    try:
        ai = int(a.lstrip("#"), 16)
        bi = int(b.lstrip("#"), 16)
    except ValueError:
        return False
    ar, ag, ab = (ai >> 16) & 0xFF, (ai >> 8) & 0xFF, ai & 0xFF
    br, bg, bb_ = (bi >> 16) & 0xFF, (bi >> 8) & 0xFF, bi & 0xFF
    return abs(ar - br) <= tol and abs(ag - bg) <= tol and abs(ab - bb_) <= tol


def _extract_dominant_colors_hex(blob_path: Path, n: int = 3) -> list[str]:
    """Return up to `n` dominant colors as #RRGGBB strings.

    Quantizes to 8 colours after a fast thumbnail downsample, sorts
    palette entries by pixel frequency, and dedupes near-identicals
    (per-channel tolerance via _color_close).
    """
    try:
        from PIL import Image
    except ImportError:
        return []
    try:
        img = Image.open(blob_path)
        if img.mode in ("RGBA", "LA"):
            bg = Image.new("RGB", img.size, (255, 255, 255))
            bg.paste(img, mask=img.split()[-1])
            img = bg
        elif img.mode != "RGB":
            img = img.convert("RGB")
        img.thumbnail((150, 150))
        quant = img.quantize(colors=8, kmeans=0)
        palette = quant.getpalette() or []
        counts = sorted(quant.getcolors() or [], reverse=True)  # [(count, idx), ...]
        out: list[str] = []
        for _count, idx in counts:
            if len(out) >= n:
                break
            base = idx * 3
            if base + 3 > len(palette):
                continue
            r, g, b = palette[base], palette[base + 1], palette[base + 2]
            hexstr = f"#{r:02X}{g:02X}{b:02X}"
            if any(_color_close(hexstr, prev) for prev in out):
                continue
            out.append(hexstr)
        return out
    except Exception:
        return []


def _extract_svg_sibling(shape) -> tuple[bytes | None, str]:
    """If the Picture shape carries an asvg:svgBlip sibling, return
    (svg_bytes, "svg"). Else (None, "")."""
    try:
        sp = shape._element
    except AttributeError:
        return None, ""
    svg_blips = sp.findall(f".//{{{_ASVG_NS}}}svgBlip")
    if not svg_blips:
        return None, ""
    embed = svg_blips[0].get(f"{{{_DOC_REL_NS}}}embed", "")
    if not embed:
        return None, ""
    try:
        slide_part = shape.part
        svg_part = slide_part.related_part(embed)
        return svg_part.blob, "svg"
    except (AttributeError, KeyError):
        return None, ""


def _extract_svg_colors(svg_bytes: bytes, n: int = 5) -> list[str]:
    """Quick regex scan of unique fill/stroke hex colors in SVG source."""
    seen: list[str] = []
    for match in _SVG_COLOR_RE.finditer(svg_bytes):
        hex_val = match.group(1).decode("ascii").upper()
        if hex_val not in seen:
            seen.append(hex_val)
        if len(seen) >= n:
            break
    return seen


# v4: structured-atom capture for non-picture, non-text-placeholder
# shape kinds — tables, charts, callouts/auto-shapes, freeforms,
# smartart. Each becomes its own asset record with a kind-specific
# descriptive block and the shape XML serialised as the "binary".
#
# CAVEAT: the saved XML is a *fragment* — it references theme/style/
# chart-data parts that live elsewhere in the source pptx. Phase D
# compose will need to pull those related parts in when instantiating
# an atom on a foreign slide. For B4 the goal is *addressability*
# (agent can list/filter/describe them), not yet *usability*.


def _serialize_shape_xml(shape) -> bytes:
    """Serialize a shape's lxml element to bytes for hashing + storage."""
    try:
        from lxml import etree as lxml_etree
        return lxml_etree.tostring(shape._element, pretty_print=False)
    except (ImportError, AttributeError):
        return b""


def _shape_atom_sha(xml_bytes: bytes) -> str:
    return hashlib.sha1(xml_bytes).hexdigest()


def _shape_is_smartart(shape) -> bool:
    """SmartArt shows up as a GraphicFrame with diagram-typed graphicData."""
    try:
        el = shape._element
    except AttributeError:
        return False
    for gd in el.findall(f".//{{{_DML_NS}}}graphicData"):
        if "diagram" in (gd.get("uri", "") or ""):
            return True
    return False


def _safe_enum_short_name(value) -> str:
    """Convert a pptx enum member to a short lowercase token, or ''.

    Example: MSO_SHAPE.ROUNDED_RECTANGLE -> 'rounded_rectangle'.
    """
    if value is None:
        return ""
    s = str(value)
    return s.rsplit(".", 1)[-1].lower() if "." in s else s.lower()


def _shape_is_recolorable(shape) -> bool:
    """True if the shape has a solid colour fill that compose could rewrite."""
    try:
        fill = shape.fill
        ftype = fill.type
    except (AttributeError, ValueError):
        return False
    return ftype is not None and str(ftype).endswith("SOLID")


def _atom_too_small(shape, slide_w: int, slide_h: int, min_frac: float = 0.005) -> bool:
    """Filter out trivially small geometric atoms (decorative dots, hairlines)."""
    w = getattr(shape, "width", 0) or 0
    h = getattr(shape, "height", 0) or 0
    slide_area = slide_w * slide_h
    if slide_area <= 0:
        return False
    return (w * h) / slide_area < min_frac


# v4.1: Designers commonly wrap multiple callouts/pictures in a Group
# shape (PowerPoint's "Group" menu item) so they can be moved as a unit.
# A flat iteration over slide.shapes encounters the GROUP but never
# descends — the children silently disappear from the atom catalog and
# the picture catalog. We descend up to MAX_GROUP_DEPTH levels; past 4
# levels is almost always abused groups (someone using groups as a
# layer system) and the atom semantics aren't interesting.
_MAX_GROUP_DEPTH = 4


def _iter_shapes_recursive(container, _depth: int = 0):
    """Yield leaf shapes from a slide or group, descending into nested groups.

    The container's own GroupShape elements are NOT yielded — groups
    are containers, not atoms in their own right. Stops descending
    past ``_MAX_GROUP_DEPTH`` to bound pathologically nested decks.
    """
    for shape in container.shapes:
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            if _depth + 1 >= _MAX_GROUP_DEPTH:
                continue
            yield from _iter_shapes_recursive(shape, _depth + 1)
        else:
            yield shape


def _extract_table_atom(shape, deck_stem: str, slide_number: int, assets_dir: Path) -> str | None:
    table = shape.table
    xml = _serialize_shape_xml(shape)
    if not xml:
        return None
    sha = _shape_atom_sha(xml)
    asset_id = f"asset_{sha[:8]}"

    rows_list = list(table.rows)
    cols_list = list(table.columns)
    cells: list[list[str]] = []
    for r in rows_list:
        cells.append([(c.text or "").strip() for c in r.cells])

    headers = cells[0] if cells else []
    sample_cells = cells[1:min(4, len(cells))]

    bin_path = assets_dir / f"{sha}.xml"
    yaml_path = assets_dir / f"{sha}.yaml"
    if not bin_path.exists():
        bin_path.write_bytes(xml)
    write_asset_yaml_stub(
        yaml_path, asset_id, sha, deck_stem, slide_number,
        kind="table",
        table={
            "rows": len(rows_list),
            "cols": len(cols_list),
            "headers": headers,
            "sample_cells": sample_cells,
        },
    )
    return asset_id


def _extract_chart_atom(shape, deck_stem: str, slide_number: int, assets_dir: Path) -> str | None:
    chart = shape.chart
    xml = _serialize_shape_xml(shape)
    if not xml:
        return None
    sha = _shape_atom_sha(xml)
    asset_id = f"asset_{sha[:8]}"

    chart_type = ""
    try:
        chart_type = _safe_enum_short_name(chart.chart_type)
    except (AttributeError, ValueError):
        pass
    series_count = 0
    try:
        series_count = len(list(chart.series))
    except (AttributeError, ValueError):
        pass
    categories_count = 0
    try:
        plots = list(chart.plots)
        if plots:
            categories_count = len(list(plots[0].categories))
    except (AttributeError, ValueError, IndexError):
        pass

    bin_path = assets_dir / f"{sha}.xml"
    yaml_path = assets_dir / f"{sha}.yaml"
    if not bin_path.exists():
        bin_path.write_bytes(xml)
    write_asset_yaml_stub(
        yaml_path, asset_id, sha, deck_stem, slide_number,
        kind="chart",
        chart={
            "type": chart_type,
            "series_count": series_count,
            "categories_count": categories_count,
        },
    )
    return asset_id


def _extract_geometric_atom(
    shape,
    deck_stem: str,
    slide_number: int,
    assets_dir: Path,
    kind: str,
) -> str | None:
    """Save a callout (auto-shape) or freeform as a kind-specific asset."""
    xml = _serialize_shape_xml(shape)
    if not xml:
        return None
    sha = _shape_atom_sha(xml)
    asset_id = f"asset_{sha[:8]}"

    geometry = ""
    if kind == "callout":
        try:
            geometry = _safe_enum_short_name(shape.auto_shape_type)
        except (AttributeError, ValueError):
            pass

    bin_path = assets_dir / f"{sha}.xml"
    yaml_path = assets_dir / f"{sha}.yaml"
    if not bin_path.exists():
        bin_path.write_bytes(xml)
    write_asset_yaml_stub(
        yaml_path, asset_id, sha, deck_stem, slide_number,
        kind=kind,
        shape={"geometry": geometry, "is_recolorable": _shape_is_recolorable(shape)},
    )
    return asset_id


def _extract_smartart_atom(shape, deck_stem: str, slide_number: int, assets_dir: Path) -> str | None:
    xml = _serialize_shape_xml(shape)
    if not xml:
        return None
    sha = _shape_atom_sha(xml)
    asset_id = f"asset_{sha[:8]}"

    el = shape._element
    nodes = [(t.text or "").strip() for t in el.findall(f".//{{{_DML_NS}}}t")]
    nodes = [n for n in nodes if n]

    # Layout name lives in the data-model part (dgm:dataModel) which
    # isn't in the slide XML — leave empty for B4. Phase D can resolve
    # it when instantiating the atom.
    bin_path = assets_dir / f"{sha}.xml"
    yaml_path = assets_dir / f"{sha}.yaml"
    if not bin_path.exists():
        bin_path.write_bytes(xml)
    write_asset_yaml_stub(
        yaml_path, asset_id, sha, deck_stem, slide_number,
        kind="smartart",
        smartart={"layout": "", "nodes": nodes},
    )
    return asset_id


def extract_structured_atoms(
    slide,
    deck_stem: str,
    slide_number: int,
    assets_dir: Path,
    slide_w: int,
    slide_h: int,
    slot_shape_ids: set | None = None,
) -> list[dict]:
    """Save non-picture, non-text-placeholder shapes as typed atom assets.

    Captures tables, charts, smartart, auto-shapes (callouts), freeforms.
    Skips pictures (extract_picture_assets handles them), textual
    placeholders (those are template slots, not addressable atoms),
    any shape already promoted to a slot (`slot_shape_ids`), and atoms
    below ~0.5% slide area (decorative hairlines / single-pixel dots).
    Descends into Group shapes up to ``_MAX_GROUP_DEPTH`` levels so
    grouped callouts aren't silently lost.

    v4.2: returns ``[{atom, kind, x, y, w, h, region}, ...]`` — one
    per slide occurrence with positional info, for the slide's
    inventory. Previously returned bare asset ids.
    """
    slot_shape_ids = slot_shape_ids or set()
    extracted: list[dict] = []

    def _add(asset_id: str | None, shape, kind: str) -> None:
        if not asset_id:
            return
        extracted.append({
            "atom": asset_id,
            "kind": kind,
            **_shape_geometry(shape, slide_w, slide_h),
        })

    for shape in list(_iter_shapes_recursive(slide)):
        # Shapes that detect_slots claimed as template slots aren't
        # addressable as standalone atoms — the agent fills them via
        # the slot interface.
        if shape.shape_id in slot_shape_ids:
            continue
        if _shape_is_picture(shape):
            continue
        # Skip text/image placeholders — those drive template slots,
        # they aren't reusable atoms in their own right.
        if getattr(shape, "is_placeholder", False):
            ph_type = shape.placeholder_format.type
            if ph_type in PLACEHOLDER_TEXTUAL or ph_type == PP_PLACEHOLDER.PICTURE:
                continue

        # Order matters: SmartArt is a GraphicFrame with diagram URI;
        # tables and charts are also GraphicFrames distinguished by
        # has_table / has_chart. Check the more specific ones first.
        try:
            if getattr(shape, "has_table", False):
                _add(_extract_table_atom(shape, deck_stem, slide_number, assets_dir), shape, "table")
                continue
            if getattr(shape, "has_chart", False):
                _add(_extract_chart_atom(shape, deck_stem, slide_number, assets_dir), shape, "chart")
                continue
        except (AttributeError, ValueError):
            pass

        if _shape_is_smartart(shape):
            _add(_extract_smartart_atom(shape, deck_stem, slide_number, assets_dir), shape, "smartart")
            continue

        stype = getattr(shape, "shape_type", None)
        if stype == MSO_SHAPE_TYPE.AUTO_SHAPE:
            if _atom_too_small(shape, slide_w, slide_h):
                continue
            _add(
                _extract_geometric_atom(shape, deck_stem, slide_number, assets_dir, "callout"),
                shape, "callout",
            )
            continue
        if stype == MSO_SHAPE_TYPE.FREEFORM:
            if _atom_too_small(shape, slide_w, slide_h):
                continue
            _add(
                _extract_geometric_atom(shape, deck_stem, slide_number, assets_dir, "freeform"),
                shape, "freeform",
            )
            continue

    return extracted


def extract_picture_assets(
    slide,
    deck_stem: str,
    slide_number: int,
    assets_dir: Path,
    slide_w: int = 0,
    slide_h: int = 0,
) -> list[dict]:
    """Save all picture shapes on this slide as <sha1>.<ext> binaries + yaml stubs.

    v4: for each Picture shape we extract:
      - the raster (PNG/JPG/EMF/...) with PIL-derived dominant colours
      - any SVG sibling (asvg:svgBlip) as a separate vector asset with
        kind="vector" and recolor_targets seeded from fill/stroke colours

    v4.1: descends into Group shapes up to ``_MAX_GROUP_DEPTH`` levels
    so grouped pictures aren't lost.

    v4.2: returns a list of per-slide inventory dicts ``{atom, kind, x,
    y, w, h, region}`` instead of bare ids — gives the compose-time
    agent positional info for each picture. SVG siblings share the
    parent picture's geometry (they render in place of the raster on
    capable viewers).
    """
    extracted: list[dict] = []
    for shape in list(_iter_shapes_recursive(slide)):
        if not _shape_is_picture(shape):
            continue
        try:
            image = shape.image
        except Exception:
            continue

        # --- Raster / fallback picture ---
        blob = image.blob
        ext = (image.ext or "png").lstrip(".")
        sha = hashlib.sha1(blob).hexdigest()
        asset_id = f"asset_{sha[:8]}"
        bin_path = assets_dir / f"{sha}.{ext}"
        yaml_path = assets_dir / f"{sha}.yaml"
        if not bin_path.exists():
            bin_path.write_bytes(blob)
        colors_hex = _extract_dominant_colors_hex(bin_path)
        write_asset_yaml_stub(
            yaml_path,
            asset_id,
            sha,
            deck_stem,
            slide_number,
            colors_hex=colors_hex,
        )
        geom = _shape_geometry(shape, slide_w, slide_h)
        extracted.append({"atom": asset_id, "kind": "image", **geom})

        # --- SVG sibling (vector source) ---
        svg_blob, svg_ext = _extract_svg_sibling(shape)
        if svg_blob:
            svg_sha = hashlib.sha1(svg_blob).hexdigest()
            svg_asset_id = f"asset_{svg_sha[:8]}"
            svg_bin_path = assets_dir / f"{svg_sha}.{svg_ext}"
            svg_yaml_path = assets_dir / f"{svg_sha}.yaml"
            if not svg_bin_path.exists():
                svg_bin_path.write_bytes(svg_blob)
            svg_colors = _extract_svg_colors(svg_blob)
            write_asset_yaml_stub(
                svg_yaml_path,
                svg_asset_id,
                svg_sha,
                deck_stem,
                slide_number,
                kind="vector",
                colors_hex=svg_colors,
                recolor_targets=svg_colors,
            )
            extracted.append({"atom": svg_asset_id, "kind": "vector", **geom})

    return extracted


# ---------------------------------------------------------------------------
# v4: per-deck theme extraction
# ---------------------------------------------------------------------------
#
# Reads the deck's primary theme (clrScheme + fontScheme) into a
# theme.yaml-shaped dict. INFORMATIONAL — never enforced at compose
# time. Ships in the consumer bundle so the agent has structured
# access to the source palette/fonts without having to guess.
#
# clrScheme has 12 named slots: 4 dk/lt + 6 accents + 2 hyperlink.
# Each slot wraps a colour element — usually <a:srgbClr val="XXXXXX"/>
# but sometimes <a:sysClr val="windowText" lastClr="000000"/>.

_PALETTE_SLOTS = (
    "dk1", "lt1", "dk2", "lt2",
    "accent1", "accent2", "accent3", "accent4", "accent5", "accent6",
    "hlink", "folHlink",
)


def _resolve_color_slot(slot_el) -> str:
    """Given a clrScheme slot element, return the resolved hex (#RRGGBB)
    or empty string if unresolvable."""
    for child in slot_el:
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag == "srgbClr":
            val = child.get("val", "")
            if val:
                return f"#{val.upper()}"
        elif tag == "sysClr":
            last = child.get("lastClr") or child.get("val")
            if last and last != "windowText" and last != "window":
                return f"#{last.upper()}"
    return ""


def extract_deck_theme(prs, deck_stem: str) -> dict:
    """Extract clrScheme + fontScheme from the deck's primary theme.

    Returns a dict matching authoring/schemas/theme.yaml. Always returns
    a complete shape with empty slots when extraction fails so the
    consumer schema is stable.
    """
    out: dict = {
        "deck": deck_stem,
        "palette": {slot: "" for slot in _PALETTE_SLOTS},
        "aliases": {"primary": "", "accent": "", "text": "", "background": ""},
        "fonts": {"major": "", "minor": ""},
        "aspect": aspect_ratio(prs.slide_width or 0, prs.slide_height or 0),
        "sources": [{"deck": deck_stem}],
    }

    masters = list(prs.slide_masters)
    if not masters:
        return out

    try:
        theme_part = masters[0].part.part_related_by(RT.THEME)
    except KeyError:
        return out

    try:
        theme_root = ET.fromstring(theme_part.blob)
    except ET.ParseError:
        return out

    clr_scheme = theme_root.find(f".//{{{_DML_NS}}}clrScheme")
    if clr_scheme is not None:
        for slot in _PALETTE_SLOTS:
            slot_el = clr_scheme.find(f"{{{_DML_NS}}}{slot}")
            if slot_el is None:
                continue
            hex_val = _resolve_color_slot(slot_el)
            if hex_val:
                out["palette"][slot] = hex_val

    font_scheme = theme_root.find(f".//{{{_DML_NS}}}fontScheme")
    if font_scheme is not None:
        for role in ("major", "minor"):
            latin = font_scheme.find(
                f"{{{_DML_NS}}}{role}Font/{{{_DML_NS}}}latin"
            )
            if latin is not None:
                out["fonts"][role] = latin.get("typeface", "")

    # Aliases — simple defaults for B1. Later iterations can detect
    # dark-on-light vs light-on-dark layouts and flip text/background.
    palette = out["palette"]
    out["aliases"]["text"] = "dk1"
    out["aliases"]["background"] = "lt1"
    out["aliases"]["accent"] = "accent1"
    out["aliases"]["primary"] = "dk2" if palette.get("dk2") else "accent1"

    return out


# ---------------------------------------------------------------------------
# Validate
# ---------------------------------------------------------------------------


SLIDE_REQUIRED = ("intent", "feel", "suitable_for")
# Base required-for-all-assets. `composition` is intentionally NOT in
# this list — it's only required for picture-kinds (see ASSET_PICTURE_KINDS
# below). Structured atoms (table, chart, callout, freeform, smartart)
# can have empty composition because the concept doesn't apply.
ASSET_REQUIRED = ("kind", "subject", "feel", "colors", "scope", "suitable_for")
# Kinds for which `composition` is meaningful and required. Mirrors the
# guidance in describe_asset.md ("the slot applies primarily to pictures").
ASSET_PICTURE_KINDS = frozenset(
    {"photo", "icon", "logo", "illustration", "screenshot", "vector"}
)

VOCAB_PATH = SCHEMAS / "vocab.yaml"


def _load_vocab() -> dict:
    with VOCAB_PATH.open("r", encoding="utf-8") as f:
        return yaml.safe_load(f)


VOCAB = _load_vocab()
SLIDE_FEEL_ENUM = set(VOCAB["slide"]["feel"])
SLIDE_SUITABLE_ENUM = set(VOCAB["slide"]["suitable_for"])
ASSET_KIND_ENUM = set(VOCAB["asset"]["kind"])
ASSET_FEEL_ENUM = set(VOCAB["asset"]["feel"])
ASSET_COMPOSITION_ENUM = set(VOCAB["asset"]["composition"])
ASSET_SUITABLE_ENUM = set(VOCAB["asset"]["suitable_for"])
ASSET_SCOPE_PREFIXES = set(VOCAB["asset"]["scope_prefixes"])


def _missing_or_empty(data: dict, keys: Iterable[str]) -> list[str]:
    missing = []
    for k in keys:
        v = data.get(k)
        if v is None or v == "" or v == []:
            missing.append(k)
    return missing


def validate_slide(data: dict) -> list[str]:
    errors = _missing_or_empty(data, SLIDE_REQUIRED)
    feel = data.get("feel")
    if feel and feel not in SLIDE_FEEL_ENUM:
        errors.append(f"feel '{feel}' not in {sorted(SLIDE_FEEL_ENUM)}")
    sfor = data.get("suitable_for") or []
    if isinstance(sfor, list):
        bad = [t for t in sfor if t not in SLIDE_SUITABLE_ENUM]
        if bad:
            errors.append(f"suitable_for has unknown tag(s): {bad}")
    intent = data.get("intent") or ""
    if intent and len(intent.split()) > 25:
        errors.append("intent is over 25 words; prefer <=20")
    return errors


def validate_asset(data: dict) -> list[str]:
    errors = _missing_or_empty(data, ASSET_REQUIRED)
    kind = data.get("kind") or ""
    if kind and kind not in ASSET_KIND_ENUM:
        errors.append(f"kind '{kind}' not in {sorted(ASSET_KIND_ENUM)}")
    if data.get("feel") and data["feel"] not in ASSET_FEEL_ENUM:
        errors.append(f"feel '{data['feel']}' not in {sorted(ASSET_FEEL_ENUM)}")
    # `composition` is required for picture-kinds only; structured atoms
    # (table, chart, callout, freeform, smartart) can leave it empty.
    composition = data.get("composition") or ""
    if kind in ASSET_PICTURE_KINDS and not composition:
        errors.append("composition is required for picture-kind assets")
    if composition and composition not in ASSET_COMPOSITION_ENUM:
        errors.append(
            f"composition '{composition}' not in {sorted(ASSET_COMPOSITION_ENUM)}"
        )
    sfor = data.get("suitable_for") or []
    if isinstance(sfor, list):
        bad = [t for t in sfor if t not in ASSET_SUITABLE_ENUM]
        if bad:
            errors.append(f"suitable_for has unknown tag(s): {bad}")
    scope = data.get("scope") or []
    if isinstance(scope, list):
        for entry in scope:
            if entry == "generic":
                continue
            if not isinstance(entry, str) or ":" not in entry:
                errors.append(
                    f"scope entry {entry!r} must be 'generic' or '<prefix>:<value>'"
                )
                continue
            prefix, _, value = entry.partition(":")
            if prefix not in ASSET_SCOPE_PREFIXES:
                errors.append(
                    f"scope entry {entry!r}: unknown prefix '{prefix}' "
                    f"(use one of {sorted(ASSET_SCOPE_PREFIXES)} or 'generic')"
                )
            if not value:
                errors.append(f"scope entry {entry!r}: missing value after ':'")
    subj = data.get("subject") or ""
    if subj and len(subj.split()) > 30:
        errors.append("subject is over 30 words; prefer <=25")
    return errors


# ---------------------------------------------------------------------------
# Discovery helpers
# ---------------------------------------------------------------------------


def iter_slide_yamls() -> Iterable[Path]:
    decks = WORKSPACE / "decks"
    if not decks.exists():
        return []
    return sorted(decks.glob("*/slides/slide_*.yaml"))


def iter_asset_yamls() -> Iterable[Path]:
    assets = WORKSPACE / "assets"
    if not assets.exists():
        return []
    return sorted(assets.glob("*.yaml"))


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------


@click.group()
def cli() -> None:
    """pptx-skill authoring CLI."""


# --- ingest ----------------------------------------------------------------


class IngestCollisionError(Exception):
    """Raised when a deck dir with the target stem already exists.

    The CLI shells out to a fresh process per ingest so collisions are
    rare in practice, but the Flask app's /api/ingest endpoint chose a
    reject-on-collision policy to keep accidental re-uploads from
    silently overwriting state — callers can catch this and surface it.
    """

    def __init__(self, deck_stem: str, existing_dir: Path):
        super().__init__(
            f"deck {deck_stem!r} already ingested at {existing_dir}"
        )
        self.deck_stem = deck_stem
        self.existing_dir = existing_dir


def _ingest_pptx(deck: Path, *, reject_collision: bool = False) -> dict:
    """Strip a deck into workspace/ as slide fragments + asset binaries.

    Pure-Python helper used by the CLI command and the Flask /api/ingest
    endpoint. Returns a dict with ingest stats so callers can show the
    user what landed without re-walking the workspace.

    ``reject_collision``: if True and ``workspace/decks/<stem>/`` already
    exists, raise IngestCollisionError instead of re-ingesting. The CLI
    leaves this False (re-ingest is a normal flow); the upload endpoint
    sets True so accidental double-uploads surface a clear error.
    """
    ensure_workspace()
    deck_stem = deck.stem
    deck_dir = WORKSPACE / "decks" / deck_stem
    if reject_collision and deck_dir.exists():
        raise IngestCollisionError(deck_stem, deck_dir)
    slides_dir = deck_dir / "slides"
    slides_dir.mkdir(parents=True, exist_ok=True)

    # Snapshot the original so future ingests are reproducible.
    original = deck_dir / "original.pptx"
    if not original.exists() or original.read_bytes() != deck.read_bytes():
        shutil.copyfile(deck, original)

    prs = Presentation(str(original))
    slide_w = prs.slide_width or 9144000
    slide_h = prs.slide_height or 6858000

    # v4: per-deck theme.yaml (palette + fonts). Regenerated on every
    # ingest — it is derived data, no human-edited fields to preserve.
    theme = extract_deck_theme(prs, deck_stem)
    write_yaml(deck_dir / "theme.yaml", theme)

    # v5 (additive): write workspace/themes/<deck>/{theme.yaml, master.pptx}
    # with semantic palette roles + master fragment for the new build engine.
    # Purely additive — v4 path above is unaffected.
    ingest_v5.digest_theme(original, deck_stem, theme, WORKSPACE / "themes")

    # Resolve slide-level theme_colors via aliases; pass palette into
    # slot detection so per-run colour refs become hex + role names.
    theme_palette = theme.get("palette") or {}
    theme_aliases = theme.get("aliases") or {}
    slide_theme_colors = {
        role: theme_palette.get(theme_aliases.get(role, ""), "")
        for role in ("primary", "accent", "text", "background")
    }
    # Strip empty entries so the YAML stays terse on fully-resolved
    # decks but explicit on partial extractions.
    slide_theme_colors = {k: v for k, v in slide_theme_colors.items() if v}
    slide_fonts = {k: v for k, v in (theme.get("fonts") or {}).items() if v}

    assets_dir = WORKSPACE / "assets"
    all_atom_ids: set[str] = set()
    all_picture_ids: set[str] = set()

    n_slides = len(prs.slides)
    for idx in range(n_slides):
        slide = prs.slides[idx]
        slot_defs, renames = detect_slots(slide, slide_w, slide_h, theme_palette)
        layout = infer_layout(slide, slot_defs, renames, slide_w, slide_h)

        slide_number = idx + 1
        slide_id = f"{deck_stem}_{slide_number:02d}"
        slide_pptx = slides_dir / f"slide_{slide_number:02d}.pptx"
        slide_yaml = slides_dir / f"slide_{slide_number:02d}.yaml"

        write_slide_fragment(original, idx, slide_pptx, renames)
        # Extract atoms BEFORE writing the slide yaml so we can populate
        # inventory with the per-shape positional info this template
        # carries (so the agent knows "picking this template gives you
        # these atoms for free at these positions, or you can address
        # them individually in compose mode").
        pic_entries = extract_picture_assets(
            slide, deck_stem, slide_number, assets_dir, slide_w, slide_h,
        )
        atom_entries = extract_structured_atoms(
            slide, deck_stem, slide_number, assets_dir, slide_w, slide_h,
            slot_shape_ids=set(renames.keys()),
        )
        # v4.2: inventory unions pictures + structured atoms with
        # per-shape geometry — gives the compose-time agent slide
        # anatomy beyond just the (placeholder-only) `slots` field.
        inventory = pic_entries + atom_entries
        all_picture_ids.update(e["atom"] for e in pic_entries if e.get("atom"))
        all_atom_ids.update(e["atom"] for e in atom_entries if e.get("atom"))
        write_slide_yaml_stub(
            slide_yaml,
            slide_id,
            layout,
            slot_defs,
            deck_stem,
            slide_number,
            theme_colors=slide_theme_colors,
            fonts=slide_fonts,
            inventory=inventory,
        )

        # v5 (additive): structural skeleton with fractional geometry,
        # font_role / color_role, per-kind constraints. v4 path above
        # is unaffected; v5 writes to workspace/skeletons/<deck>_<NN>/.
        v4_preview = slide_pptx.with_suffix(".png")
        ingest_v5.digest_skeleton(
            slide, slide_w, slide_h, deck_stem, slide_number,
            theme, WORKSPACE / "skeletons", v4_preview_path=v4_preview,
        )

    return {
        "deck_stem": deck_stem,
        "slides": n_slides,
        "pictures": len(all_picture_ids),
        "atoms": len(all_atom_ids),
        "slides_dir": str(slides_dir),
    }


@cli.command()
@click.argument("deck", type=click.Path(exists=True, dir_okay=False, path_type=Path))
def ingest(deck: Path) -> None:
    """Strip a deck into workspace/ as slide fragments + asset binaries."""
    result = _ingest_pptx(deck)
    click.echo(
        f"Ingested {deck.name}: {result['slides']} slide(s) → {result['slides_dir']}"
    )


# --- status ----------------------------------------------------------------


def _bucket(paths: Iterable[Path]) -> tuple[list[Path], list[Path], list[Path]]:
    pending, done, locked = [], [], []
    for p in paths:
        try:
            data = read_yaml(p)
        except Exception:
            pending.append(p)
            continue
        status = data.get("status", "pending")
        if status == "done":
            done.append(p)
        elif status == "locked":
            locked.append(p)
        else:
            pending.append(p)
    return pending, done, locked


@cli.command()
def status() -> None:
    """Counts pending/done/locked; lists pending file paths."""
    slides = list(iter_slide_yamls())
    assets = list(iter_asset_yamls())
    sp, sd, sl = _bucket(slides)
    ap, ad, al = _bucket(assets)
    click.echo(f"slides:  pending={len(sp)} done={len(sd)} locked={len(sl)} total={len(slides)}")
    click.echo(f"assets:  pending={len(ap)} done={len(ad)} locked={len(al)} total={len(assets)}")
    if sp or ap:
        click.echo("\npending:")
        for p in sp + ap:
            click.echo(f"  {p.relative_to(HERE)}")


# --- next ------------------------------------------------------------------


@cli.command(name="next")
@click.option("--kind", type=click.Choice(["asset", "slide"]))
@click.option("--open", "open_editor", is_flag=True, help="Launch $EDITOR on the YAML.")
def next_cmd(kind, open_editor: bool) -> None:
    """Print the next pending item path."""
    if kind == "asset":
        candidates = list(iter_asset_yamls())
    elif kind == "slide":
        candidates = list(iter_slide_yamls())
    else:
        candidates = list(iter_slide_yamls()) + list(iter_asset_yamls())

    target: Path | None = None
    for p in candidates:
        try:
            data = read_yaml(p)
        except Exception:
            target = p
            break
        if data.get("status", "pending") == "pending":
            target = p
            break

    if target is None:
        click.echo("nothing pending")
        return

    click.echo(str(target))

    if open_editor:
        editor = os.environ.get("EDITOR", "vi")
        try:
            subprocess.run([editor, str(target)], check=False)
        except FileNotFoundError:
            click.echo(f"warn: $EDITOR not found ({editor})", err=True)


# --- prompt ----------------------------------------------------------------


@cli.command()
@click.option("--kind", type=click.Choice(["asset", "slide"]), default="asset")
def prompt(kind: str) -> None:
    """Print the bundled describe prompt to stdout for copy/paste."""
    path = PROMPTS / f"describe_{kind}.md"
    if not path.exists():
        raise click.ClickException(f"missing prompt file: {path}")
    click.echo(path.read_text(encoding="utf-8"))


# --- validate --------------------------------------------------------------


def check_prompt_drift() -> list[str]:
    """Warn if a vocab enum value is missing from its prompt markdown."""
    pairs = [
        (PROMPTS / "describe_slide.md", "slide.feel", VOCAB["slide"]["feel"]),
        (PROMPTS / "describe_slide.md", "slide.suitable_for", VOCAB["slide"]["suitable_for"]),
        (PROMPTS / "describe_asset.md", "asset.kind", VOCAB["asset"]["kind"]),
        (PROMPTS / "describe_asset.md", "asset.feel", VOCAB["asset"]["feel"]),
        (PROMPTS / "describe_asset.md", "asset.composition", VOCAB["asset"]["composition"]),
        (PROMPTS / "describe_asset.md", "asset.suitable_for", VOCAB["asset"]["suitable_for"]),
    ]
    errs: list[str] = []
    for path, label, values in pairs:
        if not path.exists():
            errs.append(f"{path.name}: missing prompt file")
            continue
        text = path.read_text(encoding="utf-8")
        missing = [v for v in values if v not in text]
        if missing:
            errs.append(
                f"{path.name}: {label} missing value(s) {missing} "
                f"(update prompt to match schemas/vocab.yaml)"
            )
    return errs


@cli.command()
def validate() -> None:
    """Schema-check every YAML; auto-promote complete ones to done."""
    failures: list[tuple[Path, list[str]]] = []
    promoted = 0
    total = 0

    drift = check_prompt_drift()
    if drift:
        failures.append((VOCAB_PATH, drift))

    for p in iter_slide_yamls():
        total += 1
        try:
            data = read_yaml(p)
        except Exception as e:
            failures.append((p, [f"unreadable: {e}"]))
            continue
        errs = validate_slide(data)
        if errs:
            failures.append((p, errs))
            continue
        if data.get("status") == "pending":
            data["status"] = "done"
            write_yaml(p, data)
            promoted += 1

    for p in iter_asset_yamls():
        total += 1
        try:
            data = read_yaml(p)
        except Exception as e:
            failures.append((p, [f"unreadable: {e}"]))
            continue
        errs = validate_asset(data)
        if errs:
            failures.append((p, errs))
            continue
        if data.get("status") == "pending":
            data["status"] = "done"
            write_yaml(p, data)
            promoted += 1

    if failures:
        click.echo(f"FAIL: {len(failures)} sidecar(s) have errors")
        for p, errs in failures:
            click.echo(f"  {p.relative_to(HERE)}")
            for e in errs:
                click.echo(f"    - {e}")
    click.echo(f"checked {total}, promoted pending→done: {promoted}")
    if failures:
        sys.exit(1)


# --- slide rendering -------------------------------------------------------
#
# Slide fragments are rendered to PNGs by the first available backend.
# Priority: PowerPoint COM (Windows, best fidelity if installed) →
# LibreOffice headless (any OS) → macOS Quick Look. Both `cli.py preview`
# and the bulk-batch flow in `app.py` go through `render_slide_to_png`.


def _render_via_powerpoint_com(slide_pptx: Path, out_png: Path, size: int) -> bool:
    if sys.platform != "win32":
        return False
    ps = shutil.which("powershell") or shutil.which("pwsh")
    if not ps:
        return False
    # PowerShell accepts forward slashes on Windows, sidesteps quoting.
    src = str(slide_pptx.resolve()).replace("\\", "/")
    dst = str(out_png.resolve()).replace("\\", "/")
    height = int(size * 9 / 16)  # assume 16:9; aspect drift fine for thumbnails
    script = (
        '$ErrorActionPreference = "Stop"; '
        '$pp = New-Object -ComObject PowerPoint.Application; '
        f'$deck = $pp.Presentations.Open("{src}", $true, $true, $false); '
        f'$deck.Slides.Item(1).Export("{dst}", "PNG", {size}, {height}); '
        '$deck.Close(); $pp.Quit();'
    )
    try:
        subprocess.run(
            [ps, "-NoProfile", "-NonInteractive", "-Command", script],
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
            timeout=60,
        )
    except (subprocess.SubprocessError, OSError):
        return False
    return out_png.exists()


def _render_via_libreoffice(slide_pptx: Path, out_png: Path, size: int) -> bool:
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        return False
    out_dir = out_png.parent
    try:
        subprocess.run(
            [soffice, "--headless", "--convert-to", "png",
             "--outdir", str(out_dir), str(slide_pptx)],
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
            timeout=120,
        )
    except (subprocess.SubprocessError, OSError):
        return False
    produced = out_dir / f"{slide_pptx.stem}.png"
    if produced != out_png and produced.exists():
        produced.replace(out_png)
    return out_png.exists()


def _render_via_qlmanage(slide_pptx: Path, out_png: Path, size: int) -> bool:
    if sys.platform != "darwin":
        return False
    ql = shutil.which("qlmanage")
    if not ql:
        return False
    out_dir = out_png.parent
    try:
        subprocess.run(
            [ql, "-t", "-s", str(size), "-o", str(out_dir), str(slide_pptx)],
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
            timeout=20,
        )
    except (subprocess.SubprocessError, OSError):
        return False
    # qlmanage writes <name>.pptx.png; normalise to <stem>.png.
    produced = out_dir / f"{slide_pptx.name}.png"
    if produced != out_png and produced.exists():
        produced.replace(out_png)
    return out_png.exists()


_RENDERERS = (
    ("PowerPoint COM (Windows)", _render_via_powerpoint_com),
    ("LibreOffice (soffice)", _render_via_libreoffice),
    ("macOS Quick Look (qlmanage)", _render_via_qlmanage),
)


def render_slide_to_png(
    slide_pptx: Path, out_png: Path | None = None, size: int = 1200
) -> Path | None:
    """Render a single-slide .pptx to a PNG. Returns the PNG path on
    success, None if no backend could produce one. Cached by mtime."""
    if out_png is None:
        out_png = slide_pptx.with_suffix(".png")
    if out_png.exists() and out_png.stat().st_mtime >= slide_pptx.stat().st_mtime:
        return out_png
    for _name, fn in _RENDERERS:
        if fn(slide_pptx, out_png, size):
            return out_png
    return None


def available_renderers() -> list[str]:
    """Human-readable list of slide-rendering backends that *look* usable
    (binaries on PATH; doesn't verify PowerPoint is actually installed)."""
    out: list[str] = []
    if sys.platform == "win32" and (shutil.which("powershell") or shutil.which("pwsh")):
        out.append("PowerPoint COM (Windows; requires PowerPoint installed)")
    if shutil.which("soffice") or shutil.which("libreoffice"):
        out.append("LibreOffice (soffice)")
    if sys.platform == "darwin" and shutil.which("qlmanage"):
        out.append("macOS Quick Look (qlmanage)")
    return out


# --- preview ---------------------------------------------------------------


@cli.command()
def preview() -> None:
    """Best-effort PNG thumbnails of slide fragments. Tries PowerPoint
    (Windows), LibreOffice, then macOS Quick Look in order."""
    if not available_renderers():
        click.echo(
            "preview: no slide renderer available — install LibreOffice, "
            "or use PowerPoint on Windows",
            err=True,
        )
        return

    decks = WORKSPACE / "decks"
    if not decks.exists():
        click.echo("no decks ingested yet")
        return

    n_built = 0
    n_failed = 0
    for slide_pptx in sorted(decks.glob("*/slides/slide_*.pptx")):
        png = render_slide_to_png(slide_pptx)
        if png is not None:
            n_built += 1
        else:
            n_failed += 1
    if n_failed:
        click.echo(f"preview: {n_failed} slide(s) could not be rendered", err=True)
    click.echo(f"preview: built {n_built} thumbnail(s)")


# --- build -----------------------------------------------------------------


def _ext_for(blob_path: Path) -> str:
    return blob_path.suffix.lstrip(".") or "bin"


def _template_index_entry(s: dict) -> dict:
    out = {
        "id": s["id"],
        "intent": s.get("intent", ""),
        "feel": s.get("feel", ""),
        "suitable_for": s.get("suitable_for", []),
        "layout": s.get("layout", ""),
        "slots": s.get("slots", []),
    }
    # v4: include the auto-extracted theme snapshot + atom inventory
    # when present. Omit empties so v3-era templates stay terse.
    # v4.1: `interpretation` — model's speculative observations, info-only;
    # not used for filtering, just surfaced to the compose-time agent.
    for key in ("theme_colors", "fonts", "inventory", "interpretation"):
        v = s.get(key)
        if v:
            out[key] = v
    return out


def _asset_index_entry(a: dict) -> dict:
    out = {
        "id": a["id"],
        "kind": a.get("kind", ""),
        "subject": a.get("subject", ""),
        "depicts": a.get("depicts", ""),
        "feel": a.get("feel", ""),
        "composition": a.get("composition", ""),
        "colors": a.get("colors", []),
        "scope": a.get("scope", []),
        "suitable_for": a.get("suitable_for", []),
    }
    # v4: structured colour + kind-specific blocks. All optional —
    # only emit when present so the index stays focused on retrievable
    # fields.
    # v4.1: `interpretation` — model's speculative observations, info-only;
    # never a filter target, surfaced to the compose-time agent as context.
    for key in (
        "colors_hex", "recolor_targets", "table", "chart", "shape", "smartart",
        "interpretation",
    ):
        v = a.get(key)
        if v:
            out[key] = v
    return out


def build_index(slides: list[dict], assets: list[dict]) -> dict:
    return {
        "templates": [_template_index_entry(s) for s in slides],
        "assets": [_asset_index_entry(a) for a in assets],
    }


@cli.command()
@click.option("--allow-pending", is_flag=True, help="Build even if some sidecars are pending.")
@click.option("--out", "out_path", type=click.Path(path_type=Path), default=None)
@click.option(
    "--no-brand",
    is_flag=True,
    help="Build a policy-stripped bundle: brand.md is omitted and SKILL.md "
    "carries a 'policy disabled' notice. Useful for control-test runs "
    "comparing branded vs un-branded agent behaviour.",
)
def build(allow_pending: bool, out_path: Path | None, no_brand: bool) -> None:
    """Emit dist/skill.zip — the consumer artifact."""
    slide_yamls = list(iter_slide_yamls())
    asset_yamls = list(iter_asset_yamls())

    if not slide_yamls and not asset_yamls:
        raise click.ClickException("workspace is empty; run `ingest` first")

    pending = []
    slides_meta: list[dict] = []
    assets_meta: list[dict] = []

    for p in slide_yamls:
        d = read_yaml(p)
        if d.get("status", "pending") == "pending":
            pending.append(p)
        d["_yaml_path"] = p
        slides_meta.append(d)

    for p in asset_yamls:
        d = read_yaml(p)
        if d.get("status", "pending") == "pending":
            pending.append(p)
        d["_yaml_path"] = p
        assets_meta.append(d)

    if pending and not allow_pending:
        click.echo(f"refusing to build: {len(pending)} sidecar(s) still pending", err=True)
        for p in pending:
            click.echo(f"  {p.relative_to(HERE)}", err=True)
        click.echo("re-run with --allow-pending to override", err=True)
        sys.exit(1)

    DIST.mkdir(parents=True, exist_ok=True)
    zip_path = out_path or (DIST / "skill.zip")

    reader_src = CONSUMER / "reader.py"
    skill_md = CONSUMER / "SKILL.md"
    consumer_reqs = CONSUMER / "requirements.txt"
    brand_md = HERE / "brand.md"
    if not reader_src.exists() or not skill_md.exists():
        raise click.ClickException(
            "consumer/reader.py or SKILL.md missing — cannot build zip"
        )

    index = build_index(slides_meta, assets_meta)

    skill_text = skill_md.read_text(encoding="utf-8")
    if no_brand:
        notice = (
            "> **Policy disabled.** This bundle was built with "
            "`build --no-brand` — `brand.md` is intentionally omitted and "
            "the agent should not assume any deck-style constraints beyond "
            "what individual templates / assets describe.\n\n"
        )
        skill_text = notice + skill_text

    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("SKILL.md", skill_text)
        zf.writestr("reader.py", reader_src.read_text(encoding="utf-8"))
        if consumer_reqs.exists():
            zf.writestr("requirements.txt", consumer_reqs.read_text(encoding="utf-8"))
        zf.writestr("index.json", json.dumps(index, indent=2, ensure_ascii=False))
        if not no_brand and brand_md.exists():
            brand_text = brand_md.read_text(encoding="utf-8").strip()
            if brand_text:
                zf.writestr("brand.md", brand_text + "\n")

        # v4: per-deck theme.yaml (palette + fonts). Informational for
        # the agent. Skipped silently for any deck without one (pre-v4
        # ingest output).
        for theme_yaml in sorted((WORKSPACE / "decks").glob("*/theme.yaml")):
            deck_name = theme_yaml.parent.name
            zf.write(theme_yaml, f"decks/{deck_name}/theme.yaml")

        for sd in slides_meta:
            tid = sd["id"]
            yaml_path: Path = sd["_yaml_path"]
            slide_pptx = yaml_path.with_suffix(".pptx")
            if not slide_pptx.exists():
                click.echo(f"warn: missing {slide_pptx} — skipping {tid}", err=True)
                continue
            zf.write(slide_pptx, f"templates/{tid}/slide.pptx")
            clean = {k: v for k, v in sd.items() if not k.startswith("_")}
            zf.writestr(
                f"templates/{tid}/meta.yaml",
                yaml.safe_dump(clean, sort_keys=False, allow_unicode=True),
            )
            preview_png = yaml_path.with_suffix(".png")
            if preview_png.exists():
                zf.write(preview_png, f"templates/{tid}/preview.png")

        for ad in assets_meta:
            aid = ad["id"]
            sha = ad.get("sha1", "")
            yaml_path: Path = ad["_yaml_path"]
            # Find the binary alongside the yaml — same stem, any extension.
            blob: Path | None = None
            for cand in yaml_path.parent.glob(f"{yaml_path.stem}.*"):
                if cand.suffix == ".yaml":
                    continue
                blob = cand
                break
            if blob is None:
                click.echo(f"warn: no binary for asset {aid} ({sha}) — skipping", err=True)
                continue
            ext = _ext_for(blob)
            zf.write(blob, f"assets/{aid}.{ext}")
            clean = {k: v for k, v in ad.items() if not k.startswith("_")}
            zf.writestr(
                f"assets/{aid}.yaml",
                yaml.safe_dump(clean, sort_keys=False, allow_unicode=True),
            )

    brand_tag = " (no brand)" if no_brand else ""
    click.echo(
        f"built {zip_path}{brand_tag} — "
        f"{len(slides_meta)} template(s), {len(assets_meta)} asset(s)"
    )


# --- package-app ----------------------------------------------------------
#
# Produces a zip containing ONLY the application code + scaffolding —
# nothing organisation-specific. Use it to hand the tool to a teammate
# without leaking KB content, built artifacts, source decks, or
# org-tuned brand rules. Allowlist-based: anything not explicitly named
# below is excluded.

# Files included verbatim (path is relative to repo root). Globs allowed.
PACKAGE_APP_ALLOWLIST = (
    "README.md",
    ".gitignore",
    "authoring/cli.py",
    "authoring/app.py",
    "authoring/requirements.txt",
    "authoring/prompts/*.md",
    "authoring/schemas/*.yaml",
    "authoring/schemas/*.yml",
    "consumer/SKILL.md",
    "consumer/reader.py",
    "consumer/requirements.txt",
    "consumer/index.example.json",
    "consumer/helpers/*.py",
    "consumer/helpers/*.md",
)


@cli.command(name="package-app")
@click.option("--out", "out_path", type=click.Path(path_type=Path), default=None,
              help="Output zip path. Default: authoring/dist/pptx-skill-app.zip")
def package_app(out_path: Path | None) -> None:
    """Build a transferable zip of the app code only.

    Never includes workspace/, dist/, source .pptx files, brand.md, or
    other session/local-only files. Ships `brand.example.md` so the
    recipient knows where to put their own brand rules.
    """
    REPO_ROOT = HERE.parent
    out_path = out_path or (DIST / "pptx-skill-app.zip")
    DIST.mkdir(parents=True, exist_ok=True)

    seen: set[str] = set()
    with zipfile.ZipFile(out_path, "w", zipfile.ZIP_DEFLATED) as zf:
        for pattern in PACKAGE_APP_ALLOWLIST:
            matches = sorted(REPO_ROOT.glob(pattern))
            if not matches:
                continue
            for src in matches:
                if not src.is_file():
                    continue
                arc = src.relative_to(REPO_ROOT).as_posix()
                if arc in seen:
                    continue
                seen.add(arc)
                zf.write(src, arc)

        # Brand stub goes in as brand.example.md so the recipient can copy
        # it to brand.md and fill in. Always written from a clean template,
        # never from the local (possibly org-customised) brand.md.
        zf.writestr(
            "authoring/brand.example.md",
            _BRAND_STUB,
        )

        # Empty workspace dir + a README pointer for first-time users.
        zf.writestr("authoring/workspace/.gitkeep", "")
        zf.writestr("FIRST_RUN.md", _FIRST_RUN_NOTES)

    size_kb = out_path.stat().st_size // 1024
    click.echo(f"packaged {out_path} — {len(seen) + 3} file(s), {size_kb} KB")


_BRAND_STUB = """# Brand & visual rules

Copy this file to `authoring/brand.md` and fill in. The compose page
auto-includes brand.md in every prompt bundle so the LLM sees these
rules before reading the brief.

Keep it short and prescriptive. Long preambles get ignored.

## Palette

- Primary: <e.g. #0A2540 navy>
- Accent:  <e.g. #F26B38 orange>
- Avoid:   <colors that look off-brand>

## Voice

- <e.g. "factual, no hype, third person">

## Don'ts

- <e.g. "no exclamation marks in headlines">
- <e.g. "never use stock people-in-suits imagery">
"""


_FIRST_RUN_NOTES = """# First run

This zip contains the pptx-skill app code only. No KB content, no
built artifacts. Steps to get going on a fresh machine:

1. Unzip somewhere on disk.
2. `pip install -r authoring/requirements.txt`
3. `cp authoring/brand.example.md authoring/brand.md` and edit to your
   org's palette / voice / taboos.
4. Drop a source deck somewhere and `python3 authoring/cli.py ingest path/to/deck.pptx`.
5. Describe slides/assets via `python3 authoring/app.py` (compose page
   served on http://127.0.0.1:5050) or hand-edit YAML in
   `authoring/workspace/`.
6. `python3 authoring/cli.py validate` and then `... build` to produce
   the durable consumer skill.zip, or use the compose page to generate
   per-deck agent prompt bundles.

See README.md for the high-level overview. The consumer/SKILL.md file
inside is the canonical reference the agent reads.
"""


if __name__ == "__main__":
    cli()
