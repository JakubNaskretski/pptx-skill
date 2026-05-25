"""Regression tests for the multi-slide compose pipeline.

Covers the two bugs fixed alongside this file:

* Bug 1: rels corruption when grafting 2+ slides — `_copy_slide_into`
  used to deepcopy shape XML with rId references but never imported
  the underlying parts onto the destination slide, so PowerPoint
  showed the "couldn't read some content" repair dialog and dropped
  shapes.
* Bug 2: backgrounds were lost on grafted slides because the
  destination's blank layout supplied a blank background and the
  source's inheritance chain was never carried over.

Tests build synthetic decks in-memory rather than relying on any
ingested workspace fixture.
"""
from __future__ import annotations

import copy
import io
import re
import sys
import tempfile
import unittest
import zipfile
from pathlib import Path

REPO = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(REPO / "consumer"))

from lxml import etree  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.util import Inches  # noqa: E402

import reader as reader_mod  # noqa: E402


_PNG_1x1 = bytes.fromhex(
    "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c489"
    "0000000d49444154789c63f8cfc0000000030001ff1a4d6b0000000049454e44ae"
    "426082"
)


def _make_source_slide_with_picture(tmpdir: Path, name: str) -> Path:
    """Build a one-slide .pptx with a picture and a master-level
    gradient background. Returns the on-disk path."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Picture → rId on the slide rels.
    slide.shapes.add_picture(
        io.BytesIO(_PNG_1x1), Inches(1), Inches(1), Inches(2), Inches(2),
    )
    # Auto-shape with no rId, exercises the non-picture branch.
    slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(4), Inches(1), Inches(2), Inches(1),
    )
    # Master-level <p:bg> so flatten has something to inherit.
    master = slide.slide_layout.slide_master
    master_xml = master._element
    cSld = master_xml.find(reader_mod._qn("p:cSld"))
    if reader_mod._find_bg_element(cSld) is None:
        bg_xml = (
            "<p:bg xmlns:p='http://schemas.openxmlformats.org/presentationml/2006/main' "
            "xmlns:a='http://schemas.openxmlformats.org/drawingml/2006/main'>"
            "<p:bgPr>"
            "<a:gradFill flip='none' rotWithShape='1'>"
            "<a:gsLst>"
            "<a:gs pos='0'><a:srgbClr val='112233'/></a:gs>"
            "<a:gs pos='100000'><a:srgbClr val='AABBCC'/></a:gs>"
            "</a:gsLst>"
            "<a:lin ang='5400000' scaled='0'/>"
            "</a:gradFill>"
            "</p:bgPr>"
            "</p:bg>"
        )
        cSld.insert(0, etree.fromstring(bg_xml))
    out = tmpdir / f"{name}.pptx"
    prs.save(str(out))
    return out


def _broken_rids(pptx_bytes: bytes) -> list[tuple[str, str]]:
    broken: list[tuple[str, str]] = []
    with zipfile.ZipFile(io.BytesIO(pptx_bytes)) as z:
        for name in z.namelist():
            m = re.fullmatch(r"ppt/slides/(slide\d+)\.xml", name)
            if not m:
                continue
            slide_name = m.group(1)
            xml = z.read(name).decode("utf-8")
            used = set(re.findall(r'r:(?:embed|link|id)="(rId\d+)"', xml))
            try:
                rels = z.read(
                    f"ppt/slides/_rels/{slide_name}.xml.rels"
                ).decode("utf-8")
                have = set(re.findall(r'Id="(rId\d+)"', rels))
            except KeyError:
                have = set()
            for rid in used - have:
                broken.append((slide_name, rid))
    return broken


def _slide_has_explicit_bg(pptx_bytes: bytes, slide_name: str) -> bool:
    with zipfile.ZipFile(io.BytesIO(pptx_bytes)) as z:
        xml = z.read(f"ppt/slides/{slide_name}.xml").decode("utf-8")
        return "<p:bg" in xml


def _make_destination(tmpdir: Path, host_src: Path) -> Presentation:
    """Open the host as a Presentation we can graft onto. Mirrors
    cmd_compose's pattern (copy host pptx to a temp file, open it)."""
    import shutil
    host_copy = tmpdir / "host_open.pptx"
    shutil.copyfile(host_src, host_copy)
    return Presentation(str(host_copy))


class TestMultiSlideCopyRels(unittest.TestCase):
    def test_three_slide_graft_has_no_broken_rids(self):
        with tempfile.TemporaryDirectory() as td:
            tmp = Path(td)
            host = _make_source_slide_with_picture(tmp, "host")
            src_b = _make_source_slide_with_picture(tmp, "src_b")
            src_c = _make_source_slide_with_picture(tmp, "src_c")
            dest = _make_destination(tmp, host)
            reader_mod._copy_slide_into(dest, src_b)
            reader_mod._copy_slide_into(dest, src_c)
            out = tmp / "out.pptx"
            dest.save(str(out))
            broken = _broken_rids(out.read_bytes())
            self.assertEqual(
                broken, [],
                f"shape references survived without backing rels: {broken}",
            )

    def test_grafted_slide_keeps_its_picture(self):
        with tempfile.TemporaryDirectory() as td:
            tmp = Path(td)
            host = _make_source_slide_with_picture(tmp, "host")
            src = _make_source_slide_with_picture(tmp, "src")
            dest = _make_destination(tmp, host)
            reader_mod._copy_slide_into(dest, src)
            out = tmp / "out.pptx"
            dest.save(str(out))
            with zipfile.ZipFile(out) as z:
                slide2 = z.read("ppt/slides/slide2.xml").decode("utf-8")
                self.assertIn("<p:pic", slide2)
                self.assertIn("<a:blip ", slide2)
                # Each picture refs a media part — the rel must resolve.
                used = set(re.findall(r'r:embed="(rId\d+)"', slide2))
                rels = z.read("ppt/slides/_rels/slide2.xml.rels").decode("utf-8")
                have = set(re.findall(r'Id="(rId\d+)"', rels))
                self.assertTrue(used <= have, f"unresolved embed rels: {used - have}")


class TestBackgroundFlatten(unittest.TestCase):
    def test_grafted_slide_carries_explicit_bg(self):
        with tempfile.TemporaryDirectory() as td:
            tmp = Path(td)
            host = _make_source_slide_with_picture(tmp, "host")
            src = _make_source_slide_with_picture(tmp, "src_with_bg")
            dest = _make_destination(tmp, host)
            reader_mod._copy_slide_into(dest, src)
            out = tmp / "out.pptx"
            dest.save(str(out))
            self.assertTrue(
                _slide_has_explicit_bg(out.read_bytes(), "slide2"),
                "grafted slide should carry an explicit <p:bg> after flatten",
            )

    def test_flatten_is_idempotent(self):
        with tempfile.TemporaryDirectory() as td:
            tmp = Path(td)
            src = _make_source_slide_with_picture(tmp, "src")
            prs = Presentation(str(src))
            slide = prs.slides[0]
            reader_mod._flatten_slide_background(slide)
            xml_first = etree.tostring(
                slide._element.find(reader_mod._qn("p:cSld")),
                encoding="unicode",
            )
            reader_mod._flatten_slide_background(slide)
            xml_second = etree.tostring(
                slide._element.find(reader_mod._qn("p:cSld")),
                encoding="unicode",
            )
            self.assertEqual(xml_first, xml_second)

    def test_flatten_walks_to_master(self):
        with tempfile.TemporaryDirectory() as td:
            tmp = Path(td)
            src = _make_source_slide_with_picture(tmp, "src")
            prs = Presentation(str(src))
            slide = prs.slides[0]
            cSld = slide._element.find(reader_mod._qn("p:cSld"))
            self.assertIsNone(
                reader_mod._find_bg_element(cSld),
                "synthetic slide should start with no <p:bg> (bg lives on master)",
            )
            reader_mod._flatten_slide_background(slide)
            self.assertIsNotNone(
                reader_mod._find_bg_element(cSld),
                "flatten should materialise the master's <p:bg> on the slide",
            )


class TestMediaPartnameCollision(unittest.TestCase):
    """Both synthetic sources name their picture `ppt/media/image1.png`.
    A naive `relate_to` would write two zip entries with that same name
    (PowerPoint then reads only the first). Verify the import path
    either dedupes by content or renames."""

    def test_no_duplicate_zip_entries(self):
        with tempfile.TemporaryDirectory() as td:
            tmp = Path(td)
            host = _make_source_slide_with_picture(tmp, "host")
            src_b = _make_source_slide_with_picture(tmp, "src_b")
            src_c = _make_source_slide_with_picture(tmp, "src_c")
            dest = _make_destination(tmp, host)
            reader_mod._copy_slide_into(dest, src_b)
            reader_mod._copy_slide_into(dest, src_c)
            out = tmp / "out.pptx"
            dest.save(str(out))
            with zipfile.ZipFile(out) as z:
                from collections import Counter
                names = z.namelist()
                dupes = {n: c for n, c in Counter(names).items() if c > 1}
                self.assertEqual(dupes, {}, f"zip has duplicate entries: {dupes}")


class TestCanOpenInPythonPptx(unittest.TestCase):
    """Soft equivalent of 'PowerPoint opens cleanly' — python-pptx's
    parser raises on most malformed parts."""

    def test_three_slide_output_parses(self):
        with tempfile.TemporaryDirectory() as td:
            tmp = Path(td)
            host = _make_source_slide_with_picture(tmp, "host")
            src_b = _make_source_slide_with_picture(tmp, "src_b")
            src_c = _make_source_slide_with_picture(tmp, "src_c")
            dest = _make_destination(tmp, host)
            reader_mod._copy_slide_into(dest, src_b)
            reader_mod._copy_slide_into(dest, src_c)
            out = tmp / "out.pptx"
            dest.save(str(out))
            reopened = Presentation(str(out))
            self.assertEqual(len(reopened.slides), 3)


if __name__ == "__main__":
    unittest.main()
