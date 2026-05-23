"""Tests for pptx-skill v4 phases D3–D5, E, and v4.1 (font remap + groups).

Run from repo root: ``python3 -m unittest tests.test_v4``

The suite is unittest-only (no pytest dep). Synthetic decks are built
in-memory via python-pptx rather than committed as binary fixtures.
"""
from __future__ import annotations

import json
import subprocess
import sys
import tempfile
import unittest
import zipfile
from io import BytesIO
from pathlib import Path

REPO = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(REPO / "authoring"))
sys.path.insert(0, str(REPO / "consumer"))

from pptx import Presentation  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.util import Inches  # noqa: E402

import cli as cli_mod  # noqa: E402
import reader as reader_mod  # noqa: E402


# ---------------------------------------------------------------------------
# Phase D3 — detect_slots picks up table placeholders / freestanding stays atom
# ---------------------------------------------------------------------------


class TestDetectSlots(unittest.TestCase):
    def test_text_placeholder_becomes_text_slot(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "hi"
        slots, _ = cli_mod.detect_slots(slide, prs.slide_width, prs.slide_height, {})
        kinds = {s["kind"] for s in slots}
        self.assertIn("text", kinds)

    def test_freestanding_table_is_not_a_slot(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.add_table(
            rows=2, cols=2,
            left=Inches(1), top=Inches(1),
            width=Inches(4), height=Inches(2),
        )
        slots, _ = cli_mod.detect_slots(slide, prs.slide_width, prs.slide_height, {})
        # Free-standing tables stay as atoms (extract_structured_atoms).
        self.assertFalse(any(s.get("kind") == "table" for s in slots))


# ---------------------------------------------------------------------------
# Phase D3 — _fill_table_shape semantics
# ---------------------------------------------------------------------------


class TestFillTableShape(unittest.TestCase):
    def _slide_with_table(self, rows=3, cols=2):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        ts = slide.shapes.add_table(
            rows=rows, cols=cols,
            left=Inches(1), top=Inches(1),
            width=Inches(4), height=Inches(2),
        )
        ts.name = "data"
        for r in range(rows):
            for c in range(cols):
                ts.table.cell(r, c).text = f"src_r{r}c{c}"
        return slide, ts

    def test_partial_fill_leaves_extra_rows_untouched(self):
        slide, ts = self._slide_with_table()
        ws = reader_mod._apply_slot_value(
            slide, "data", [["A", "B"], ["X", "Y"]], "table"
        )
        self.assertEqual(ws, [])
        self.assertEqual(ts.table.cell(0, 0).text, "A")
        self.assertEqual(ts.table.cell(1, 1).text, "Y")
        self.assertEqual(ts.table.cell(2, 0).text, "src_r2c0")

    def test_too_many_rows_truncates_with_warning(self):
        slide, _ = self._slide_with_table(rows=2, cols=2)
        ws = reader_mod._apply_slot_value(slide, "data", [["a", "b"]] * 5, "table")
        self.assertTrue(any("truncating" in w for w in ws), ws)

    def test_too_many_cols_truncates_with_warning(self):
        slide, _ = self._slide_with_table(rows=2, cols=2)
        ws = reader_mod._apply_slot_value(slide, "data", [["a", "b", "c"]], "table")
        self.assertTrue(any("template has 2; truncating" in w for w in ws), ws)

    def test_dict_form_honors_cells_only(self):
        slide, ts = self._slide_with_table()
        ws = reader_mod._apply_slot_value(
            slide, "data", {"cells": [["dict-A"]], "ignored": 1}, "table"
        )
        self.assertEqual(ts.table.cell(0, 0).text, "dict-A")
        self.assertTrue(any("only honors 'cells'" in w for w in ws))

    def test_bad_value_shape_warns_and_leaves_cells(self):
        slide, ts = self._slide_with_table()
        ws = reader_mod._apply_slot_value(slide, "data", "not a list", "table")
        self.assertTrue(any("expects list-of-lists" in w for w in ws))
        self.assertEqual(ts.table.cell(0, 0).text, "src_r0c0")

    def test_kind_hint_none_falls_back_to_has_table(self):
        slide, ts = self._slide_with_table()
        reader_mod._apply_slot_value(slide, "data", [["auto"]], None)
        self.assertEqual(ts.table.cell(0, 0).text, "auto")


# ---------------------------------------------------------------------------
# Phase D4 — native text + atom placement helpers
# ---------------------------------------------------------------------------


class TestComposeShapes(unittest.TestCase):
    def test_native_text_renders_at_fractional_position(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        ws = reader_mod._place_native_text(
            slide,
            {"value": "Hello", "x": 0.1, "y": 0.1, "w": 0.5, "h": 0.1, "bold": True},
            prs.slide_width,
            prs.slide_height,
        )
        self.assertEqual(ws, [])
        textboxes = [
            s for s in slide.shapes if s.has_text_frame and s.text_frame.text == "Hello"
        ]
        self.assertEqual(len(textboxes), 1)
        # bold honored
        runs = textboxes[0].text_frame.paragraphs[0].runs
        self.assertTrue(runs and runs[0].font.bold)

    def test_native_text_warns_for_unhonored_keys(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        ws = reader_mod._place_native_text(
            slide,
            {"value": "Hi", "color_role": "accent", "font_role": "major"},
            prs.slide_width,
            prs.slide_height,
        )
        self.assertTrue(any("color_role" in w for w in ws))
        self.assertTrue(any("font_role" in w for w in ws))


# ---------------------------------------------------------------------------
# Phase D5 — semantic clrScheme remap + role-target recolor
# ---------------------------------------------------------------------------


class TestSchemeRemap(unittest.TestCase):
    def test_build_remap_only_for_mismatched_aliases(self):
        src = {"aliases": {"primary": "dk2", "accent": "accent3"}}
        host = {"aliases": {"primary": "dk2", "accent": "accent1"}}
        self.assertEqual(reader_mod._build_scheme_remap(src, host), {"accent3": "accent1"})

    def test_build_remap_empty_when_aliases_align(self):
        src = {"aliases": {"primary": "dk2", "accent": "accent1"}}
        host = {"aliases": {"primary": "dk2", "accent": "accent1"}}
        self.assertEqual(reader_mod._build_scheme_remap(src, host), {})

    def test_apply_remap_normalises_tx1_dk1(self):
        from lxml import etree
        xml = (
            b'<root xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            b'<a:schemeClr val="tx1"/></root>'
        )
        el = etree.fromstring(xml)
        count = reader_mod._apply_scheme_remap(el, {"dk1": "accent1"})
        self.assertEqual(count, 1)
        self.assertEqual(el[0].get("val"), "accent1")

    def test_recolor_resolves_role_token_against_host_theme(self):
        from lxml import etree
        host = {
            "palette": {"dk2": "#0E2841", "accent1": "#156082"},
            "aliases": {"primary": "dk2", "accent": "accent1"},
        }
        xml = (
            b'<root xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            b'<a:srgbClr val="ABCDEF"/></root>'
        )
        el = etree.fromstring(xml)
        n, ws = reader_mod._apply_recolor_xml(el, {"#abcdef": "accent"}, host)
        self.assertEqual(n, 1)
        self.assertEqual(el[0].get("val"), "156082")

    def test_recolor_unresolved_role_emits_warning(self):
        from lxml import etree
        host = {"palette": {}, "aliases": {}}
        xml = (
            b'<root xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            b'<a:srgbClr val="ABCDEF"/></root>'
        )
        el = etree.fromstring(xml)
        n, ws = reader_mod._apply_recolor_xml(el, {"#abcdef": "no-such-role"}, host)
        self.assertEqual(n, 0)
        self.assertTrue(any("could not be resolved" in w for w in ws))


# ---------------------------------------------------------------------------
# Phase E — build --no-brand kill switch
# ---------------------------------------------------------------------------


class TestBuildNoBrand(unittest.TestCase):
    def _run_build(self, *extra):
        with tempfile.TemporaryDirectory() as d:
            out = Path(d) / "skill.zip"
            cmd = [
                sys.executable, str(REPO / "authoring/cli.py"),
                "build", "--allow-pending", "--out", str(out), *extra,
            ]
            r = subprocess.run(cmd, capture_output=True, text=True, cwd=str(REPO))
            self.assertEqual(r.returncode, 0, f"build failed:\n{r.stderr}")
            return out.read_bytes()

    def test_no_brand_omits_brand_md_and_adds_notice(self):
        data = self._run_build("--no-brand")
        with zipfile.ZipFile(BytesIO(data)) as z:
            self.assertNotIn("brand.md", z.namelist())
            self.assertIn("Policy disabled", z.read("SKILL.md").decode("utf-8"))

    def test_branded_build_keeps_brand_md_when_present(self):
        if not (REPO / "authoring/brand.md").exists():
            self.skipTest("no brand.md in workspace")
        data = self._run_build()
        with zipfile.ZipFile(BytesIO(data)) as z:
            self.assertIn("brand.md", z.namelist())
            self.assertNotIn("Policy disabled", z.read("SKILL.md").decode("utf-8"))


# ---------------------------------------------------------------------------
# End-to-end: compose regression + mixed-mode + compose-only
# ---------------------------------------------------------------------------


class TestComposeRoundTrip(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        # One build for the whole class — costs ~1s.
        r = subprocess.run(
            [sys.executable, str(REPO / "authoring/cli.py"), "build", "--allow-pending"],
            capture_output=True, text=True, cwd=str(REPO),
        )
        if r.returncode != 0:
            raise unittest.SkipTest(f"build failed:\n{r.stderr}")
        cls.work = Path(tempfile.mkdtemp(prefix="rtb_"))
        with zipfile.ZipFile(REPO / "authoring/dist/skill.zip") as z:
            z.extractall(cls.work)
        cls.idx = json.loads((cls.work / "index.json").read_text())

    def _compose(self, plan, name):
        plan_path = self.work / f"plan_{name}.json"
        out_path = self.work / f"out_{name}.pptx"
        plan_path.write_text(json.dumps(plan))
        r = subprocess.run(
            [sys.executable, str(self.work / "reader.py"), "compose",
             str(plan_path), str(out_path)],
            capture_output=True, text=True,
        )
        self.assertEqual(r.returncode, 0, f"{r.stdout}\n{r.stderr}")
        return Presentation(str(out_path)), json.loads(r.stdout)

    def _first_text_template(self):
        for t in self.idx["templates"]:
            slots = [s for s in t.get("slots", []) if s["kind"] == "text"]
            if slots:
                return t["id"], slots[0]["id"]
        self.skipTest("no text template in bundle")

    def test_v3_pure_template_plan_still_renders(self):
        tid, slot = self._first_text_template()
        prs, _ = self._compose([{"template": tid, "slots": {slot: "regression"}}], "v3")
        self.assertEqual(len(prs.slides), 1)

    def test_compose_only_plan_drops_host_first_slide(self):
        prs, _ = self._compose(
            [{"compose": True, "shapes": [
                {"kind": "text", "value": "alone", "x": 0.1, "y": 0.1, "w": 0.5, "h": 0.1},
            ]}],
            "compose_only",
        )
        self.assertEqual(len(prs.slides), 1)
        texts = [
            shp.text_frame.text for slide in prs.slides for shp in slide.shapes
            if shp.has_text_frame
        ]
        self.assertTrue(any("alone" in t for t in texts), texts)

    def test_mixed_plan_template_then_compose(self):
        tid, slot = self._first_text_template()
        prs, _ = self._compose(
            [
                {"template": tid, "slots": {slot: "first"}},
                {"compose": True, "shapes": [
                    {"kind": "text", "value": "second", "x": 0.1, "y": 0.1, "w": 0.5, "h": 0.1},
                ]},
            ],
            "mixed",
        )
        self.assertEqual(len(prs.slides), 2)

    def test_table_atom_placement_if_available(self):
        table_asset = next(
            (a for a in self.idx["assets"] if a.get("kind") == "table"), None
        )
        if table_asset is None:
            self.skipTest("no table atom in bundle")
        prs, result = self._compose(
            [{"compose": True, "shapes": [
                {"atom": table_asset["id"], "kind": "table",
                 "x": 0.1, "y": 0.1, "w": 0.8, "h": 0.6,
                 "cells": [["roundtrip-A", "roundtrip-B"]]},
            ]}],
            "table_atom",
        )
        slide = prs.slides[-1]
        placed = [s for s in slide.shapes if getattr(s, "has_table", False)]
        self.assertEqual(len(placed), 1)
        t = placed[0].table
        self.assertEqual(t.cell(0, 0).text.strip(), "roundtrip-A")
        self.assertEqual(t.cell(0, 1).text.strip(), "roundtrip-B")


# ---------------------------------------------------------------------------
# v4.1 — surgical theme-font remap (D5 extension)
# ---------------------------------------------------------------------------


class TestFontRemap(unittest.TestCase):
    DML = "http://schemas.openxmlformats.org/drawingml/2006/main"

    def _wrap(self, *inner: str) -> bytes:
        return (
            f'<root xmlns:a="{self.DML}">' + "".join(inner) + "</root>"
        ).encode("utf-8")

    def test_build_remap_only_includes_differing_roles(self):
        # major matches, minor differs → remap only contains the minor entry.
        src = {"fonts": {"major": "Aptos Display", "minor": "Helvetica"}}
        host = {"fonts": {"major": "Aptos Display", "minor": "Aptos"}}
        self.assertEqual(
            reader_mod._build_font_remap(src, host),
            {"helvetica": "Aptos"},
        )

    def test_build_remap_distinct_major_minor(self):
        src = {"fonts": {"major": "Helvetica", "minor": "Century Gothic"}}
        host = {"fonts": {"major": "Aptos Display", "minor": "Aptos"}}
        self.assertEqual(
            reader_mod._build_font_remap(src, host),
            {"helvetica": "Aptos Display", "century gothic": "Aptos"},
        )

    def test_build_remap_empty_when_fonts_match(self):
        src = {"fonts": {"major": "Aptos Display", "minor": "Aptos"}}
        host = {"fonts": {"major": "aptos display", "minor": "APTOS"}}
        self.assertEqual(reader_mod._build_font_remap(src, host), {})

    def test_build_remap_empty_when_either_side_missing_fonts(self):
        self.assertEqual(reader_mod._build_font_remap({}, {"fonts": {"major": "X"}}), {})
        self.assertEqual(reader_mod._build_font_remap({"fonts": {"major": "X"}}, {}), {})
        self.assertEqual(reader_mod._build_font_remap(None, None), {})

    def test_apply_remap_rewrites_explicit_latin(self):
        from lxml import etree
        el = etree.fromstring(self._wrap('<a:latin typeface="Helvetica Neue"/>'))
        n = reader_mod._apply_font_remap(el, {"helvetica neue": "Inter"})
        self.assertEqual(n, 1)
        self.assertEqual(el[0].get("typeface"), "Inter")

    def test_apply_remap_preserves_one_off_explicit_font(self):
        """A non-theme explicit font (Courier code snippet, etc.) survives unchanged."""
        from lxml import etree
        el = etree.fromstring(self._wrap('<a:latin typeface="Courier New"/>'))
        n = reader_mod._apply_font_remap(el, {"helvetica neue": "Inter"})
        self.assertEqual(n, 0)
        self.assertEqual(el[0].get("typeface"), "Courier New")

    def test_apply_remap_skips_theme_refs(self):
        """+mj-lt and +mn-lt already self-resolve; never rewrite them."""
        from lxml import etree
        el = etree.fromstring(self._wrap(
            '<a:latin typeface="+mj-lt"/>',
            '<a:latin typeface="+mn-lt"/>',
        ))
        n = reader_mod._apply_font_remap(el, {"+mj-lt": "Inter", "helvetica": "Inter"})
        self.assertEqual(n, 0)
        self.assertEqual(el[0].get("typeface"), "+mj-lt")
        self.assertEqual(el[1].get("typeface"), "+mn-lt")

    def test_apply_remap_handles_ea_and_cs(self):
        from lxml import etree
        el = etree.fromstring(self._wrap(
            '<a:latin typeface="Helvetica"/>',
            '<a:ea typeface="Helvetica"/>',
            '<a:cs typeface="Helvetica"/>',
        ))
        n = reader_mod._apply_font_remap(el, {"helvetica": "Aptos"})
        self.assertEqual(n, 3)
        for node in el:
            self.assertEqual(node.get("typeface"), "Aptos")

    def test_apply_remap_case_insensitive_match(self):
        from lxml import etree
        el = etree.fromstring(self._wrap('<a:latin typeface="HELVETICA"/>'))
        n = reader_mod._apply_font_remap(el, {"helvetica": "Aptos"})
        self.assertEqual(n, 1)
        self.assertEqual(el[0].get("typeface"), "Aptos")

    def test_apply_remap_returns_zero_on_empty_remap(self):
        from lxml import etree
        el = etree.fromstring(self._wrap('<a:latin typeface="Helvetica"/>'))
        self.assertEqual(reader_mod._apply_font_remap(el, {}), 0)
        self.assertEqual(el[0].get("typeface"), "Helvetica")


# ---------------------------------------------------------------------------
# v4.1 — extract_structured_atoms / extract_picture_assets group recursion
# ---------------------------------------------------------------------------


class TestGroupRecursion(unittest.TestCase):
    def _add_callout(self, container, left_in: int, top_in: int):
        return container.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left_in), Inches(top_in), Inches(2), Inches(1),
        )

    def test_iter_descends_into_group(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        grp = slide.shapes.add_group_shape()
        self._add_callout(grp, 1, 1)
        self._add_callout(grp, 4, 1)
        leafs = list(cli_mod._iter_shapes_recursive(slide))
        # Title placeholder + 2 grouped callouts; no GROUP itself.
        kinds = [str(getattr(s, "shape_type", "")) for s in leafs]
        self.assertEqual(sum("AUTO_SHAPE" in k for k in kinds), 2)
        self.assertFalse(any("GROUP" in k for k in kinds))

    def test_iter_descends_into_nested_groups(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        outer = slide.shapes.add_group_shape()
        inner = outer.shapes.add_group_shape()
        self._add_callout(inner, 1, 1)
        leafs = list(cli_mod._iter_shapes_recursive(slide))
        autos = [s for s in leafs if "AUTO_SHAPE" in str(getattr(s, "shape_type", ""))]
        self.assertEqual(len(autos), 1)

    def test_iter_depth_limit_skips_too_deep_shapes(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        # Build _MAX_GROUP_DEPTH+2 nested groups, shape at the bottom.
        container = slide.shapes
        for _ in range(cli_mod._MAX_GROUP_DEPTH + 2):
            container = container.add_group_shape().shapes
        container.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1), Inches(1), Inches(2), Inches(1),
        )
        leafs = list(cli_mod._iter_shapes_recursive(slide))
        autos = [s for s in leafs if "AUTO_SHAPE" in str(getattr(s, "shape_type", ""))]
        self.assertEqual(autos, [], "shape below depth limit should be skipped")

    def test_extract_structured_atoms_descends_into_groups(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        grp = slide.shapes.add_group_shape()
        self._add_callout(grp, 1, 1)
        self._add_callout(grp, 4, 1)
        with tempfile.TemporaryDirectory() as d:
            entries = cli_mod.extract_structured_atoms(
                slide, deck_stem="syn_group", slide_number=1,
                assets_dir=Path(d),
                slide_w=prs.slide_width, slide_h=prs.slide_height,
            )
            # Two grouped callouts should both surface as callout atoms.
            self.assertEqual(len(entries), 2)
            for entry in entries:
                self.assertIn("atom", entry)
                self.assertEqual(entry["kind"], "callout")
                # Geometry is per-shape fractional; group-local coords
                # still produce valid fractions (sometimes inflated for
                # scaled groups; here the groups have default transforms).
                for key in ("x", "y", "w", "h", "region"):
                    self.assertIn(key, entry)
                sha = entry["atom"].replace("asset_", "")
                xmls = list(Path(d).glob(f"{sha}*.xml"))
                yamls = list(Path(d).glob(f"{sha}*.yaml"))
                self.assertEqual(len(xmls), 1, f"missing xml for {entry['atom']}")
                self.assertEqual(len(yamls), 1, f"missing yaml for {entry['atom']}")

    def test_extract_structured_atoms_mixes_grouped_and_flat(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        # Flat callout.
        self._add_callout(slide, 1, 1)
        # Grouped callout pair.
        grp = slide.shapes.add_group_shape()
        self._add_callout(grp, 4, 1)
        self._add_callout(grp, 7, 1)
        with tempfile.TemporaryDirectory() as d:
            entries = cli_mod.extract_structured_atoms(
                slide, deck_stem="syn_mixed", slide_number=1,
                assets_dir=Path(d),
                slide_w=prs.slide_width, slide_h=prs.slide_height,
            )
            self.assertEqual(len(entries), 3)
            self.assertTrue(all(e["kind"] == "callout" for e in entries))


# ---------------------------------------------------------------------------
# v4.2 — slide anatomy: enriched layout + positional inventory
# ---------------------------------------------------------------------------


class TestSlideAnatomy(unittest.TestCase):
    def test_shape_kind_label_identifies_common_kinds(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        callout = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1), Inches(1), Inches(2), Inches(1),
        )
        tbl = slide.shapes.add_table(
            rows=2, cols=2,
            left=Inches(4), top=Inches(1), width=Inches(3), height=Inches(2),
        )
        self.assertEqual(cli_mod._shape_kind_label(callout), "callout")
        self.assertEqual(cli_mod._shape_kind_label(tbl), "table")

    def test_shape_geometry_yields_fractions_and_region(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        shp = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(7), Inches(5), Inches(2), Inches(1.5),  # right-bottom area
        )
        geom = cli_mod._shape_geometry(shp, prs.slide_width, prs.slide_height)
        # Fractions in (0, 1) for a real shape on a real slide.
        for k in ("x", "y", "w", "h"):
            self.assertIsInstance(geom[k], float)
            self.assertGreaterEqual(geom[k], 0.0)
            self.assertLess(geom[k], 1.0)
        # Center of the shape lands in the right-bottom third → region
        # should be "bottom-right".
        self.assertEqual(geom["region"], "bottom-right")

    def test_infer_layout_includes_non_placeholder_atoms(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        # Title placeholder is already present from layout 5.
        slide.shapes.title.text = "hi"
        # Add a callout in the bottom-right.
        slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(7), Inches(5), Inches(2), Inches(1.5),
        )
        slots, renames = cli_mod.detect_slots(
            slide, prs.slide_width, prs.slide_height, {},
        )
        layout = cli_mod.infer_layout(
            slide, slots, renames, prs.slide_width, prs.slide_height,
        )
        # The title slot should appear with its slot id; the callout
        # should appear as "callout@bottom-right".
        self.assertIn("@", layout)
        self.assertIn("callout@bottom-right", layout)

    def test_infer_layout_disambiguates_repeated_kinds(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        # Two callouts at different positions → expect "callout@…, callout#2@…".
        slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1), Inches(1), Inches(2), Inches(2),
        )
        slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(6), Inches(5), Inches(2), Inches(1.5),
        )
        slots, renames = cli_mod.detect_slots(
            slide, prs.slide_width, prs.slide_height, {},
        )
        layout = cli_mod.infer_layout(
            slide, slots, renames, prs.slide_width, prs.slide_height,
        )
        self.assertIn("callout@", layout)
        self.assertIn("callout#2@", layout)

    def test_infer_layout_skips_decorative_tiny_shapes(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        # Hairline-sized shape (well under 0.5% area) — should be filtered.
        slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1), Inches(1), Inches(0.05), Inches(0.05),
        )
        slots, renames = cli_mod.detect_slots(
            slide, prs.slide_width, prs.slide_height, {},
        )
        layout = cli_mod.infer_layout(
            slide, slots, renames, prs.slide_width, prs.slide_height,
        )
        # No "callout@" in the layout — the tiny shape was skipped.
        self.assertNotIn("callout@", layout)

    def test_extract_structured_atoms_inventory_carries_geometry(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(7), Inches(5), Inches(2), Inches(1.5),
        )
        with tempfile.TemporaryDirectory() as d:
            entries = cli_mod.extract_structured_atoms(
                slide, deck_stem="syn_geom", slide_number=1,
                assets_dir=Path(d),
                slide_w=prs.slide_width, slide_h=prs.slide_height,
            )
            self.assertEqual(len(entries), 1)
            e = entries[0]
            self.assertTrue(e["atom"].startswith("asset_"))
            self.assertEqual(e["kind"], "callout")
            self.assertEqual(e["region"], "bottom-right")
            self.assertGreater(e["x"], 0.5)  # right-half
            self.assertGreater(e["y"], 0.5)  # bottom-half


# ---------------------------------------------------------------------------
# v4.2 — composition required for picture-kinds only
# ---------------------------------------------------------------------------


class TestAssetValidation(unittest.TestCase):
    BASE = {
        "subject": "x",
        "feel": "formal",
        "colors": ["gray"],
        "scope": ["generic"],
        "suitable_for": ["data"],
    }

    def test_table_atom_with_empty_composition_validates(self):
        data = {**self.BASE, "kind": "table", "composition": ""}
        self.assertEqual(cli_mod.validate_asset(data), [])

    def test_callout_atom_with_empty_composition_validates(self):
        data = {**self.BASE, "kind": "callout", "composition": ""}
        self.assertEqual(cli_mod.validate_asset(data), [])

    def test_photo_with_empty_composition_fails(self):
        data = {**self.BASE, "kind": "photo", "composition": ""}
        errs = cli_mod.validate_asset(data)
        self.assertTrue(
            any("composition is required" in e for e in errs),
            errs,
        )

    def test_photo_with_valid_composition_validates(self):
        data = {**self.BASE, "kind": "photo", "composition": "centered"}
        self.assertEqual(cli_mod.validate_asset(data), [])

    def test_bad_composition_value_fails_regardless_of_kind(self):
        # Even for a structured atom, a NON-EMPTY but invalid composition
        # still fails — we don't accept arbitrary strings.
        data = {**self.BASE, "kind": "table", "composition": "lopsided"}
        errs = cli_mod.validate_asset(data)
        self.assertTrue(any("not in" in e for e in errs), errs)


if __name__ == "__main__":
    unittest.main()
