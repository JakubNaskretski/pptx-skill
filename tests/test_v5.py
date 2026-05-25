"""Tests for pptx-skill v5 — structural-skeleton pipeline.

Run from repo root: ``python3 -m unittest tests.test_v5``

unittest-only (no pytest dep). Synthetic decks built in-memory via
python-pptx; no binary fixtures committed.

Covers:
- ingest_v5.digest_skeleton — slot inventory + unmapped_shapes + shape_id
- ingest_v5.compute_repeated_picture_info — median area + threshold
- ingest_v5._propose_categories — opening / data / closing rules
- reader.cmd_v5_match_skeletons — ranking + zero-match suggested_action
- reader.cmd_v5_validate_plan — required_unfilled error + overflow:shrink warning
- reader.cmd_v5_compose — end-to-end build + shape placement
- Re-ingest preservation of user-set status / categories / overlap_decision
"""
from __future__ import annotations

import io
import json
import sys
import tempfile
import unittest
from pathlib import Path
from unittest import mock

REPO = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(REPO / "authoring"))
sys.path.insert(0, str(REPO / "consumer"))

from pptx import Presentation  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.util import Inches  # noqa: E402

import ingest_v5  # noqa: E402
import reader as reader_mod  # noqa: E402


# ---------------------------------------------------------------------------
# Synthetic deck helpers
# ---------------------------------------------------------------------------


def make_title_body_deck(n_slides: int = 1):
    """Build an in-memory deck with N title+body slides."""
    prs = Presentation()
    for i in range(n_slides):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Title {i+1}"
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                ph.text = f"Body line A {i+1}\nBody line B {i+1}"
    return prs


def stub_theme_v4():
    """Minimal v4 theme dict matching what extract_deck_theme produces."""
    return {
        "deck": "synthetic",
        "palette": {"dk1": "#000000", "lt1": "#FFFFFF", "accent1": "#FF0000"},
        "aliases": {"primary": "accent1", "accent": "accent1",
                    "text": "dk1", "background": "lt1"},
        "fonts": {"major": "Calibri", "minor": "Calibri"},
    }


def add_picture_to_slide(slide, png_bytes: bytes, left, top, width, height):
    """Add a Picture shape from raw PNG bytes."""
    return slide.shapes.add_picture(io.BytesIO(png_bytes), left, top, width, height)


# Tiny valid PNG (1x1 black pixel) — minimal binary for test pictures.
_TINY_PNG = bytes.fromhex(
    "89504e470d0a1a0a0000000d49484452"
    "00000001000000010802000000907753"
    "de0000000c4944415478da6300010000"
    "0500010d0a2db40000000049454e44ae"
    "426082"
)


# ---------------------------------------------------------------------------
# digest_skeleton — slot inventory
# ---------------------------------------------------------------------------


class TestDigestSkeleton(unittest.TestCase):
    def test_writes_skeleton_yaml_with_slots(self):
        prs = make_title_body_deck(1)
        with tempfile.TemporaryDirectory() as tmp:
            skel_root = Path(tmp) / "skeletons"
            out = ingest_v5.digest_skeleton(
                prs.slides[0], prs.slide_width, prs.slide_height,
                "test", 1, stub_theme_v4(), skel_root,
            )
        self.assertEqual(out["id"], "test_01")
        self.assertEqual(out["status"], "pending")
        self.assertTrue(any(s["kind"] == "heading" for s in out["slots"]))
        self.assertTrue(any(s["kind"] == "bullets" for s in out["slots"]))

    def test_every_slot_has_shape_id(self):
        prs = make_title_body_deck(1)
        with tempfile.TemporaryDirectory() as tmp:
            out = ingest_v5.digest_skeleton(
                prs.slides[0], prs.slide_width, prs.slide_height,
                "test", 1, stub_theme_v4(), Path(tmp),
            )
        for slot in out["slots"]:
            self.assertIn("shape_id", slot,
                          f"slot {slot.get('id')} missing shape_id (re-ingest preservation breaks without it)")

    def test_unmapped_shapes_captures_skipped(self):
        prs = make_title_body_deck(1)
        # Add a freeform (auto-shape) that the heuristic drops.
        slide = prs.slides[0]
        slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.1), Inches(0.1), Inches(0.3), Inches(0.3),
        )
        with tempfile.TemporaryDirectory() as tmp:
            out = ingest_v5.digest_skeleton(
                slide, prs.slide_width, prs.slide_height,
                "test", 1, stub_theme_v4(), Path(tmp),
            )
        self.assertGreater(len(out["unmapped_shapes"]), 0)
        kinds = {u["kind_hint"] for u in out["unmapped_shapes"]}
        self.assertIn("auto_shape", kinds)


# ---------------------------------------------------------------------------
# compute_repeated_picture_info — brand-mark detection
# ---------------------------------------------------------------------------


class TestRepeatedPictureDetection(unittest.TestCase):
    def test_picture_on_every_slide_is_flagged(self):
        prs = make_title_body_deck(3)
        # Add the same picture to all 3 slides at the same small size.
        for slide in prs.slides:
            add_picture_to_slide(
                slide, _TINY_PNG,
                Inches(0.1), Inches(0.1), Inches(0.5), Inches(0.5),
            )
        info = ingest_v5.compute_repeated_picture_info(prs)
        self.assertEqual(len(info), 1)
        median = list(info.values())[0]
        self.assertGreater(median, 0)

    def test_unique_pictures_not_flagged(self):
        prs = make_title_body_deck(3)
        # Different bytes per slide → different SHAs.
        for i, slide in enumerate(prs.slides):
            png = _TINY_PNG + bytes([i])  # vary by 1 byte
            add_picture_to_slide(
                slide, png,
                Inches(0.1), Inches(0.1), Inches(0.5), Inches(0.5),
            )
        info = ingest_v5.compute_repeated_picture_info(prs)
        self.assertEqual(info, {})


# ---------------------------------------------------------------------------
# Auto-classifier
# ---------------------------------------------------------------------------


class TestCategoryProposal(unittest.TestCase):
    def test_table_slot_triggers_data_category(self):
        slots = [
            {"kind": "heading", "geometry": {"h": 0.1}},
            {"kind": "table", "geometry": {"h": 0.5}},
        ]
        cats = ingest_v5._propose_categories(slots)
        self.assertIn("data", cats)

    def test_closing_pattern_in_text(self):
        slots = [
            {"kind": "heading", "geometry": {"h": 0.2},
             "source_excerpt": "Thank you for your attention"},
        ]
        cats = ingest_v5._propose_categories(slots)
        self.assertIn("closing", cats)

    def test_fallback_to_content(self):
        slots = [{"kind": "paragraph", "geometry": {"h": 0.5},
                  "source_excerpt": "Some paragraph text"}]
        cats = ingest_v5._propose_categories(slots)
        self.assertEqual(cats, ["content"])


# ---------------------------------------------------------------------------
# Re-ingest preservation
# ---------------------------------------------------------------------------


class TestReingestPreservation(unittest.TestCase):
    def test_user_set_status_survives(self):
        prs = make_title_body_deck(1)
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            ingest_v5.digest_skeleton(
                prs.slides[0], prs.slide_width, prs.slide_height,
                "deckA", 1, stub_theme_v4(), root,
            )
            # User marks done + sets overlap_decision via the UI
            sk_path = root / "deckA_01" / "skeleton.yaml"
            import yaml
            data = yaml.safe_load(sk_path.read_text())
            data["status"] = "done"
            data["overlap_decision"] = "image_slot"
            data["categories"] = ["data"]
            sk_path.write_text(yaml.safe_dump(data, sort_keys=False))
            # Re-ingest
            out = ingest_v5.digest_skeleton(
                prs.slides[0], prs.slide_width, prs.slide_height,
                "deckA", 1, stub_theme_v4(), root,
            )
        self.assertEqual(out["status"], "done")
        self.assertEqual(out["overlap_decision"], "image_slot")
        self.assertEqual(out["categories"], ["data"])

    def test_user_edited_slot_kind_survives_reingest(self):
        """User clicks a kind button to reclassify a slot; the new
        kind must NOT be reverted on re-ingest. Mirrors what the
        C-actions reclassify endpoint does, then re-runs digest.
        """
        prs = make_title_body_deck(1)
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            out = ingest_v5.digest_skeleton(
                prs.slides[0], prs.slide_width, prs.slide_height,
                "deckA", 1, stub_theme_v4(), root,
            )
            # Find a slot the heuristic put as one kind; flip it to another
            sk_path = root / "deckA_01" / "skeleton.yaml"
            import yaml
            data = yaml.safe_load(sk_path.read_text())
            target = data["slots"][0]
            original_kind = target["kind"]
            new_kind = "paragraph" if original_kind == "heading" else "heading"
            target["kind"] = new_kind
            target["user_edited"] = True
            sk_path.write_text(yaml.safe_dump(data, sort_keys=False))
            # Re-ingest
            out2 = ingest_v5.digest_skeleton(
                prs.slides[0], prs.slide_width, prs.slide_height,
                "deckA", 1, stub_theme_v4(), root,
            )
        slot = next(s for s in out2["slots"] if s["id"] == target["id"])
        self.assertEqual(slot["kind"], new_kind,
                         "user-edited kind must survive re-ingest")
        self.assertTrue(slot.get("user_edited"))


# ---------------------------------------------------------------------------
# Constraint helpers
# ---------------------------------------------------------------------------


class TestConstraintHelpers(unittest.TestCase):
    def test_text_fit_under(self):
        fits, _, hr = reader_mod._v5_check_text_fit("hello", {"max_chars": 10})
        self.assertTrue(fits)
        self.assertEqual(hr, 5)

    def test_text_fit_over(self):
        fits, reason, hr = reader_mod._v5_check_text_fit("hello world", {"max_chars": 5})
        self.assertFalse(fits)
        self.assertLess(hr, 0)
        self.assertIn("max_chars", reason)

    def test_bullets_too_many_items(self):
        fits, reason, _ = reader_mod._v5_check_bullets_fit(["a", "b", "c"], {"max_items": 2})
        self.assertFalse(fits)
        self.assertIn("max_items", reason)


# ---------------------------------------------------------------------------
# Reader CLI commands — exercised via cmd_* functions directly to stdout
# ---------------------------------------------------------------------------


class _StubArgs:
    """Minimal stand-in for argparse.Namespace."""
    def __init__(self, **kw):
        self.__dict__.update(kw)


def _capture_stdout(fn, *args):
    """Run a cmd_v5_* function; capture and parse its stdout JSON."""
    buf = io.StringIO()
    real_stdout = sys.stdout
    sys.stdout = buf
    try:
        fn(*args)
    finally:
        sys.stdout = real_stdout
    return json.loads(buf.getvalue())


class TestMatchSkeletons(unittest.TestCase):
    def setUp(self):
        # Build a workspace and patch _v5_bundle_root.
        self.tmpdir = tempfile.TemporaryDirectory()
        root = Path(self.tmpdir.name)
        (root / "themes").mkdir()
        (root / "skeletons").mkdir()
        # Two skeletons: tight title (max 30) vs loose (max 200)
        for sid, max_chars in [("tight_01", 30), ("loose_01", 200)]:
            d = root / "skeletons" / sid
            d.mkdir()
            (d / "skeleton.yaml").write_text(
                f"id: {sid}\n"
                f"source_deck: synth\n"
                f"source_slide_index: 1\n"
                f"status: pending\n"
                f"categories: [content]\n"
                f"slots:\n"
                f"  - id: title\n"
                f"    kind: heading\n"
                f"    geometry: {{x: 0.05, y: 0.05, w: 0.9, h: 0.1}}\n"
                f"    constraints: {{max_chars: {max_chars}, required: true}}\n",
            )
        self._patcher = mock.patch.object(reader_mod, "_v5_bundle_root",
                                          return_value=root)
        self._patcher.start()

    def tearDown(self):
        self._patcher.stop()
        self.tmpdir.cleanup()

    def test_tighter_skeleton_ranks_higher(self):
        out = _capture_stdout(
            reader_mod.cmd_v5_match_skeletons,
            _StubArgs(content='{"title": "Q4 results beat consensus"}',
                      category=None, has_slot=None),
        )
        self.assertEqual(len(out["matches"]), 2)
        # The 30-char-max skeleton fits a 27-char title more tightly
        # than the 200-char-max one → ranks first.
        self.assertEqual(out["matches"][0]["skeleton_id"], "tight_01")

    def test_zero_match_returns_suggested_action(self):
        long_title = "X" * 250  # exceeds both 30 and 200 char limits
        out = _capture_stdout(
            reader_mod.cmd_v5_match_skeletons,
            _StubArgs(content=json.dumps({"title": long_title}),
                      category=None, has_slot=None),
        )
        self.assertEqual(out["matches"], [])
        self.assertGreater(len(out["issues"]), 0)
        issue = out["issues"][0]
        self.assertIn("suggested_action", issue)
        # Tightest constraint across both candidates is 30
        self.assertEqual(issue["tightest_constraint"], 30)


class TestValidatePlan(unittest.TestCase):
    def setUp(self):
        self.tmpdir = tempfile.TemporaryDirectory()
        root = Path(self.tmpdir.name)
        (root / "themes").mkdir()
        (root / "skeletons" / "sk_01").mkdir(parents=True)
        (root / "skeletons" / "sk_01" / "skeleton.yaml").write_text(
            "id: sk_01\n"
            "source_deck: synth\n"
            "source_slide_index: 1\n"
            "status: pending\n"
            "categories: [content]\n"
            "slots:\n"
            "  - id: title\n"
            "    kind: heading\n"
            "    geometry: {x: 0, y: 0, w: 1, h: 0.2}\n"
            "    constraints: {max_chars: 20, required: true}\n"
            "  - id: body\n"
            "    kind: paragraph\n"
            "    geometry: {x: 0, y: 0.3, w: 1, h: 0.5}\n"
            "    constraints: {max_chars: 100, required: false}\n",
        )
        self._patcher = mock.patch.object(reader_mod, "_v5_bundle_root",
                                          return_value=root)
        self._patcher.start()
        self.root = root

    def tearDown(self):
        self._patcher.stop()
        self.tmpdir.cleanup()

    def test_required_unfilled_is_error(self):
        plan = self.root / "plan.json"
        plan.write_text(json.dumps([{"skeleton_id": "sk_01", "slots": {"body": "ok"}}]))
        out = _capture_stdout(
            reader_mod.cmd_v5_validate_plan,
            _StubArgs(plan=str(plan)),
        )
        self.assertFalse(out["ok"])
        self.assertTrue(any(e["violation"] == "required_unfilled" for e in out["errors"]))

    def test_overflow_shrink_is_warning_not_error(self):
        plan = self.root / "plan.json"
        plan.write_text(json.dumps([{
            "skeleton_id": "sk_01",
            "slots": {
                "title": {"value": "Way too long for the 20-char title slot",
                          "overflow": "shrink"},
            },
        }]))
        out = _capture_stdout(
            reader_mod.cmd_v5_validate_plan,
            _StubArgs(plan=str(plan)),
        )
        self.assertGreater(len(out["warnings"]), 0)
        warning_slots = [w["slot_id"] for w in out["warnings"]]
        self.assertIn("title", warning_slots)
        # Title isn't a hard error because overflow:shrink covers it
        error_slots = [e.get("slot_id") for e in out["errors"]]
        self.assertNotIn("title", error_slots)


# ---------------------------------------------------------------------------
# v5.1 — slot role inference + role-aware match-skeletons
# ---------------------------------------------------------------------------


class TestSlotRoleInference(unittest.TestCase):
    def test_title_placeholder_gets_page_title_role(self):
        prs = make_title_body_deck(1)
        with tempfile.TemporaryDirectory() as tmp:
            out = ingest_v5.digest_skeleton(
                prs.slides[0], prs.slide_width, prs.slide_height,
                "deck", 1, stub_theme_v4(), Path(tmp),
            )
        title = next(s for s in out["slots"] if s["kind"] == "heading")
        self.assertEqual(title.get("role"), "page_title")

    def test_body_placeholder_with_short_bullets_is_key_points(self):
        prs = make_title_body_deck(1)
        # Default body has 2 short bullets → key_points (≤4 items, ≤60 chars/item)
        with tempfile.TemporaryDirectory() as tmp:
            out = ingest_v5.digest_skeleton(
                prs.slides[0], prs.slide_width, prs.slide_height,
                "deck", 1, stub_theme_v4(), Path(tmp),
            )
        body = next((s for s in out["slots"] if s["kind"] == "bullets"), None)
        self.assertIsNotNone(body)
        self.assertEqual(body.get("role"), "key_points")

    def test_revert_flag_disables_inference(self):
        prs = make_title_body_deck(1)
        with mock.patch.object(ingest_v5, "_ENABLE_SLOT_ROLES", False):
            with tempfile.TemporaryDirectory() as tmp:
                out = ingest_v5.digest_skeleton(
                    prs.slides[0], prs.slide_width, prs.slide_height,
                    "deck", 1, stub_theme_v4(), Path(tmp),
                )
        for slot in out["slots"]:
            self.assertNotIn("role", slot,
                             "role must not be emitted when _ENABLE_SLOT_ROLES=False")


class _StubFreeShape:
    """Minimal stub for _infer_slot_role free-shape path — only needs
    is_placeholder=False to skip the placeholder branch."""
    is_placeholder = False


class TestFreeShapeRoleInference(unittest.TestCase):
    """Direct unit tests for the tightened free-shape rules in
    _infer_slot_role + the cross-slot refinement in _refine_slot_roles.
    """

    def _make_slot(self, kind, *, x=0.1, y=0.1, w=0.8, h=0.1,
                   size_pt=None, max_chars=None, max_lines=None,
                   max_items=None, max_chars_per_item=None,
                   excerpt=""):
        slot = {"kind": kind,
                "geometry": {"x": x, "y": y, "w": w, "h": h},
                "style": {},
                "constraints": {},
                "source_excerpt": excerpt}
        if size_pt is not None:
            slot["style"]["size_pt"] = size_pt
        if max_chars is not None:
            slot["constraints"]["max_chars"] = max_chars
        if max_lines is not None:
            slot["constraints"]["max_lines"] = max_lines
        if max_items is not None:
            slot["constraints"]["max_items"] = max_items
        if max_chars_per_item is not None:
            slot["constraints"]["max_chars_per_item"] = max_chars_per_item
        return slot

    def _infer(self, slot):
        return ingest_v5._infer_slot_role(_StubFreeShape(), slot,
                                          slide_w=9144000, slide_h=6858000)

    def test_top_large_heading_is_page_title(self):
        slot = self._make_slot("heading", y=0.05, size_pt=36)
        self.assertEqual(self._infer(slot), "page_title")

    def test_top_medium_heading_is_section_header(self):
        slot = self._make_slot("heading", y=0.10, size_pt=20)
        self.assertEqual(self._infer(slot), "section_header")

    def test_huge_short_text_is_kpi_value(self):
        slot = self._make_slot("paragraph", y=0.40, size_pt=48, max_chars=10)
        self.assertEqual(self._infer(slot), "kpi_value")

    def test_small_bottom_paragraph_is_footnote(self):
        slot = self._make_slot("paragraph", y=0.92, size_pt=9, max_chars=120)
        self.assertEqual(self._infer(slot), "footnote")

    def test_cta_action_verb_prefix(self):
        slot = self._make_slot("paragraph", y=0.6, size_pt=18,
                               max_chars=40, excerpt="Visit our website today")
        self.assertEqual(self._infer(slot), "cta")

    def test_caption_requires_below_top(self):
        # Same small short text — at top is NOT a caption (likely byline-ish);
        # below the middle IS a caption.
        top = self._make_slot("paragraph", y=0.05, size_pt=11, max_chars=50)
        mid = self._make_slot("paragraph", y=0.55, size_pt=11, max_chars=50)
        self.assertNotEqual(self._infer(top), "caption")
        self.assertEqual(self._infer(mid), "caption")

    def test_long_paragraph_is_body(self):
        slot = self._make_slot("paragraph", y=0.4, size_pt=18, max_chars=300)
        self.assertEqual(self._infer(slot), "body")

    def test_short_uncertain_paragraph_falls_through_to_none(self):
        # Medium-sized, short, mid-slide, no CTA verbs — not enough
        # signal for any single role.
        slot = self._make_slot("paragraph", y=0.4, size_pt=16, max_chars=40,
                               excerpt="Some neutral text")
        self.assertIsNone(self._infer(slot))

    def test_refinement_assigns_kpi_label_near_kpi_value(self):
        slots = [
            {"kind": "paragraph", "role": "kpi_value",
             "geometry": {"x": 0.1, "y": 0.4, "w": 0.2, "h": 0.15},
             "style": {"size_pt": 48}, "constraints": {"max_chars": 5}},
            # Small text directly below the KPI — should get kpi_label
            {"kind": "paragraph",
             "geometry": {"x": 0.1, "y": 0.56, "w": 0.2, "h": 0.05},
             "style": {"size_pt": 12}, "constraints": {"max_chars": 30}},
        ]
        ingest_v5._refine_slot_roles(slots)
        self.assertEqual(slots[1].get("role"), "kpi_label")

    def test_refinement_assigns_byline_below_page_title(self):
        slots = [
            {"kind": "heading", "role": "page_title",
             "geometry": {"x": 0.05, "y": 0.05, "w": 0.9, "h": 0.15},
             "style": {"size_pt": 40}, "constraints": {"max_chars": 50}},
            # Small short text directly under the title — byline pattern
            {"kind": "paragraph",
             "geometry": {"x": 0.05, "y": 0.22, "w": 0.5, "h": 0.04},
             "style": {"size_pt": 12}, "constraints": {"max_chars": 40}},
        ]
        ingest_v5._refine_slot_roles(slots)
        self.assertEqual(slots[1].get("role"), "byline")

    def test_refinement_does_not_override_existing_role(self):
        slots = [
            {"kind": "paragraph", "role": "kpi_value",
             "geometry": {"x": 0.1, "y": 0.4, "w": 0.2, "h": 0.15},
             "style": {"size_pt": 48}, "constraints": {"max_chars": 5}},
            # Pre-assigned as 'caption' — refinement must NOT downgrade
            # it to kpi_label even though it sits adjacent to a kpi_value.
            {"kind": "paragraph", "role": "caption",
             "geometry": {"x": 0.1, "y": 0.56, "w": 0.2, "h": 0.05},
             "style": {"size_pt": 12}, "constraints": {"max_chars": 30}},
        ]
        ingest_v5._refine_slot_roles(slots)
        self.assertEqual(slots[1]["role"], "caption")


class TestRoleAwareMatching(unittest.TestCase):
    """match-skeletons prefers role match over first-of-kind when both
    fire. So if a content key names a role and the skeleton has a slot
    with that role, that slot is picked even if an earlier slot of the
    same kind exists.
    """

    def setUp(self):
        self.tmpdir = tempfile.TemporaryDirectory()
        root = Path(self.tmpdir.name)
        (root / "themes").mkdir()
        (root / "skeletons" / "sk_01").mkdir(parents=True)
        # Skeleton with two heading slots: first is "section_header",
        # second is "page_title". Without role-matching, "title" content
        # would pick the first (section_header). With role-matching,
        # it picks the second.
        (root / "skeletons" / "sk_01" / "skeleton.yaml").write_text(
            "id: sk_01\n"
            "source_deck: synth\n"
            "source_slide_index: 1\n"
            "status: pending\n"
            "categories: [content]\n"
            "slots:\n"
            "  - id: header_a\n"
            "    kind: heading\n"
            "    role: section_header\n"
            "    geometry: {x: 0.05, y: 0.05, w: 0.9, h: 0.1}\n"
            "    constraints: {max_chars: 30, required: false}\n"
            "  - id: title_b\n"
            "    kind: heading\n"
            "    role: page_title\n"
            "    geometry: {x: 0.05, y: 0.2, w: 0.9, h: 0.15}\n"
            "    constraints: {max_chars: 60, required: true}\n",
        )
        self._patcher = mock.patch.object(reader_mod, "_v5_bundle_root",
                                          return_value=root)
        self._patcher.start()

    def tearDown(self):
        self._patcher.stop()
        self.tmpdir.cleanup()

    def test_title_content_key_picks_page_title_role(self):
        out = _capture_stdout(
            reader_mod.cmd_v5_match_skeletons,
            _StubArgs(content='{"title": "Q4 results"}',
                      category=None, has_slot=None),
        )
        self.assertEqual(len(out["matches"]), 1)
        # "title" content key → page_title role → title_b slot,
        # not header_a (which is the first heading-kind slot)
        self.assertEqual(out["matches"][0]["slot_mapping"]["title"], "title_b")


# ---------------------------------------------------------------------------
# Aspect-aware image crop (compose-v5 _v5_place_image)
# ---------------------------------------------------------------------------


def _wide_png_bytes() -> bytes:
    """Generate a 200x100 (2:1 wide) PNG."""
    from PIL import Image
    img = Image.new("RGB", (200, 100), color="red")
    buf = io.BytesIO()
    img.save(buf, "PNG")
    return buf.getvalue()


def _tall_png_bytes() -> bytes:
    """Generate a 100x200 (1:2 tall) PNG."""
    from PIL import Image
    img = Image.new("RGB", (100, 200), color="blue")
    buf = io.BytesIO()
    img.save(buf, "PNG")
    return buf.getvalue()


class TestImageAutoFit(unittest.TestCase):
    """The aspect-aware crop in _v5_place_image. Slot is square; asset
    is non-square. cover should center-crop preserving aspect;
    contain should letterbox; stretch should distort (legacy).
    """

    def _build_slot(self, auto_fit: str):
        return {
            "id": "hero", "kind": "image",
            "geometry": {"x": 0.1, "y": 0.1, "w": 0.2, "h": 0.2},  # square slot
            "constraints": {"auto_fit": auto_fit, "required": True},
        }

    def _place(self, png_bytes: bytes, slot: dict):
        from pptx import Presentation
        with tempfile.TemporaryDirectory() as tmp:
            d = Path(tmp)
            asset_path = d / "asset_test.png"
            asset_path.write_bytes(png_bytes)
            (d / "assets").mkdir()
            (d / "assets" / "asset_test.png").write_bytes(png_bytes)

            with mock.patch.object(reader_mod, "_v5_bundle_root", return_value=d):
                prs = Presentation()
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                ws = reader_mod._v5_place_image(
                    slide, slot, "asset_test", prs.slide_width, prs.slide_height, {},
                )
            pics = [s for s in slide.shapes if s.shape_type == 13]  # PICTURE
            return pics, ws

    def test_cover_wide_asset_in_square_slot_crops_sides(self):
        pics, _ = self._place(_wide_png_bytes(), self._build_slot("cover"))
        self.assertEqual(len(pics), 1)
        pic = pics[0]
        # 2:1 asset in 1:1 slot, cover policy → height-matched, sides cropped.
        # crop_left + crop_right should be > 0 (we lopped off width).
        self.assertGreater(pic.crop_left, 0)
        self.assertGreater(pic.crop_right, 0)
        # Top/bottom shouldn't be touched.
        self.assertEqual(pic.crop_top, 0)
        self.assertEqual(pic.crop_bottom, 0)

    def test_cover_tall_asset_in_square_slot_crops_top_bottom(self):
        pics, _ = self._place(_tall_png_bytes(), self._build_slot("cover"))
        pic = pics[0]
        self.assertGreater(pic.crop_top, 0)
        self.assertGreater(pic.crop_bottom, 0)
        self.assertEqual(pic.crop_left, 0)
        self.assertEqual(pic.crop_right, 0)

    def test_contain_wide_asset_letterboxes(self):
        pics, _ = self._place(_wide_png_bytes(), self._build_slot("contain"))
        pic = pics[0]
        # Letterbox: no cropping, image just shrunk to fit
        self.assertEqual(pic.crop_left, 0)
        self.assertEqual(pic.crop_right, 0)
        self.assertEqual(pic.crop_top, 0)
        # Placed shape height should be smaller than the slot height
        # because asset aspect is wider.
        # Slot is 0.2 * slide_height tall; placed should be less.
        from pptx.util import Emu
        # No precise numeric check — just that it's smaller than the slot.
        slot_h_emu = int(0.2 * 6858000)
        self.assertLess(pic.height, slot_h_emu)

    def test_stretch_uses_slot_dims_directly(self):
        pics, _ = self._place(_wide_png_bytes(), self._build_slot("stretch"))
        pic = pics[0]
        self.assertEqual(pic.crop_left, 0)
        self.assertEqual(pic.crop_right, 0)
        # Stretched: placed dimensions match the slot exactly.
        slot_w_emu = int(0.2 * 9144000)
        slot_h_emu = int(0.2 * 6858000)
        self.assertEqual(pic.width, slot_w_emu)
        self.assertEqual(pic.height, slot_h_emu)


# ---------------------------------------------------------------------------
# Compose round-trip
# ---------------------------------------------------------------------------


class TestComposeRoundTrip(unittest.TestCase):
    def setUp(self):
        self.tmpdir = tempfile.TemporaryDirectory()
        self.root = Path(self.tmpdir.name)
        (self.root / "skeletons" / "sk_01").mkdir(parents=True)
        # Write a real master.pptx
        themes = self.root / "themes" / "synth"
        themes.mkdir(parents=True)
        synth_prs = make_title_body_deck(1)
        synth_prs.save(str(themes / "master.pptx"))
        (themes / "theme.yaml").write_text(
            "id: synth\n"
            "palette: {primary: '#FF0000', accent: '#00FF00'}\n"
            "fonts: {major: Calibri, minor: Calibri}\n"
            "master_pptx: master.pptx\n",
        )
        (self.root / "skeletons" / "sk_01" / "skeleton.yaml").write_text(
            "id: sk_01\n"
            "source_deck: synth\n"
            "source_slide_index: 1\n"
            "status: pending\n"
            "categories: [content]\n"
            "slots:\n"
            "  - id: title\n"
            "    kind: heading\n"
            "    geometry: {x: 0.05, y: 0.05, w: 0.9, h: 0.15}\n"
            "    style: {font_role: major, size_pt: 32}\n"
            "    constraints: {max_chars: 50, required: true}\n"
            "  - id: body\n"
            "    kind: bullets\n"
            "    geometry: {x: 0.05, y: 0.3, w: 0.9, h: 0.6}\n"
            "    constraints: {max_items: 5, max_chars_per_item: 80}\n",
        )
        self._patcher = mock.patch.object(reader_mod, "_v5_bundle_root",
                                          return_value=self.root)
        self._patcher.start()

    def tearDown(self):
        self._patcher.stop()
        self.tmpdir.cleanup()

    def test_compose_produces_real_pptx(self):
        plan = self.root / "plan.json"
        plan.write_text(json.dumps([{
            "skeleton_id": "sk_01",
            "slots": {"title": "Hello world", "body": ["item a", "item b"]},
        }]))
        out_pptx = self.root / "out.pptx"
        _capture_stdout(
            reader_mod.cmd_v5_compose,
            _StubArgs(plan=str(plan), out=str(out_pptx), theme="synth"),
        )
        self.assertTrue(out_pptx.exists())
        prs = Presentation(str(out_pptx))
        self.assertEqual(len(prs.slides), 1)
        # Two text boxes added (title + body bullets)
        text_boxes = [s for s in prs.slides[0].shapes
                      if s.has_text_frame and s.text_frame.text.strip()]
        all_text = "\n".join(s.text_frame.text for s in text_boxes)
        self.assertIn("Hello world", all_text)
        self.assertIn("item a", all_text)


if __name__ == "__main__":
    unittest.main()
