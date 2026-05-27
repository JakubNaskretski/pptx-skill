"""pptx-skill consumer reader.

Three commands an agent can call to use a built skill bundle:

  python reader.py list [--filter key=value,key=value]
  python reader.py get <id>
  python reader.py compose <plan.json> <output.pptx>

All commands write JSON to stdout (compose also writes the deck).
No state. No vision required. Read SKILL.md for the contract.
"""

from __future__ import annotations

import argparse
import copy
import json
import re
import shutil
import sys
import tempfile
from pathlib import Path
from typing import Any

import yaml
from pptx import Presentation


HERE = Path(__file__).resolve().parent


# ---------------------------------------------------------------------------
# Bundle layout helpers
# ---------------------------------------------------------------------------


def bundle_root() -> Path:
    return HERE


def load_index() -> dict:
    p = bundle_root() / "index.json"
    if not p.exists():
        raise SystemExit(f"index.json not found at {p}")
    return json.loads(p.read_text(encoding="utf-8"))


def template_dir(tid: str) -> Path:
    return bundle_root() / "templates" / tid


def asset_path(aid: str) -> Path | None:
    """Find the asset binary by id (asset_<sha8>) — any extension."""
    assets_dir = bundle_root() / "assets"
    if not assets_dir.exists():
        return None
    for cand in assets_dir.glob(f"{aid}.*"):
        if cand.suffix == ".yaml":
            continue
        return cand
    return None


def asset_meta_path(aid: str) -> Path:
    return bundle_root() / "assets" / f"{aid}.yaml"


# ---------------------------------------------------------------------------
# Filter
# ---------------------------------------------------------------------------


UNSET_SENTINEL = "none"


def parse_filter(s: str | None) -> dict[str, list[str]]:
    """Parse `k1=v1,k2=v2|v3` into {k1: [v1], k2: [v2, v3]}.

    Pipe (`|`) expresses OR within a key. The literal value `none`
    matches items where the field is missing or empty — useful to
    include un-tagged items (`feel=warm|none`) or to audit them
    explicitly (`feel=none`).
    """
    if not s:
        return {}
    out: dict[str, list[str]] = {}
    for chunk in s.split(","):
        chunk = chunk.strip()
        if not chunk:
            continue
        if "=" not in chunk:
            raise SystemExit(f"bad --filter token (need key=value): {chunk!r}")
        k, v = chunk.split("=", 1)
        out[k.strip()] = [x.strip() for x in v.split("|") if x.strip()]
    return out


def _field_is_unset(got: Any) -> bool:
    if got is None:
        return True
    if isinstance(got, (str, list, dict)) and not got:
        return True
    return False


def matches_filter(item: dict, flt: dict[str, list[str]]) -> bool:
    for k, wants in flt.items():
        unset_ok = UNSET_SENTINEL in wants
        explicit = [w for w in wants if w != UNSET_SENTINEL]
        if k not in item or _field_is_unset(item.get(k)):
            if unset_ok:
                continue
            return False
        if not explicit:
            # Filter was `k=none` only and the field IS set → exclude.
            return False
        got = item[k]
        if isinstance(got, list):
            got_strs = [str(x) for x in got]
            if not any(w in got_strs for w in explicit):
                return False
        else:
            if str(got) not in explicit:
                return False
    return True


# ---------------------------------------------------------------------------
# Compose
# ---------------------------------------------------------------------------


def _find_shape_by_name(slide, name: str):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


_BULLET_PREFIX_RE = re.compile(r"^\s*[•·◦▪▫●◆■□*\-–—]\s+")


def _strip_bullet_prefix(text: str) -> str:
    """Strip leading bullet glyphs from each line.

    Templates render bullets via layout formatting; if the caller also
    prepended a glyph, the slide ends up with two bullets per line.
    """
    if not text:
        return text
    return "\n".join(_BULLET_PREFIX_RE.sub("", ln) for ln in text.split("\n"))


def _copy_run_color(src_font, dst_font) -> None:
    """Copy a run's font.color from src to dst, handling RGB + theme.

    ColorFormat assignment switches the colour type on assignment, so
    we try RGB first (most common) then theme-color. Silent fallback
    if neither resolves — leaves dst at its inherited default.
    """
    try:
        ctype = src_font.color.type
    except (AttributeError, ValueError):
        return
    if ctype is None:
        return
    try:
        rgb = src_font.color.rgb
        if rgb is not None:
            dst_font.color.rgb = rgb
            return
    except (AttributeError, ValueError):
        pass
    try:
        tc = src_font.color.theme_color
        if tc is not None:
            dst_font.color.theme_color = tc
    except (AttributeError, ValueError):
        pass


def _copy_run_font(src_font, dst_font) -> None:
    """Copy size/bold/italic/name/color from a template run to a new run."""
    try:
        if src_font.size is not None:
            dst_font.size = src_font.size
        if src_font.bold is not None:
            dst_font.bold = src_font.bold
        if src_font.italic is not None:
            dst_font.italic = src_font.italic
        if src_font.name is not None:
            dst_font.name = src_font.name
    except (AttributeError, ValueError):
        pass
    _copy_run_color(src_font, dst_font)


def _fill_text_shape(shape, value: str) -> None:
    """Replace the shape's text frame content with `value`, keeping its first
    paragraph's font/style (incl. colour) as the template."""
    if not shape.has_text_frame:
        return
    value = _strip_bullet_prefix(str(value))
    tf = shape.text_frame
    # Capture the first run's formatting cues if available.
    first_para = tf.paragraphs[0] if tf.paragraphs else None
    template_run = first_para.runs[0] if (first_para and first_para.runs) else None

    tf.clear()
    p = tf.paragraphs[0]
    if template_run is not None:
        # text_frame.clear() leaves an empty first paragraph — reuse it.
        run = p.add_run()
        run.text = value
        _copy_run_font(template_run.font, run.font)
    else:
        p.text = value


def _fill_bullets_shape(shape, values: list[str]) -> None:
    """Replace bullets in a text-frame placeholder. Each value becomes one
    paragraph (newlines within a value become soft breaks)."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    first_para = tf.paragraphs[0] if tf.paragraphs else None
    template_run = first_para.runs[0] if (first_para and first_para.runs) else None
    template_level = first_para.level if first_para is not None else 0

    tf.clear()
    if not values:
        return
    first = True
    for value in values:
        lines = _strip_bullet_prefix(str(value)).split("\n")
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.level = template_level
        run = p.add_run()
        run.text = lines[0]
        if template_run is not None:
            _copy_run_font(template_run.font, run.font)
        # Sub-runs for additional lines within a bullet inherit the
        # same font + colour as the primary run. (FINDINGS A3.17 noted
        # that without this they fall back to paragraph defaults and
        # show visible drift mid-bullet. Sticking with "\n" in run.text
        # for the soft break — the proper <a:br/> rewrite is a Phase
        # D follow-up.)
        for extra in lines[1:]:
            sub = p.add_run()
            sub.text = "\n" + extra
            if template_run is not None:
                _copy_run_font(template_run.font, sub.font)


def _fill_table_shape(shape, cells: list[list]) -> list[str]:
    """Fill a table placeholder's cells from a list-of-lists.

    The agent supplies cells as `[[row0_col0, row0_col1, ...], ...]`.
    Cells beyond what the template table holds are dropped with a
    warning. Cells the agent doesn't provide are left untouched —
    that keeps decorative template rows (e.g. footer summary) intact
    when the agent only wants to overwrite a subset. Per-cell font /
    colour formatting is inherited from each cell's existing first
    run.

    Returns a list of non-fatal warning strings.
    """
    warnings: list[str] = []
    if not getattr(shape, "has_table", False):
        warnings.append(f"slot '{shape.name}': not a table shape; skipping table fill")
        return warnings
    table = shape.table
    rows = list(table.rows)
    nrows = len(rows)

    if len(cells) > nrows:
        warnings.append(
            f"slot '{shape.name}': got {len(cells)} table rows, "
            f"template has {nrows}; truncating extras"
        )

    for i, row_vals in enumerate(cells[:nrows]):
        row_cells = list(rows[i].cells)
        ncols = len(row_cells)
        if not isinstance(row_vals, list):
            warnings.append(
                f"slot '{shape.name}': row {i} is not a list (got {type(row_vals).__name__}); skipping"
            )
            continue
        if len(row_vals) > ncols:
            warnings.append(
                f"slot '{shape.name}': row {i} got {len(row_vals)} cols, "
                f"template has {ncols}; truncating extras"
            )
        for j, val in enumerate(row_vals[:ncols]):
            cell = row_cells[j]
            tf = cell.text_frame
            first_para = tf.paragraphs[0] if tf.paragraphs else None
            template_run = first_para.runs[0] if (first_para and first_para.runs) else None
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = _strip_bullet_prefix(str(val))
            if template_run is not None:
                _copy_run_font(template_run.font, run.font)

    return warnings


def _replace_image_shape_legacy(slide, shape, image_path: Path) -> None:
    """Remove the existing Picture shape and add a fresh one at the same
    geometry. Loses crop / border / shadow / rotation / transparency /
    effects / alt text. Used as a fallback when the in-place rewire
    can't find what it needs.

    Placeholder shapes inherit geometry from their layout/master and may
    return None for left/top/width/height. We resolve those via the
    placeholder chain when possible, and fall back to slide-fraction
    defaults so add_picture never sees None (which raises a cryptic
    "a real number is required" error inside python-pptx).
    """
    left, top, width, height = _resolve_shape_geometry(slide, shape)
    name = shape.name

    sp = shape._element
    sp.getparent().remove(sp)

    new_pic = slide.shapes.add_picture(
        str(image_path), left, top, width=width, height=height
    )
    new_pic.name = name


def _resolve_shape_geometry(slide, shape) -> tuple[int, int, int, int]:
    """Return (left, top, width, height) in EMU for any shape, including
    placeholders that inherit geometry. Walks placeholder -> layout
    placeholder -> master placeholder, then falls back to slide-fraction
    defaults so the caller never has to deal with None.
    """
    def chain(attr: str):
        v = getattr(shape, attr, None)
        if v is not None:
            return v
        try:
            ph_format = shape.placeholder_format
            if ph_format is not None and ph_format.idx is not None:
                layout = slide.slide_layout
                for layout_ph in layout.placeholders:
                    if layout_ph.placeholder_format.idx == ph_format.idx:
                        v = getattr(layout_ph, attr, None)
                        if v is not None:
                            return v
                        master = layout.slide_master
                        for master_ph in master.placeholders:
                            if master_ph.placeholder_format.idx == ph_format.idx:
                                mv = getattr(master_ph, attr, None)
                                if mv is not None:
                                    return mv
        except Exception:
            pass
        return None

    left = chain("left")
    top = chain("top")
    width = chain("width")
    height = chain("height")

    prs = slide.part.package.presentation_part.presentation
    slide_w = prs.slide_width or 9144000
    slide_h = prs.slide_height or 6858000
    if left is None:
        left = 0
    if top is None:
        top = 0
    if width is None:
        width = int(slide_w * 0.4)
    if height is None:
        height = int(slide_h * 0.3)
    return int(left), int(top), int(width), int(height)


def _replace_image_shape(slide, shape, image_path: Path) -> None:
    """Swap a Picture shape's image while preserving everything else.

    Strategy: register the new image as a slide-part rel (use the
    public add_picture API as the registration mechanism, then
    immediately discard the temporary shape), then rewire the
    existing shape's <a:blip r:embed> at the new rel id. Crop,
    border, rotation, shadow, transparency, picture effects, alt
    text — all preserved.

    Falls back to remove+re-add (the v3 behaviour) if the existing
    shape isn't a standard blipFill picture or if anything else
    unexpected happens.
    """
    try:
        from pptx.oxml.ns import qn
    except ImportError:
        _replace_image_shape_legacy(slide, shape, image_path)
        return

    blipFill = shape._element.find(qn("p:blipFill"))
    if blipFill is None:
        _replace_image_shape_legacy(slide, shape, image_path)
        return
    blip = blipFill.find(qn("a:blip"))
    if blip is None:
        _replace_image_shape_legacy(slide, shape, image_path)
        return

    embed_attr = qn("r:embed")

    # Use add_picture as the registration vehicle: it imports the
    # image into the slide's rels properly across all python-pptx
    # versions. We grab the rel id, then remove the temporary shape
    # — the image part stays bound to the slide via the rel.
    try:
        _l, _t, _w, _h = _resolve_shape_geometry(slide, shape)
        tmp_pic = slide.shapes.add_picture(
            str(image_path),
            left=0, top=0,
            width=_w, height=_h,
        )
        tmp_blip = tmp_pic._element.find(".//" + qn("a:blip"))
        if tmp_blip is None:
            tmp_pic._element.getparent().remove(tmp_pic._element)
            _replace_image_shape_legacy(slide, shape, image_path)
            return
        new_rid = tmp_blip.get(embed_attr)
        tmp_pic._element.getparent().remove(tmp_pic._element)
    except Exception:
        _replace_image_shape_legacy(slide, shape, image_path)
        return

    if not new_rid:
        _replace_image_shape_legacy(slide, shape, image_path)
        return

    blip.set(embed_attr, new_rid)


def _degrade_styled_value(value: Any) -> tuple[Any, list[str]]:
    """Reduce a v4-styled slot value to its v3-compatible primitive.

    The agent can emit (per SKILL.md v4):
      - {"text": "...", "color_role": "...", "bold": true, ...}
      - {"runs": [{"text": "X", "bold": true}, {"text": " Y"}]}
      - {"asset": "asset_<id>", "recolor": {"#ff0000": "accent"}}

    Full handling (colour overrides, per-run formatting, image recolour)
    lands in Phase D — until then we extract the primitive payload and
    emit a one-line warning about what's being dropped. Existing
    string/list/asset_id values pass through unchanged.
    """
    if not isinstance(value, dict):
        return value, []

    warnings: list[str] = []
    if "runs" in value and isinstance(value["runs"], list):
        text = "".join(
            str(r.get("text", "")) for r in value["runs"] if isinstance(r, dict)
        )
        warnings.append(
            f"per-run formatting not yet honored — flattened to plain text {text!r}"
        )
        return text, warnings
    if "asset" in value:
        aid = str(value["asset"])
        ignored = sorted(k for k in value if k != "asset")
        if ignored:
            warnings.append(
                f"image overrides not yet honored — ignoring {ignored}"
            )
        return aid, warnings
    if "text" in value:
        text = str(value["text"])
        ignored = sorted(k for k in value if k != "text")
        if ignored:
            warnings.append(
                f"text styling not yet honored — ignoring {ignored}"
            )
        return text, warnings

    warnings.append(f"unrecognised slot value shape: {value!r}")
    return value, warnings


def _apply_slot_value(slide, slot_id: str, value: Any, kind_hint: str | None) -> list[str]:
    """Apply one slot. Returns a list of warning strings (not fatal)."""
    warnings: list[str] = []
    shape = _find_shape_by_name(slide, slot_id)
    if shape is None:
        warnings.append(f"slot '{slot_id}' not found on slide (no shape with that name)")
        return warnings

    # Table slots use a list-of-lists shape that the degrade-and-flatten
    # pre-pass below would corrupt — handle them up-front. We trust the
    # slot's kind hint, but also detect has_table on the shape as a
    # fallback when the meta is silent (pre-D3 templates).
    is_table_slot = kind_hint == "table" or (
        kind_hint is None and getattr(shape, "has_table", False)
    )
    if is_table_slot:
        cells = value
        if isinstance(value, dict):
            ignored = sorted(k for k in value if k != "cells")
            if ignored:
                warnings.append(
                    f"slot '{slot_id}': table slot only honors 'cells'; ignoring {ignored}"
                )
            cells = value.get("cells")
        if not isinstance(cells, list) or not all(isinstance(r, list) for r in cells):
            warnings.append(
                f"slot '{slot_id}': table slot expects list-of-lists "
                f"(got {type(value).__name__}); leaving template cells unchanged"
            )
            return warnings
        warnings.extend(_fill_table_shape(shape, cells))
        return warnings

    # v4: degrade styled/per-run/image-override dicts to v3 primitives.
    # Bullets lists may carry styled dicts per item — degrade each one.
    if isinstance(value, list):
        flat: list = []
        for item in value:
            primitive, item_warns = _degrade_styled_value(item)
            flat.append(primitive)
            warnings.extend(item_warns)
        value = flat
    else:
        value, scalar_warns = _degrade_styled_value(value)
        warnings.extend(scalar_warns)

    kind = kind_hint
    if kind is None:
        # Infer from shape / value.
        if isinstance(value, list):
            kind = "bullets"
        elif (
            getattr(shape, "shape_type", None) is not None
            and str(shape.shape_type).endswith("PICTURE")
        ) or (
            isinstance(value, str) and value.startswith("asset_")
        ):
            kind = "image"
        else:
            kind = "text"

    if kind == "image":
        aid = value if isinstance(value, str) else ""
        if not aid.startswith("asset_"):
            warnings.append(f"slot '{slot_id}': image slot expects asset_<id>, got {value!r}")
            return warnings
        bin_path = asset_path(aid)
        if bin_path is None:
            warnings.append(f"slot '{slot_id}': asset {aid} not found in bundle")
            return warnings
        try:
            _replace_image_shape(slide, shape, bin_path)
        except Exception as e:
            warnings.append(f"slot '{slot_id}': image replace failed: {type(e).__name__}: {e}")
        return warnings

    if kind == "bullets":
        items = value if isinstance(value, list) else [str(value)]
        _fill_bullets_shape(shape, [str(v) for v in items])
        return warnings

    # text
    _fill_text_shape(shape, str(value))
    return warnings


# --- cross-deck slide copy --------------------------------------------------


# ---------------------------------------------------------------------------
# Deck-theme normalisation (v4 — D5)
# ---------------------------------------------------------------------------
#
# When we copy a foreign deck's slide / atom into the host package, any
# `<a:schemeClr val="..."/>` reference resolves against the HOST's theme
# at render time — so a "brand pink accent1" in the source might paint as
# a "navy accent1" in the host. We mitigate by remapping the small set of
# semantically-named slots (primary / accent / text / background) per the
# decks' aliases, so a shape that used "the source deck's accent" still
# uses "the host deck's accent" after the copy. Other clrScheme slots
# (accent2-6, dk3, …) we accept as-is — there is no canonical semantic
# mapping for them.

# Both names are valid in <a:schemeClr val>: theme-side (dk1/lt1/dk2/lt2)
# and use-site (tx1/bg1/tx2/bg2). Normalize both forms to the theme-side.
_SCHEME_NORMALIZE = {
    "tx1": "dk1", "bg1": "lt1", "tx2": "dk2", "bg2": "lt2",
}


def _norm_scheme_slot(val: str) -> str:
    v = (val or "").lower()
    return _SCHEME_NORMALIZE.get(v, v)


def _load_deck_theme(deck_name: str) -> dict | None:
    """Load decks/<deck>/theme.yaml from the bundle; None if missing."""
    if not deck_name:
        return None
    p = bundle_root() / "decks" / deck_name / "theme.yaml"
    if not p.exists():
        return None
    try:
        return yaml.safe_load(p.read_text(encoding="utf-8")) or {}
    except Exception:
        return None


def _resolve_template_deck_theme(template_meta: dict) -> dict:
    """Return the full deck theme for a template; fall back to a
    minimal synthesized dict from template_meta.theme_colors if the
    deck's theme.yaml isn't in the bundle (e.g. older builds)."""
    src = (template_meta.get("sources") or [{}])[0] if template_meta.get("sources") else {}
    deck_name = src.get("deck") or ""
    deck_theme = _load_deck_theme(deck_name)
    if deck_theme:
        return deck_theme
    palette: dict = {}
    aliases: dict = {}
    for role, hex_val in (template_meta.get("theme_colors") or {}).items():
        if hex_val:
            palette[role] = hex_val
            aliases[role] = role
    return {"palette": palette, "aliases": aliases}


def _build_scheme_remap(
    source_theme: dict | None, host_theme: dict | None
) -> dict[str, str]:
    """Build a clrScheme-slot remap aligning source.aliases with host.aliases.

    For each shared role where the slot differs between decks, add a
    rule. Keys are normalised (dk1/lt1/dk2/lt2) so both naming forms
    in the XML get matched.
    """
    if not source_theme or not host_theme:
        return {}
    src_aliases = source_theme.get("aliases") or {}
    host_aliases = host_theme.get("aliases") or {}
    remap: dict[str, str] = {}
    for role in ("primary", "accent", "text", "background"):
        src_slot = src_aliases.get(role)
        host_slot = host_aliases.get(role)
        if not src_slot or not host_slot:
            continue
        src_norm = _norm_scheme_slot(src_slot)
        host_norm = _norm_scheme_slot(host_slot)
        if src_norm != host_norm:
            remap[src_norm] = host_slot
    return remap


def _apply_scheme_remap(el, remap: dict[str, str]) -> int:
    """Rewrite <a:schemeClr val=...> per remap; returns count of substitutions."""
    if not remap:
        return 0
    try:
        from pptx.oxml.ns import qn
    except ImportError:
        return 0
    count = 0
    for sc in el.findall(".//" + qn("a:schemeClr")):
        val = (sc.get("val") or "").lower()
        norm = _norm_scheme_slot(val)
        if norm in remap:
            sc.set("val", remap[norm])
            count += 1
    return count


# ---------------------------------------------------------------------------
# v4.1 — surgical theme-font remap (D5 extension)
# ---------------------------------------------------------------------------
#
# Theme refs like <a:latin typeface="+mj-lt"/> already self-resolve to the
# host's major font after a cross-deck copy. Explicit typefaces don't —
# they survive the copy verbatim. Keynote-exported decks are the common
# offender: Keynote bakes "Helvetica Neue" in as an explicit typeface
# even though it IS that deck's theme major font. On a host whose theme
# major is e.g. Inter, the copied text then stubbornly renders Helvetica.
#
# Surgical policy: only rewrite an explicit typeface when it matches the
# *source* deck theme's major or minor. Anything else (Courier code,
# Comic Sans header) is preserved — assume the author meant it.

def _build_font_remap(
    source_theme: dict | None, host_theme: dict | None
) -> dict[str, str]:
    """Build a typeface remap aligning source major/minor to host's.

    Returns ``{lowercased_source_typeface: host_typeface}`` for the
    major and minor roles only when both sides have a font defined and
    they differ. One-off explicit fonts (i.e. typefaces that don't
    match the source theme's major/minor) are intentionally *not* in
    the remap and survive the copy unchanged.
    """
    if not source_theme or not host_theme:
        return {}
    src_fonts = source_theme.get("fonts") or {}
    host_fonts = host_theme.get("fonts") or {}
    remap: dict[str, str] = {}
    for role in ("major", "minor"):
        src = (src_fonts.get(role) or "").strip()
        host = (host_fonts.get(role) or "").strip()
        if not src or not host:
            continue
        if src.lower() == host.lower():
            continue
        remap[src.lower()] = host
    return remap


def _apply_font_remap(el, remap: dict[str, str]) -> int:
    """Rewrite ``<a:latin|ea|cs typeface="..."/>`` per remap; returns count.

    Only touches explicit typefaces. Theme refs (typefaces starting
    with ``+`` like ``+mj-lt``) self-resolve at render time and are
    skipped. Matching is case-insensitive; the substituted value
    preserves the host theme's casing.
    """
    if not remap:
        return 0
    try:
        from pptx.oxml.ns import qn
    except ImportError:
        return 0
    count = 0
    for tag in ("a:latin", "a:ea", "a:cs"):
        for node in el.findall(".//" + qn(tag)):
            typeface = (node.get("typeface") or "").strip()
            if not typeface or typeface.startswith("+"):
                continue
            host = remap.get(typeface.lower())
            if host and host != typeface:
                node.set("typeface", host)
                count += 1
    return count


# ---------------------------------------------------------------------------
# Rel-aware shape import (Phase 1 of the multi-slide compose fix)
# ---------------------------------------------------------------------------
#
# A naive `deepcopy(shape._element)` brings the source's rId references
# along verbatim. Those rIds live in the source slide's *local* rels file
# and mean nothing in the destination — PowerPoint then refuses to render
# the shape and pops the "couldn't read some content" dialog.
#
# _import_shape_xml handles this by:
#   1) deepcopy the shape XML,
#   2) for each rId-bearing attribute, look up the rel on the source part
#      and re-register it on the destination part (which auto-imports the
#      target part and any rels that come with it),
#   3) rewrite the attribute to the destination-local rId.
#
# Transitive rels (chart → embedded xlsx → image) ride along automatically
# because python-pptx packages parts together with their own rels graph
# when `relate_to(target_part, reltype)` registers a new internal rel.

_OPC_REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_REL_ATTRS = (
    f"{{{_OPC_REL_NS}}}embed",
    f"{{{_OPC_REL_NS}}}link",
    f"{{{_OPC_REL_NS}}}id",
)


def _iter_rel_attr_holders(el):
    """Yield (element, attr_qname) for every rId-bearing attribute under el."""
    for descendant in el.iter():
        for q in _REL_ATTRS:
            if q in descendant.attrib:
                yield descendant, q


def _existing_part_for_blob(package, partname, blob: bytes):
    """Return the dest-package part already at `partname` iff its blob
    matches `blob` (so the import can dedupe instead of colliding).
    None when no part exists there OR the existing part has different
    bytes — caller must then allocate a fresh partname."""
    for part in package.iter_parts():
        if part.partname == partname and getattr(part, "blob", None) == blob:
            return part
    return None


def _has_colliding_part(package, partname) -> bool:
    for part in package.iter_parts():
        if part.partname == partname:
            return True
    return False


def _clone_part_into(dest_package, src_part):
    """Materialise a copy of `src_part` inside `dest_package`. When the
    partname is free (or already holds the same bytes), reuse it as-is.
    When it would collide with different bytes, allocate a fresh
    partname so the OPC package stays valid.

    Returns a Part bound to dest_package."""
    from pptx.opc.package import Part
    from pptx.opc.packuri import PackURI

    src_partname = src_part.partname
    src_blob = src_part.blob
    src_ctype = src_part.content_type

    existing = _existing_part_for_blob(dest_package, src_partname, src_blob)
    if existing is not None:
        return existing

    if not _has_colliding_part(dest_package, src_partname):
        # Free slot — rebind by re-creating the part inside dest_package
        # using the same class so chart/image/etc. behaviours are kept.
        return type(src_part).load(src_partname, src_ctype, dest_package, src_blob)

    # Collision: allocate a fresh partname matching the original prefix
    # (e.g. /ppt/media/image%d.jpeg). PackURI exposes .baseURI + .ext.
    base_uri = src_partname.baseURI            # e.g. /ppt/media
    ext = src_partname.ext                     # e.g. jpeg
    # Strip trailing digits from the original filename to get the tmpl
    # stem ("image" from "image1"). Falls back to "part" if the source
    # used a non-numbered name.
    stem = src_partname.filename.split(".")[0]
    stem = re.sub(r"\d+$", "", stem) or "part"
    tmpl = f"{base_uri}/{stem}%d.{ext}" if ext else f"{base_uri}/{stem}%d"
    new_partname = dest_package.next_partname(tmpl)
    return type(src_part).load(
        PackURI(str(new_partname)), src_ctype, dest_package, src_blob,
    )


def _import_part_rel(src_part, dest_part, old_rid: str) -> str | None:
    """Look up old_rid on src_part, register the same target on dest_part,
    return the new dest-local rId. None if the rel cannot be resolved.

    Handles partname collisions across source decks by cloning the
    target onto dest_package with a fresh partname when needed."""
    try:
        src_rel = src_part.rels[old_rid]
    except KeyError:
        return None
    reltype = src_rel.reltype
    try:
        if src_rel.is_external:
            return dest_part.relate_to(src_rel.target_ref, reltype, is_external=True)
        target = src_rel.target_part
        dest_pkg = dest_part.package
        # If target was loaded from a different package OR its partname
        # would clash with an existing part in dest, clone into dest.
        if target.package is not dest_pkg or _has_colliding_part(
            dest_pkg, target.partname
        ):
            target = _clone_part_into(dest_pkg, target)
        return dest_part.relate_to(target, reltype)
    except Exception:
        # Unknown reltype or part class python-pptx can't handle — log
        # and skip so the rest of the shape import still succeeds.
        return None


def _rewrite_rel_attrs_in_subtree(src_part, dest_part, subtree) -> list[str]:
    """For every rId attr in `subtree`, import the rel onto dest_part and
    rewrite the attr value. Returns a list of warnings for refs that
    couldn't be resolved."""
    warnings: list[str] = []
    rid_subs: dict[str, str] = {}
    for el, attr_q in _iter_rel_attr_holders(subtree):
        old_rid = el.get(attr_q)
        if old_rid is None or old_rid in rid_subs:
            continue
        new_rid = _import_part_rel(src_part, dest_part, old_rid)
        if new_rid is None:
            warnings.append(
                f"unresolved rel {old_rid!r} on <{el.tag.rpartition('}')[2]} "
                f"{attr_q.rpartition('}')[2]}=...> — shape kept but ref dropped"
            )
            continue
        rid_subs[old_rid] = new_rid
    for el, attr_q in _iter_rel_attr_holders(subtree):
        old_rid = el.get(attr_q)
        if old_rid in rid_subs:
            el.set(attr_q, rid_subs[old_rid])
        elif old_rid is not None:
            # Drop the broken ref so OOXML doesn't see a phantom rId.
            del el.attrib[attr_q]
    return warnings


def _import_shape_xml(src_slide_part, dest_slide_part, dest_sptree, src_shape_el):
    """Import one shape: deepcopy, rewrite rIds against dest part, append.

    Returns a list of warnings (empty on the happy path)."""
    new_el = copy.deepcopy(src_shape_el)
    warnings = _rewrite_rel_attrs_in_subtree(src_slide_part, dest_slide_part, new_el)
    dest_sptree.append(new_el)
    return new_el, warnings


# ---------------------------------------------------------------------------
# Background flatten (Phase 3 of the multi-slide compose fix)
# ---------------------------------------------------------------------------
#
# PowerPoint resolves slide backgrounds through an inheritance chain:
#   slide → layout → master → theme (via <p:bgRef idx>)
# When we graft slides from one deck onto another, the destination's
# layout/master/theme is used — so the source's background is lost.
#
# _flatten_slide_background materialises the effective background as an
# explicit <p:bg> on the source slide before graft, so the background
# copies along with the shape tree.

_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


def _qn(tag: str) -> str:
    prefix, _, local = tag.partition(":")
    ns = {"a": _A_NS, "p": _P_NS, "r": _OPC_REL_NS}.get(prefix)
    return f"{{{ns}}}{local}" if ns else tag


def _find_bg_element(cSld_el):
    if cSld_el is None:
        return None
    for child in cSld_el:
        if child.tag == _qn("p:bg"):
            return child
    return None


def _theme_part_for(slide_or_layout_or_master):
    """Return the SlideMaster theme part (the only theme that defines
    bgFillStyleLst entries used by <p:bgRef>)."""
    try:
        if hasattr(slide_or_layout_or_master, "slide_layout"):
            master = slide_or_layout_or_master.slide_layout.slide_master
        elif hasattr(slide_or_layout_or_master, "slide_master"):
            master = slide_or_layout_or_master.slide_master
        else:
            master = slide_or_layout_or_master
        # python-pptx exposes the theme part as master.element.part's
        # related theme part; easiest is to read its rels for the theme.
        master_part = master.part
        for rel in master_part.rels.values():
            if rel.reltype.endswith("/theme"):
                return rel.target_part
    except Exception:
        return None
    return None


def _resolve_theme_bg_ref(theme_part, idx: int):
    """Return the theme's bgFillStyleLst[idx-1001] element (deepcopy),
    suitable for embedding under <p:bgPr>. None on miss."""
    if theme_part is None or idx is None:
        return None
    try:
        from lxml import etree
        root = etree.fromstring(theme_part.blob)
    except Exception:
        return None
    bg_lst = root.find(
        f".//{_qn('a:fmtScheme')}/{_qn('a:bgFillStyleLst')}"
    )
    if bg_lst is None:
        return None
    children = list(bg_lst)
    # OOXML spec: bgRef@idx is 1001-based for bgFillStyleLst.
    pos = idx - 1001
    if pos < 0 or pos >= len(children):
        return None
    return copy.deepcopy(children[pos])


def _build_bg_pr_from_fill(fill_el):
    """Wrap a fill element (a:solidFill/a:gradFill/a:blipFill/a:pattFill)
    in a fresh <p:bgPr> ready to be the sole child of <p:bg>."""
    from lxml import etree
    bg_pr = etree.SubElement(etree.Element(_qn("p:bg")), _qn("p:bgPr"))
    bg_pr.append(fill_el)
    # OOXML schema: <p:bgPr> must end with <a:effectLst/> (can be empty).
    etree.SubElement(bg_pr, _qn("a:effectLst"))
    return bg_pr.getparent()  # the <p:bg> wrapper


def _find_effective_bg(slide):
    """Walk slide → layout → master and return the first <p:bg> element
    found (with the part it came from, for rel-import).

    Returns (bg_el_deepcopy, source_part) or (None, None)."""
    chain = []
    try:
        chain.append((slide._element.find(_qn("p:cSld")), slide.part))
    except Exception:
        pass
    try:
        layout = slide.slide_layout
        chain.append((layout._element.find(_qn("p:cSld")), layout.part))
    except Exception:
        pass
    try:
        master = slide.slide_layout.slide_master
        chain.append((master._element.find(_qn("p:cSld")), master.part))
    except Exception:
        pass
    for cSld_el, part in chain:
        bg = _find_bg_element(cSld_el)
        if bg is not None:
            return copy.deepcopy(bg), part
    return None, None


def _flatten_slide_background(slide) -> list[str]:
    """Materialise the slide's effective background into an explicit
    <p:bg> on the slide. Idempotent. Returns warnings list."""
    warnings: list[str] = []
    sld_el = slide._element
    cSld_el = sld_el.find(_qn("p:cSld"))
    if cSld_el is None:
        return warnings
    if _find_bg_element(cSld_el) is not None:
        return warnings  # already explicit

    bg_el, src_part = _find_effective_bg(slide)
    if bg_el is None or src_part is None:
        return warnings  # nothing to flatten

    # Resolve <p:bgRef> → inline <p:bgPr> from the theme.
    bg_ref = bg_el.find(_qn("p:bgRef"))
    if bg_ref is not None:
        try:
            idx = int(bg_ref.get("idx") or 0)
        except ValueError:
            idx = 0
        theme_part = _theme_part_for(slide)
        fill_el = _resolve_theme_bg_ref(theme_part, idx)
        if fill_el is None:
            warnings.append(
                f"bgRef idx={idx} could not be resolved against theme — "
                f"background may render blank"
            )
            return warnings
        # Carry over scheme-colour overrides from the bgRef (e.g. when
        # the slide says "use theme bg #2 but recolor accent1 → bg2").
        # bgRef wraps an inline colour spec as last child; preserve it.
        new_bg = _build_bg_pr_from_fill(fill_el)
        src_part = theme_part  # rels inside the theme blob, if any
    else:
        new_bg = bg_el

    # Import any rels referenced inside the bg subtree (e.g. blipFill
    # r:embed pointing at a media part on the source layout/master).
    warnings.extend(
        _rewrite_rel_attrs_in_subtree(src_part, slide.part, new_bg)
    )

    # <p:bg> must be the first child of <p:cSld> per the OOXML schema.
    cSld_el.insert(0, new_bg)
    return warnings


def _copy_slide_into(
    dest_prs: Presentation,
    src_slide_pptx: Path,
    source_theme: dict | None = None,
    host_theme: dict | None = None,
):
    """Copy the (single) slide from src_slide_pptx into dest_prs and return
    the new slide.

    Strategy: open the source deck, flatten its background so it's self-
    contained, copy its first slide's shape tree onto a blank layout in
    the destination via rel-aware import, then remap semantic clrScheme
    slots (D5) so the copied shapes stay consistent with the host's
    brand colours.
    """
    src_prs = Presentation(str(src_slide_pptx))
    if len(src_prs.slides) == 0:
        raise SystemExit(f"{src_slide_pptx}: no slides")
    src_slide = src_prs.slides[0]

    # Phase 3: make the source slide's background self-contained before
    # we graft. Otherwise the destination's blank layout supplies a
    # blank background and the original gradient/image is lost.
    _flatten_slide_background(src_slide)

    # Pick a blank-ish layout in dest. Prefer the layout with the fewest
    # placeholders to minimise interference with the copied content.
    dest_layout = None
    best_count = None
    for layout in dest_prs.slide_layouts:
        try:
            count = len(layout.placeholders)
        except Exception:
            count = 99
        if best_count is None or count < best_count:
            best_count = count
            dest_layout = layout
    if dest_layout is None:
        dest_layout = dest_prs.slide_layouts[0]

    new_slide = dest_prs.slides.add_slide(dest_layout)

    # Clear any placeholders the layout brought along — we want a clean canvas.
    for shp in list(new_slide.shapes):
        shp._element.getparent().remove(shp._element)

    # Carry the flattened background over to the new slide's <p:cSld>.
    src_cSld = src_slide._element.find(_qn("p:cSld"))
    dest_cSld = new_slide._element.find(_qn("p:cSld"))
    if src_cSld is not None and dest_cSld is not None:
        src_bg = _find_bg_element(src_cSld)
        if src_bg is not None and _find_bg_element(dest_cSld) is None:
            new_bg = copy.deepcopy(src_bg)
            # Rels inside the bg (blip fills) were imported onto
            # src_slide.part during flatten; re-import onto new_slide.part.
            _rewrite_rel_attrs_in_subtree(
                src_slide.part, new_slide.part, new_bg
            )
            dest_cSld.insert(0, new_bg)

    remap = _build_scheme_remap(source_theme, host_theme)
    font_remap = _build_font_remap(source_theme, host_theme)

    # Rel-aware shape import. Handles pictures, group shapes, hyperlinks,
    # picture-filled auto-shapes, charts, SmartArt — anything that carries
    # an rId reference inside its XML — by re-registering the underlying
    # parts on the destination slide and rewriting rIds in place.
    for shape in src_slide.shapes:
        new_el, _ws = _import_shape_xml(
            src_slide_part=src_slide.part,
            dest_slide_part=new_slide.part,
            dest_sptree=new_slide.shapes._spTree,
            src_shape_el=shape._element,
        )
        if remap:
            _apply_scheme_remap(new_el, remap)
        if font_remap:
            _apply_font_remap(new_el, font_remap)

    return new_slide


# ---------------------------------------------------------------------------
# Compose-mode (v4) — custom slides assembled from atoms + native shapes
# ---------------------------------------------------------------------------


def _frac_to_emu(spec: dict, key: str, total_emu: int, fallback: float) -> int:
    v = spec.get(key)
    if not isinstance(v, (int, float)) or v < 0:
        v = fallback
    return int(total_emu * float(v))


def _set_atom_geometry(el, x_emu: int, y_emu: int, w_emu: int, h_emu: int) -> None:
    """Rewrite (or insert) the offset+extent on a captured shape's XML.

    Tables and other GraphicFrame atoms use a direct ``p:xfrm`` child.
    Auto-shapes / freeforms (``p:sp``) use ``p:spPr/a:xfrm``. We handle
    both shapes silently — if neither path exists the atom keeps its
    source geometry rather than crashing.
    """
    try:
        from pptx.oxml.ns import qn
        from lxml import etree
    except ImportError:
        return
    xfrm = el.find(qn("p:xfrm"))
    if xfrm is None:
        sp_pr = el.find(qn("p:spPr"))
        if sp_pr is not None:
            xfrm = sp_pr.find(qn("a:xfrm"))
            if xfrm is None:
                xfrm = etree.SubElement(sp_pr, qn("a:xfrm"))
    if xfrm is None:
        return
    off = xfrm.find(qn("a:off"))
    if off is None:
        off = etree.SubElement(xfrm, qn("a:off"))
    off.set("x", str(x_emu))
    off.set("y", str(y_emu))
    ext = xfrm.find(qn("a:ext"))
    if ext is None:
        ext = etree.SubElement(xfrm, qn("a:ext"))
    ext.set("cx", str(w_emu))
    ext.set("cy", str(h_emu))


def _apply_recolor_xml(el, recolor: dict, deck_theme: dict | None = None) -> tuple[int, list[str]]:
    """Rewrite ``<a:srgbClr val="...">`` matches per a recolor map.

    Keys are source hex codes (with or without '#'). Values can be:
      - another hex code → direct substitution
      - a colour-role token ('primary', 'accent', 'text', 'background'
        or any clrScheme slot) → resolved against ``deck_theme`` if
        provided; warns and skips if unresolvable.
    Returns ``(substitutions_applied, warnings)``.
    """
    try:
        from pptx.oxml.ns import qn
    except ImportError:
        return 0, []
    warnings: list[str] = []
    palette = (deck_theme or {}).get("palette") or {}
    aliases = (deck_theme or {}).get("aliases") or {}

    def _resolve_target(raw: str) -> str | None:
        if not isinstance(raw, str):
            return None
        s = raw.strip().lstrip("#").upper()
        if len(s) == 6 and all(c in "0123456789ABCDEF" for c in s):
            return s
        # role / clrScheme token
        token = raw.strip().lower()
        alias_target = aliases.get(token)
        hex_val = palette.get(alias_target) if alias_target else palette.get(token)
        if not hex_val:
            return None
        return hex_val.lstrip("#").upper()

    norm_map: dict[str, str] = {}
    for src, tgt in recolor.items():
        if not isinstance(src, str):
            continue
        src_n = src.strip().lstrip("#").upper()
        if len(src_n) != 6:
            warnings.append(f"recolor source {src!r} is not a 6-digit hex; skipping")
            continue
        tgt_n = _resolve_target(tgt)
        if tgt_n is None:
            warnings.append(
                f"recolor target {tgt!r} could not be resolved (need hex or known role); skipping"
            )
            continue
        norm_map[src_n] = tgt_n

    if not norm_map:
        return 0, warnings

    applied = 0
    for clr in el.findall(".//" + qn("a:srgbClr")):
        val = (clr.get("val") or "").upper()
        if val in norm_map:
            clr.set("val", norm_map[val])
            applied += 1
    return applied, warnings


def _place_atom(
    slide,
    spec: dict,
    slide_w_emu: int,
    slide_h_emu: int,
    host_theme: dict | None = None,
) -> list[str]:
    """Place one atom on a slide at a fractional bbox.

    Supported spec keys:
      atom              asset id (required)
      x, y, w, h        fractions of slide; default to a top-left 40x30% box
      kind              optional override; defaults to the asset's stored kind
      recolor           {src_hex: target_hex_or_role} — rewrites srgbClr fills
      cells             list-of-lists (kind=table only) — post-place cell fill

    Picture atoms (raster/vector) re-import via ``add_picture`` so the
    image part lives in the host package. XML atoms (table, callout,
    freeform) graft as fragments. Charts / smartart are NOT yet
    placeable — they reference external parts that need a related-parts
    copy pass; we warn + skip rather than emit a half-broken slide.
    """
    warnings: list[str] = []
    aid = spec.get("atom")
    if not isinstance(aid, str) or not aid:
        warnings.append(f"compose-mode shape missing 'atom' id: {spec!r}")
        return warnings

    meta_p = asset_meta_path(aid)
    bin_p = asset_path(aid)
    if not meta_p.exists() or bin_p is None:
        warnings.append(f"compose-mode: asset {aid!r} not found in bundle")
        return warnings
    meta = yaml.safe_load(meta_p.read_text(encoding="utf-8")) or {}
    kind = spec.get("kind") or meta.get("kind") or ""

    x_emu = _frac_to_emu(spec, "x", slide_w_emu, 0.05)
    y_emu = _frac_to_emu(spec, "y", slide_h_emu, 0.05)
    w_emu = _frac_to_emu(spec, "w", slide_w_emu, 0.4)
    h_emu = _frac_to_emu(spec, "h", slide_h_emu, 0.3)

    if bin_p.suffix.lower() != ".xml":
        # Picture / vector — straightforward add_picture path.
        try:
            new_pic = slide.shapes.add_picture(
                str(bin_p), x_emu, y_emu, width=w_emu, height=h_emu
            )
            new_pic.name = aid
        except Exception as e:
            warnings.append(f"compose-mode: failed to place picture {aid!r}: {e}")
            return warnings
        if spec.get("recolor"):
            warnings.append(
                f"compose-mode {aid!r}: recolor on picture atoms not yet honored"
            )
        return warnings

    # XML atom path — graft as fragment.
    if kind in ("chart", "smartart"):
        warnings.append(
            f"compose-mode: {kind} atom {aid!r} not yet placeable "
            f"(related-parts copying deferred); skipping"
        )
        return warnings

    # Use pptx's parser so the element wraps with the right custom
    # class (CT_GraphicalObjectFrame, CT_Shape, …); a raw lxml element
    # would crash SlideShapeFactory's `has_ph_elm` check on lookup.
    try:
        from pptx.oxml import parse_xml
    except ImportError:
        warnings.append(f"compose-mode: pptx.oxml unavailable; cannot graft atom {aid!r}")
        return warnings

    try:
        new_el = parse_xml(bin_p.read_bytes())
    except Exception as e:
        warnings.append(f"compose-mode: atom {aid!r} XML failed to parse: {e}")
        return warnings

    _set_atom_geometry(new_el, x_emu, y_emu, w_emu, h_emu)

    # D5: remap semantic clrScheme slots when the atom came from a
    # different deck than the host — so e.g. an accent1 reference
    # painted in the source deck's accent still paints in the host
    # deck's accent on the rendered slide.
    asset_src = (meta.get("sources") or [{}])[0] if meta.get("sources") else {}
    asset_deck_theme = _load_deck_theme(asset_src.get("deck") or "")
    if asset_deck_theme and host_theme:
        scheme_remap = _build_scheme_remap(asset_deck_theme, host_theme)
        if scheme_remap:
            _apply_scheme_remap(new_el, scheme_remap)
        # v4.1: rewrite explicit typefaces that match the source's
        # major/minor → host's major/minor. Preserves intentional
        # one-off fonts (Courier, etc.).
        font_remap = _build_font_remap(asset_deck_theme, host_theme)
        if font_remap:
            _apply_font_remap(new_el, font_remap)

    recolor = spec.get("recolor")
    if isinstance(recolor, dict) and recolor:
        applied, ws = _apply_recolor_xml(new_el, recolor, host_theme)
        warnings.extend(ws)
        if applied == 0 and not ws:
            warnings.append(
                f"compose-mode: recolor for {aid!r} matched 0 colours in atom XML"
            )

    slide.shapes._spTree.append(new_el)

    if kind == "table" and spec.get("cells") is not None:
        new_shape = slide.shapes[-1]
        cells = spec.get("cells")
        if isinstance(cells, list) and all(isinstance(r, list) for r in cells):
            warnings.extend(_fill_table_shape(new_shape, cells))
        else:
            warnings.append(
                f"compose-mode: atom {aid!r} 'cells' must be list-of-lists; skipping fill"
            )

    return warnings


def _place_native_text(slide, spec: dict, slide_w_emu: int, slide_h_emu: int) -> list[str]:
    """Add a native textbox at fractional position from a compose-mode shape spec.

    Supported spec keys (in addition to x/y/w/h):
      value or text     text to render
      bold              bool
      font_role         informational — warns (not yet honored)
      color_role        informational — warns (not yet honored)
    """
    warnings: list[str] = []
    text = spec.get("value", spec.get("text", ""))
    if not isinstance(text, str):
        text = str(text)

    x_emu = _frac_to_emu(spec, "x", slide_w_emu, 0.05)
    y_emu = _frac_to_emu(spec, "y", slide_h_emu, 0.05)
    w_emu = _frac_to_emu(spec, "w", slide_w_emu, 0.5)
    h_emu = _frac_to_emu(spec, "h", slide_h_emu, 0.1)

    try:
        tb = slide.shapes.add_textbox(x_emu, y_emu, w_emu, h_emu)
    except Exception as e:
        warnings.append(f"compose-mode text: failed to add textbox: {e}")
        return warnings
    tb.text_frame.text = _strip_bullet_prefix(text)

    if spec.get("bold"):
        for p in tb.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True

    for k in ("font_role", "color_role"):
        if spec.get(k):
            warnings.append(
                f"compose-mode text: {k}={spec[k]!r} not yet honored"
            )
    return warnings


def _compose_custom_slide(
    dest_prs: Presentation,
    entry: dict,
    host_theme: dict | None = None,
):
    """Append a blank slide and populate from ``entry.shapes``.

    Each shape spec is one of:
      - ``{"atom": <id>, ...}``  → routed to ``_place_atom``
      - ``{"kind": "text", "value": "...", ...}``  → routed to
        ``_place_native_text``

    Returns ``(new_slide, warnings)``.
    """
    warnings: list[str] = []

    dest_layout = None
    best_count = None
    for layout in dest_prs.slide_layouts:
        try:
            count = len(layout.placeholders)
        except Exception:
            count = 99
        if best_count is None or count < best_count:
            best_count = count
            dest_layout = layout
    if dest_layout is None:
        dest_layout = dest_prs.slide_layouts[0]

    new_slide = dest_prs.slides.add_slide(dest_layout)
    # Clean canvas — drop any layout-inherited placeholders.
    for shp in list(new_slide.shapes):
        shp._element.getparent().remove(shp._element)

    slide_w_emu = dest_prs.slide_width
    slide_h_emu = dest_prs.slide_height

    shapes = entry.get("shapes") or []
    for spec in shapes:
        if not isinstance(spec, dict):
            warnings.append(f"compose-mode entry: bad shape spec {spec!r}")
            continue
        if "atom" in spec:
            warnings.extend(_place_atom(new_slide, spec, slide_w_emu, slide_h_emu, host_theme))
        elif spec.get("kind") == "text":
            warnings.extend(_place_native_text(new_slide, spec, slide_w_emu, slide_h_emu))
        else:
            warnings.append(
                f"compose-mode entry: shape spec lacks 'atom' or kind=text: {spec!r}"
            )

    return new_slide, warnings


def _drop_first_slide(prs: Presentation) -> None:
    """Remove the host's original first slide after appending plan slides.

    Used when a compose-mode entry is the first plan entry — we still
    need a host pptx (for its theme/master), but we don't want its
    original first slide to leak into the output.
    """
    try:
        from pptx.oxml.ns import qn
    except ImportError:
        return
    sld_id_list = prs.slides._sldIdLst
    sld_ids = list(sld_id_list)
    if not sld_ids:
        return
    first = sld_ids[0]
    rid = first.get(qn("r:id"))
    sld_id_list.remove(first)
    if rid:
        try:
            prs.part.drop_rel(rid)
        except KeyError:
            pass


# ---------------------------------------------------------------------------
# Commands
# ---------------------------------------------------------------------------



# ===========================================================================
# v5 redesign — read-side methods (phase D)
#
# Self-contained block. Operates on either a built v5 bundle (next to
# reader.py with themes/ + skeletons/ + assets/ siblings, per phase F)
# OR directly on workspace/themes/ + workspace/skeletons/ during dev
# (so we can exercise the API end-to-end before phase E + F ship).
#
# Read-only — no deck building. Phase E owns compose-v5.
# ===========================================================================


def _v5_bundle_root() -> Path:
    """Find the v5 data root. Tries built-bundle layout first
    (themes/skeletons siblings next to reader.py), falls back to the
    authoring workspace for dev. Returns None if neither is present.
    """
    here = bundle_root()
    if (here / "themes").is_dir() and (here / "skeletons").is_dir():
        return here
    # Dev fallback: repo/authoring/workspace/
    ws = here.parent / "authoring" / "workspace"
    if (ws / "themes").is_dir() and (ws / "skeletons").is_dir():
        return ws
    return here  # caller will see empty results


def _v5_themes_dir() -> Path:
    return _v5_bundle_root() / "themes"


def _v5_skeletons_dir() -> Path:
    return _v5_bundle_root() / "skeletons"


def _v5_assets_dir() -> Path:
    # Assets live next to themes/skeletons in workspace; in a built
    # bundle they're under "assets/" sibling. Same path either way.
    return _v5_bundle_root() / "assets"


def _v5_load_yaml(path: Path) -> dict | None:
    if not path.exists():
        return None
    try:
        return yaml.safe_load(path.read_text(encoding="utf-8")) or {}
    except Exception:
        return None


def _v5_iter_skeletons() -> list[dict]:
    root = _v5_skeletons_dir()
    if not root.exists():
        return []
    out: list[dict] = []
    for d in sorted(root.iterdir()):
        if not d.is_dir():
            continue
        sk = _v5_load_yaml(d / "skeleton.yaml")
        if sk:
            out.append(sk)
    return out


def _v5_load_skeleton(skeleton_id: str) -> dict | None:
    return _v5_load_yaml(_v5_skeletons_dir() / skeleton_id / "skeleton.yaml")


def _v5_load_theme(theme_id: str) -> dict | None:
    return _v5_load_yaml(_v5_themes_dir() / theme_id / "theme.yaml")


def _v5_load_asset_meta(asset_id: str) -> dict | None:
    return _v5_load_yaml(_v5_assets_dir() / f"{asset_id}.yaml")


# --- Constraint helpers (single source of truth for fit logic) -------------


def _v5_check_text_fit(text: str, constraints: dict) -> tuple[bool, str, int]:
    """Returns (fits, reason, headroom_chars). Headroom positive = under
    constraint, negative = over.
    """
    max_chars = constraints.get("max_chars")
    if max_chars is None:
        return True, "", 0
    n = len(text or "")
    headroom = max_chars - n
    if headroom < 0:
        return False, f"{n} chars > max_chars {max_chars}", headroom
    return True, "", headroom


def _v5_check_bullets_fit(items: list, constraints: dict) -> tuple[bool, str, dict]:
    max_items = constraints.get("max_items")
    max_chars_per_item = constraints.get("max_chars_per_item")
    n_items = len(items or [])
    headroom = {"items": (max_items - n_items) if max_items is not None else 0}
    if max_items is not None and n_items > max_items:
        return False, f"{n_items} items > max_items {max_items}", headroom
    if max_chars_per_item is not None:
        for i, it in enumerate(items or []):
            if len(str(it)) > max_chars_per_item:
                return (
                    False,
                    f"item {i} has {len(str(it))} chars > max_chars_per_item {max_chars_per_item}",
                    headroom,
                )
    return True, "", headroom


def _v5_check_table_fit(table_dict: dict, constraints: dict) -> tuple[bool, str, dict]:
    rows = table_dict.get("rows", 0)
    cols = table_dict.get("cols", 0)
    max_rows = constraints.get("max_rows")
    max_cols = constraints.get("max_cols")
    headroom = {
        "rows": (max_rows - rows) if max_rows is not None else 0,
        "cols": (max_cols - cols) if max_cols is not None else 0,
    }
    if max_rows is not None and rows > max_rows:
        return False, f"{rows} rows > max_rows {max_rows}", headroom
    if max_cols is not None and cols > max_cols:
        return False, f"{cols} cols > max_cols {max_cols}", headroom
    return True, "", headroom


def _v5_check_image_fit(asset_meta: dict, constraints: dict) -> tuple[bool, str, dict]:
    transform = {"will_crop": False, "will_resize": False}
    slot_aspect = (constraints.get("aspect") or "free").lower()
    if slot_aspect == "free":
        return True, "", transform
    aspect_targets = {"1:1": 1.0, "16:9": 16/9, "4:3": 4/3, "3:4": 3/4, "9:16": 9/16}
    target = aspect_targets.get(slot_aspect)
    if target is None:
        return True, f"unknown aspect {slot_aspect!r} — letting through", transform
    w = asset_meta.get("width") or asset_meta.get("dimensions", {}).get("width") or 0
    h = asset_meta.get("height") or asset_meta.get("dimensions", {}).get("height") or 0
    if w <= 0 or h <= 0:
        return True, "asset dims unknown — best-effort fit", transform
    asset_aspect = w / h
    if abs(asset_aspect - target) / target > 0.05:
        transform["will_crop"] = True
        return False, f"asset aspect {asset_aspect:.2f} vs slot {slot_aspect} ({target:.2f}); would crop", transform
    return True, "", transform


# --- Slot lookup helpers ---------------------------------------------------


_CONTENT_KEY_TO_KIND = {
    "title": "heading",
    "heading": "heading",
    "subtitle": "heading",
    "paragraph": "paragraph",
    "bullets": "bullets",
    "image": "image",
    "hero": "image",
    "table": "table",
    "chart": "chart",
    "footer": "footer",
}

# v5.1 — role preference in slot mapping. When the agent uses a content
# key that names a role, prefer slots with matching role over first-of-kind.
# Flip _V5_ENABLE_ROLE_MATCHING to False to disable purely on the reader
# side without touching the ingest output.
_V5_ENABLE_ROLE_MATCHING = True

_CONTENT_KEY_TO_ROLE = {
    "title": "page_title",
    "page_title": "page_title",
    "subtitle": "subtitle",
    "body": "body",
    "footer": "footer",
    "footnote": "footnote",
    "caption": "caption",
    "key_points": "key_points",
    "detailed_list": "detailed_list",
    "cta": "cta",
    "section_header": "section_header",
}


def _v5_first_slot_of_kind(skeleton: dict, kind: str, exclude_ids: set[str]) -> dict | None:
    for s in skeleton.get("slots") or []:
        if s.get("kind") == kind and s.get("id") not in exclude_ids:
            return s
    return None


def _v5_first_slot_by_role(skeleton: dict, role: str, exclude_ids: set[str]) -> dict | None:
    for s in skeleton.get("slots") or []:
        if s.get("role") == role and s.get("id") not in exclude_ids:
            return s
    return None


def _v5_build_slot_mapping(content: dict, skeleton: dict) -> tuple[dict, list[str]]:
    """Map content keys to slot ids. Returns (mapping, unmapped_keys).
    Unmapped keys = content the skeleton has no slot for.

    Preference order:
    1. Role match (when _V5_ENABLE_ROLE_MATCHING and the content key
       names a role and the skeleton has a slot with that role)
    2. First slot of matching kind (legacy behaviour)
    """
    mapping: dict = {}
    used_ids: set[str] = set()
    unmapped: list[str] = []
    for key in content:
        target_kind = _CONTENT_KEY_TO_KIND.get(key)
        target_role = _CONTENT_KEY_TO_ROLE.get(key) if _V5_ENABLE_ROLE_MATCHING else None
        slot = None
        if target_role is not None:
            slot = _v5_first_slot_by_role(skeleton, target_role, used_ids)
        if slot is None and target_kind is not None:
            slot = _v5_first_slot_of_kind(skeleton, target_kind, used_ids)
        if slot is None:
            unmapped.append(key)
            continue
        mapping[key] = slot["id"]
        used_ids.add(slot["id"])
    return mapping, unmapped


def _v5_check_slot_fit(content_value: Any, slot: dict) -> tuple[bool, str, Any]:
    """Validates a content value against a slot. Returns (fits, reason,
    headroom). Headroom shape depends on kind.
    """
    constraints = slot.get("constraints") or {}
    kind = slot.get("kind")
    if kind in ("heading", "paragraph", "footer"):
        if isinstance(content_value, dict) and "value" in content_value:
            content_value = content_value["value"]
        return _v5_check_text_fit(str(content_value), constraints)
    if kind == "bullets":
        items = content_value if isinstance(content_value, list) else [content_value]
        return _v5_check_bullets_fit(items, constraints)
    if kind == "table":
        if isinstance(content_value, dict):
            return _v5_check_table_fit(content_value, constraints)
        return False, "table value must be a dict {rows, cols, has_header}", {}
    if kind == "chart":
        # Light validation: type whitelist + series/categories counts.
        if not isinstance(content_value, dict):
            return False, "chart value must be a dict {series, categories, type}", {}
        type_key = str(content_value.get("type") or "column").strip().lower()
        if type_key not in _V5_CHART_TYPE_MAP:
            return False, (f"chart type {type_key!r} not supported "
                           f"(use one of {sorted(_V5_CHART_TYPE_MAP)})"), {}
        n_series = content_value.get("n_series") or len(content_value.get("series", []) or [])
        n_cats = content_value.get("n_categories") or len(content_value.get("categories", []) or [])
        max_series = constraints.get("max_series", 99)
        max_cats = constraints.get("max_categories", 99)
        if n_series > max_series:
            return False, f"{n_series} series > max_series {max_series}", {}
        if n_cats > max_cats:
            return False, f"{n_cats} categories > max_categories {max_cats}", {}
        return True, "", {"series": max_series - n_series, "categories": max_cats - n_cats}
    if kind == "image":
        # Image content is typically just an asset_id string; full fit
        # check uses _v5_check_image_fit with the asset meta loaded.
        # Here we only validate the value's shape — content fit happens
        # in cmd_v5_check_asset_fit.
        if isinstance(content_value, str) and content_value.startswith("asset_"):
            return True, "", {}
        if isinstance(content_value, dict) and "asset" in content_value:
            return True, "", {}
        return False, "image value must be 'asset_<id>' or {asset: ...}", {}
    return True, f"unknown kind {kind!r} — passing through", {}


def _v5_required_slots_filled(skeleton: dict, mapping: dict) -> list[str]:
    """Return ids of required slots that are NOT in the mapping (i.e.
    would be left empty by this content)."""
    missing = []
    mapped_ids = set(mapping.values())
    for s in skeleton.get("slots") or []:
        if (s.get("constraints") or {}).get("required") and s.get("id") not in mapped_ids:
            missing.append(s.get("id"))
    return missing


def _v5_headroom_summary(content_value: Any, slot: dict) -> str:
    """Human-friendly headroom string per kind. Used in match-skeletons."""
    kind = slot.get("kind")
    c = slot.get("constraints") or {}
    if kind in ("heading", "paragraph", "footer"):
        v = content_value["value"] if isinstance(content_value, dict) and "value" in content_value else content_value
        if c.get("max_chars"):
            return f"{c['max_chars'] - len(str(v))} chars to spare"
    if kind == "bullets":
        items = content_value if isinstance(content_value, list) else [content_value]
        if c.get("max_items"):
            return f"{c['max_items'] - len(items)} items to spare"
    if kind == "table" and isinstance(content_value, dict):
        if c.get("max_rows"):
            return f"{c['max_rows'] - content_value.get('rows', 0)} rows to spare"
    return ""


# --- CLI command implementations -------------------------------------------


def cmd_v5_list_themes(args: argparse.Namespace) -> None:
    out = []
    root = _v5_themes_dir()
    if root.exists():
        for d in sorted(root.iterdir()):
            if not d.is_dir():
                continue
            t = _v5_load_yaml(d / "theme.yaml")
            if not t:
                continue
            out.append({
                "id": t.get("id"),
                "palette": t.get("palette", {}),
                "fonts": t.get("fonts", {}),
                "decoration_count": len(t.get("decorations") or []),
                "preview_path": str(d / "preview.png") if (d / "preview.png").exists() else None,
            })
    json.dump({"themes": out}, sys.stdout, indent=2, ensure_ascii=False)
    sys.stdout.write("\n")


def cmd_v5_list_skeletons(args: argparse.Namespace) -> None:
    cats = set(args.category) if args.category else None
    has_slot = set(args.has_slot) if args.has_slot else None
    statuses = set(args.status) if args.status else {"pending", "done"}

    out = []
    for sk in _v5_iter_skeletons():
        if sk.get("status", "pending") not in statuses:
            continue
        sk_cats = set(sk.get("categories") or [])
        if cats and not (sk_cats & cats):
            continue
        sk_kinds = {s.get("kind") for s in (sk.get("slots") or [])}
        if has_slot and not (sk_kinds & has_slot):
            continue
        sk_dir = _v5_skeletons_dir() / sk.get("id", "")
        out.append({
            "id": sk.get("id"),
            "source_deck": sk.get("source_deck"),
            "source_slide_index": sk.get("source_slide_index"),
            "status": sk.get("status", "pending"),
            "categories": sk.get("categories") or [],
            "slot_count": len(sk.get("slots") or []),
            "slot_kinds": sorted(sk_kinds),
            "preview_path": str(sk_dir / "preview.png") if (sk_dir / "preview.png").exists() else None,
        })
    json.dump({"skeletons": out}, sys.stdout, indent=2, ensure_ascii=False)
    sys.stdout.write("\n")


def cmd_v5_get_skeleton(args: argparse.Namespace) -> None:
    sk = _v5_load_skeleton(args.id)
    if sk is None:
        raise SystemExit(f"skeleton not found: {args.id}")
    json.dump(sk, sys.stdout, indent=2, ensure_ascii=False)
    sys.stdout.write("\n")


def cmd_v5_get_theme(args: argparse.Namespace) -> None:
    t = _v5_load_theme(args.id)
    if t is None:
        raise SystemExit(f"theme not found: {args.id}")
    json.dump(t, sys.stdout, indent=2, ensure_ascii=False)
    sys.stdout.write("\n")


def cmd_v5_match_skeletons(args: argparse.Namespace) -> None:
    try:
        content = json.loads(args.content)
    except json.JSONDecodeError as e:
        raise SystemExit(f"--content must be valid JSON: {e}")
    if not isinstance(content, dict):
        raise SystemExit("--content must be a JSON object")

    filter_cats = set(args.category) if args.category else None
    filter_has_slot = set(args.has_slot) if args.has_slot else None

    candidates = []
    # Track tightest constraint per content key across all candidates
    # — used to drive the rephrase suggestion on zero-match.
    tightest_per_key: dict[str, dict] = {}

    for sk in _v5_iter_skeletons():
        if sk.get("status") == "rejected":
            continue
        sk_cats = set(sk.get("categories") or [])
        if filter_cats and not (sk_cats & filter_cats):
            continue
        sk_kinds = {s.get("kind") for s in (sk.get("slots") or [])}
        if filter_has_slot and not (sk_kinds & filter_has_slot):
            continue

        mapping, unmapped_keys = _v5_build_slot_mapping(content, sk)
        if unmapped_keys:
            # Content key has no matching slot kind in this skeleton.
            for key in unmapped_keys:
                kind = _CONTENT_KEY_TO_KIND.get(key, key)
                cur = tightest_per_key.get(key, {})
                if "no_such_slot" not in cur:
                    cur["no_such_slot"] = True
                    cur["suggested_action"] = f"no skeleton offers a '{kind}' slot for content key '{key}'"
                    tightest_per_key[key] = cur
            continue

        # Required-slot gate
        missing = _v5_required_slots_filled(sk, mapping)
        if missing:
            continue

        # Fit each content piece
        all_fit = True
        slot_headroom: dict = {}
        slots_by_id = {s["id"]: s for s in (sk.get("slots") or [])}
        for key, slot_id in mapping.items():
            slot = slots_by_id[slot_id]
            fits, reason, headroom = _v5_check_slot_fit(content[key], slot)
            if not fits:
                all_fit = False
                # Track the tightest version of this constraint
                cur = tightest_per_key.get(key, {})
                c = slot.get("constraints") or {}
                if slot.get("kind") in ("heading", "paragraph", "footer"):
                    your_len = len(str(content[key] if not isinstance(content[key], dict) else content[key].get("value", "")))
                    constraint = c.get("max_chars", 0)
                    if "tightest_constraint" not in cur or constraint < cur["tightest_constraint"]:
                        cur.update({
                            "slot": key, "your_value": str(content[key])[:80],
                            "your_length": your_len, "tightest_constraint": constraint,
                            "suggested_action": f"rephrase to ≤{constraint} chars (drop {your_len - constraint})",
                        })
                        tightest_per_key[key] = cur
                elif slot.get("kind") == "bullets":
                    items = content[key] if isinstance(content[key], list) else [content[key]]
                    n = len(items)
                    constraint = c.get("max_items", 0)
                    if "tightest_constraint" not in cur or constraint < cur["tightest_constraint"]:
                        cur.update({
                            "slot": key, "your_count": n, "tightest_constraint": constraint,
                            "suggested_action": f"consolidate to ≤{constraint} items",
                        })
                        tightest_per_key[key] = cur
                break  # one issue per skeleton is enough
            slot_headroom[slot_id] = _v5_headroom_summary(content[key], slot)

        if not all_fit:
            continue

        # Compute fit_score
        tightness_scores = []
        for key, slot_id in mapping.items():
            slot = slots_by_id[slot_id]
            c = slot.get("constraints") or {}
            value = content[key]
            if slot.get("kind") in ("heading", "paragraph", "footer"):
                v = value["value"] if isinstance(value, dict) and "value" in value else value
                m = c.get("max_chars") or len(str(v))
                if m > 0:
                    tightness_scores.append(min(1.0, len(str(v)) / m))
            elif slot.get("kind") == "bullets":
                items = value if isinstance(value, list) else [value]
                m = c.get("max_items") or len(items)
                if m > 0:
                    tightness_scores.append(min(1.0, len(items) / m))
        tightness = sum(tightness_scores) / len(tightness_scores) if tightness_scores else 0.5

        cat_bonus = 0.10 if filter_cats and (sk_cats & filter_cats) else 0
        extra_slots = max(0, len(sk.get("slots") or []) - len(mapping))
        extra_bonus = min(0.20, 0.05 * extra_slots)

        fit_score = min(1.0, tightness * 0.70 + cat_bonus + extra_bonus)

        candidates.append({
            "skeleton_id": sk["id"],
            "categories": list(sk_cats),
            "fit_score": round(fit_score, 3),
            "slot_mapping": mapping,
            "headroom": slot_headroom,
        })

    candidates.sort(key=lambda c: c["fit_score"], reverse=True)

    if candidates:
        result = {"matches": candidates, "issues": []}
    else:
        result = {"matches": [], "issues": list(tightest_per_key.values())}

    json.dump(result, sys.stdout, indent=2, ensure_ascii=False)
    sys.stdout.write("\n")


def cmd_v5_validate_plan(args: argparse.Namespace) -> None:
    plan_path = Path(args.plan)
    if not plan_path.exists():
        raise SystemExit(f"plan not found: {plan_path}")
    plan = json.loads(plan_path.read_text(encoding="utf-8"))
    if not isinstance(plan, list):
        raise SystemExit("plan must be a JSON array")

    errors: list[dict] = []
    warnings: list[dict] = []

    for i, entry in enumerate(plan):
        sk_id = entry.get("skeleton_id")
        sk = _v5_load_skeleton(sk_id) if sk_id else None
        if sk is None:
            errors.append({
                "slide_index": i, "slot_id": None,
                "violation": "skeleton_not_found",
                "message": f"no skeleton with id {sk_id!r}",
            })
            continue
        slots_by_id = {s["id"]: s for s in (sk.get("slots") or [])}
        filled = entry.get("slots") or {}

        # Required slots
        for slot in sk.get("slots") or []:
            if (slot.get("constraints") or {}).get("required") and slot["id"] not in filled:
                errors.append({
                    "slide_index": i, "slot_id": slot["id"],
                    "violation": "required_unfilled",
                    "message": f"required slot {slot['id']!r} not in plan",
                })

        # Constraint checks on filled slots
        for slot_id, value in filled.items():
            slot = slots_by_id.get(slot_id)
            if slot is None:
                errors.append({
                    "slide_index": i, "slot_id": slot_id,
                    "violation": "unknown_slot",
                    "message": f"slot {slot_id!r} not in skeleton {sk_id!r}",
                })
                continue
            is_overflow_shrink = isinstance(value, dict) and value.get("overflow") == "shrink"
            inner = value.get("value", value) if isinstance(value, dict) else value
            fits, reason, _ = _v5_check_slot_fit(inner, slot)
            if not fits:
                if is_overflow_shrink:
                    warnings.append({
                        "slide_index": i, "slot_id": slot_id,
                        "overflow_kind": "shrink", "message": reason,
                    })
                else:
                    errors.append({
                        "slide_index": i, "slot_id": slot_id,
                        "violation": "constraint", "message": reason,
                    })

    result = {"ok": len(errors) == 0, "errors": errors, "warnings": warnings}
    json.dump(result, sys.stdout, indent=2, ensure_ascii=False)
    sys.stdout.write("\n")


def cmd_v5_check_asset_fit(args: argparse.Namespace) -> None:
    sk = _v5_load_skeleton(args.skeleton_id)
    if sk is None:
        raise SystemExit(f"skeleton not found: {args.skeleton_id}")
    slot = next((s for s in (sk.get("slots") or []) if s.get("id") == args.slot_id), None)
    if slot is None:
        raise SystemExit(f"slot {args.slot_id!r} not in {args.skeleton_id}")
    if slot.get("kind") != "image":
        json.dump({
            "fits": False, "will_resize_to": None, "will_crop": False,
            "reason": f"slot is kind={slot.get('kind')!r}, not image",
            "suggestion": "pick an image slot or change the slot kind",
        }, sys.stdout, indent=2, ensure_ascii=False)
        sys.stdout.write("\n")
        return
    meta = _v5_load_asset_meta(args.asset_id) or {}
    fits, reason, transform = _v5_check_image_fit(meta, slot.get("constraints") or {})
    result = {
        "fits": fits,
        "will_resize_to": None,
        "will_crop": transform.get("will_crop", False),
        "reason": reason or None,
        "suggestion": ("would crop to slot aspect" if transform.get("will_crop") else None) if not fits else None,
    }
    json.dump(result, sys.stdout, indent=2, ensure_ascii=False)
    sys.stdout.write("\n")


def cmd_v5_find_asset(args: argparse.Namespace) -> None:
    """Deterministic shortlist filter over the asset library.

    Filter dimensions:
      --kind   required; matches the asset's ``kind`` exactly
      --tags   optional, repeatable; an asset matches if it carries
               every requested tag (AND, not OR)

    `description` is included in each match for picking the final 1-of-N
    by topic fit, but is never a filter input. `width`/`height`/`aspect`
    ride along so the agent can call check-asset-fit on the shortlist
    without a second round-trip.

    Two runs against the same library + same query produce the same
    shortlist in the same order (sorted by id, ascending). Idempotent
    by construction.

    If the filter is too tight to return any candidates, ``suggestion``
    points to the broadening step.
    """
    index = load_index()
    assets = index.get("assets", []) or []

    pool = [a for a in assets if str(a.get("kind", "")) == args.kind]
    total_kind = len(pool)

    wanted_tags = list(args.tags or [])
    if wanted_tags:
        pool = [
            a for a in pool
            if all(t in (a.get("tags") or []) for t in wanted_tags)
        ]

    pool.sort(key=lambda a: str(a.get("id", "")))
    limit = max(1, int(args.limit or 5))
    shortlist = pool[:limit]

    suggestion: str | None = None
    if not shortlist:
        if wanted_tags:
            suggestion = (
                f"no match — drop --tags and retry with --kind {args.kind} "
                f"only. If still empty, either stage a new asset "
                f'(POST /api/asset/add) or pass "placeholder" as the '
                f"asset_id to render a labeled grey box for manual "
                f"replacement post-build."
            )
        else:
            suggestion = (
                f"no assets of kind={args.kind!r} in the library. "
                f'Stage one via /api/asset/add or use "placeholder".'
            )

    out = {
        "query": {
            "kind": args.kind,
            "tags": wanted_tags,
        },
        "matches": [
            {
                "id": a.get("id"),
                "kind": a.get("kind"),
                "tags": list(a.get("tags") or []),
                "description": a.get("description", ""),
                "width": int(a.get("width") or 0),
                "height": int(a.get("height") or 0),
                "aspect": float(a.get("aspect") or 0.0),
                "colors_hex": list(a.get("colors_hex") or []),
            }
            for a in shortlist
        ],
        "count": len(shortlist),
        "total_of_kind": total_kind,
        "tag_vocab": list(index.get("tag_vocab") or []),
        "suggestion": suggestion,
    }
    json.dump(out, sys.stdout, indent=2, ensure_ascii=False)
    sys.stdout.write("\n")


def cmd_v5_compose(args: argparse.Namespace) -> None:
    """Build a deck from a v5 plan on a chosen host theme's master.

    Opens themes/<theme_id>/master.pptx as host, strips its slides,
    then for each plan entry creates a new blank slide and places
    primitives per the skeleton's slot inventory.

    Plan shape (same as validate-plan):
      [{"skeleton_id": "...", "slots": {"slot_id": value, ...}}, ...]

    Slot value shapes:
      string → text content
      list[string] → bullets
      dict {value, overflow: "shrink"} → text with autofit
      dict {rows, cols, has_header, data: [[...]]} → table
      dict {type, series, categories} → chart (deferred — emits warning)
      "asset_<id>" or dict {asset: "..."} → image

    Writes a JSON result to stdout with output path + warnings,
    plus a <out>.warnings.json sidecar with the same warnings for
    the user to triage overflow:shrink events after the deck opens.
    """
    from pptx import Presentation
    from pptx.util import Emu, Pt

    plan_path = Path(args.plan)
    if not plan_path.exists():
        raise SystemExit(f"plan not found: {plan_path}")
    plan = json.loads(plan_path.read_text(encoding="utf-8"))
    if not isinstance(plan, list) or not plan:
        raise SystemExit("plan must be a non-empty JSON array")

    theme = _v5_load_theme(args.theme)
    if theme is None:
        raise SystemExit(f"theme not found: {args.theme}")
    master_path = _v5_themes_dir() / args.theme / theme.get("master_pptx", "master.pptx")
    if not master_path.exists():
        raise SystemExit(f"theme master.pptx missing: {master_path}")

    out_path = Path(args.out)
    warnings: list[dict] = []

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        shutil.copyfile(master_path, tmp.name)
        host_path = Path(tmp.name)

    try:
        prs = Presentation(str(host_path))
        slide_w = prs.slide_width or 9144000
        slide_h = prs.slide_height or 6858000

        # Strip the source's existing slides — we build fresh.
        _v5_drop_all_slides(prs)

        # Pick a blank-ish layout to host each new slide.
        blank_layout = _v5_pick_blank_layout(prs)

        for i, entry in enumerate(plan):
            sk_id = entry.get("skeleton_id")
            sk = _v5_load_skeleton(sk_id) if sk_id else None
            if sk is None:
                warnings.append({"slide_index": i, "violation": "skeleton_not_found",
                                 "message": f"no skeleton {sk_id!r}"})
                continue

            slide = prs.slides.add_slide(blank_layout)
            slots_by_id = {s["id"]: s for s in (sk.get("slots") or [])}
            filled = entry.get("slots") or {}

            # Apply background_image (B4-render) if the skeleton has
            # one. Paint full-bleed FIRST so subsequent slot shapes
            # stack on top via python-pptx's natural z-order. Fail-soft:
            # missing file or any add_picture error → warn and proceed
            # without the underlay (deck still renders, just without
            # the structural illustration baked in).
            bg_rel = sk.get("background_image")
            if bg_rel:
                bg_path = _v5_skeletons_dir() / sk_id / bg_rel
                if not bg_path.exists():
                    warnings.append({
                        "slide_index": i, "slot_id": "_background",
                        "violation": "background_missing",
                        "message": f"background_image set to {bg_rel!r} "
                                   f"but file not in bundle; slide built without underlay",
                    })
                else:
                    try:
                        slide.shapes.add_picture(
                            str(bg_path), 0, 0,
                            width=slide_w, height=slide_h,
                        )
                    except Exception as e:
                        warnings.append({
                            "slide_index": i, "slot_id": "_background",
                            "violation": "background_place_failed",
                            "message": f"{type(e).__name__}: {e}; "
                                       f"slide built without underlay",
                        })

            for slot_id, value in filled.items():
                slot = slots_by_id.get(slot_id)
                if slot is None:
                    warnings.append({
                        "slide_index": i, "slot_id": slot_id,
                        "violation": "unknown_slot",
                        "message": f"slot {slot_id!r} not in skeleton {sk_id!r}",
                    })
                    continue
                ws = _v5_place_slot(slide, slot, value, slide_w, slide_h, theme)
                for w in ws:
                    w.setdefault("slide_index", i)
                    w.setdefault("slot_id", slot_id)
                warnings.extend(ws)

        out_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(out_path))

        # Sidecar for human review of overflow/warning events.
        if warnings:
            sidecar = out_path.with_suffix(out_path.suffix + ".warnings.json")
            sidecar.write_text(
                json.dumps({"warnings": warnings}, indent=2, ensure_ascii=False),
                encoding="utf-8",
            )
    finally:
        try:
            host_path.unlink()
        except OSError:
            pass

    result = {
        "output": str(out_path),
        "slides": len(plan),
        "warnings": warnings,
        "warnings_sidecar": str(out_path.with_suffix(out_path.suffix + ".warnings.json")) if warnings else None,
    }
    json.dump(result, sys.stdout, indent=2, ensure_ascii=False)
    sys.stdout.write("\n")


def _v5_drop_all_slides(prs) -> None:
    """Remove every slide from the host master so we can add fresh
    ones for each plan entry. Mirrors the pattern in v4's
    _drop_first_slide but applied repeatedly.
    """
    sldIdLst = prs.slides._sldIdLst
    rId_to_drop = []
    for sldId in list(sldIdLst):
        rId_to_drop.append(sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"))
        sldIdLst.remove(sldId)
    for rId in rId_to_drop:
        if rId:
            try:
                prs.part.drop_rel(rId)
            except KeyError:
                pass


def _v5_pick_blank_layout(prs):
    """Pick the simplest available layout (fewest placeholders).
    Mirrors v4's "fewest placeholders" heuristic.
    """
    layouts = list(prs.slide_layouts)
    if not layouts:
        raise SystemExit("host master has no slide layouts")
    return min(layouts, key=lambda lo: len(list(lo.placeholders)))


def _v5_emu_geometry(slot: dict, slide_w: int, slide_h: int) -> tuple[int, int, int, int]:
    """Fractional → EMU. Returns (left, top, width, height) in EMU."""
    g = slot.get("geometry") or {}
    return (
        int((g.get("x", 0)) * slide_w),
        int((g.get("y", 0)) * slide_h),
        int((g.get("w", 0)) * slide_w),
        int((g.get("h", 0)) * slide_h),
    )


def _v5_resolve_font_name(slot: dict, theme: dict) -> str | None:
    """Resolve a slot's style.font_role against the host theme."""
    style = slot.get("style") or {}
    role = style.get("font_role")
    fonts = theme.get("fonts") or {}
    if role == "major":
        return fonts.get("major")
    if role == "minor":
        return fonts.get("minor")
    if role == "explicit":
        return style.get("typeface")
    return None


def _v5_resolve_color(slot: dict, theme: dict):
    """Resolve a slot's style.color_role or .color → RGBColor or None."""
    from pptx.dml.color import RGBColor
    style = slot.get("style") or {}
    role = style.get("color_role")
    palette = theme.get("palette") or {}
    hex_val = None
    if role and role in palette:
        hex_val = palette[role]
    elif style.get("color"):
        hex_val = style["color"]
    if not hex_val:
        return None
    try:
        return RGBColor.from_string(hex_val.lstrip("#"))
    except Exception:
        return None


def _v5_place_slot(slide, slot: dict, value, slide_w: int, slide_h: int, theme: dict) -> list[dict]:
    """Dispatch to per-kind placers."""
    kind = slot.get("kind")
    warnings: list[dict] = []
    overflow = None
    if isinstance(value, dict) and "overflow" in value:
        overflow = value.get("overflow")
        value = value.get("value", value)
    try:
        if kind in ("heading", "paragraph", "footer"):
            warnings.extend(_v5_place_text(slide, slot, value, slide_w, slide_h, theme, overflow))
        elif kind == "bullets":
            items = value if isinstance(value, list) else [value]
            warnings.extend(_v5_place_bullets(slide, slot, items, slide_w, slide_h, theme, overflow))
        elif kind == "image":
            warnings.extend(_v5_place_image(slide, slot, value, slide_w, slide_h, theme))
        elif kind == "table":
            warnings.extend(_v5_place_table(slide, slot, value, slide_w, slide_h, theme))
        elif kind == "chart":
            warnings.extend(_v5_place_chart(slide, slot, value, slide_w, slide_h, theme))
        else:
            warnings.append({"violation": "unknown_kind",
                             "message": f"unknown slot kind {kind!r}"})
    except Exception as e:
        warnings.append({"violation": "place_failed",
                         "message": f"{type(e).__name__}: {e}"})
    return warnings


def _v5_place_text(slide, slot: dict, value, slide_w, slide_h, theme, overflow) -> list[dict]:
    from pptx.util import Pt
    from pptx.enum.text import MSO_ANCHOR
    warnings: list[dict] = []
    left, top, w, h = _v5_emu_geometry(slot, slide_w, slide_h)
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    if overflow == "shrink":
        # MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE doesn't exist universally;
        # use word_wrap + autofit via tf.auto_size if available.
        try:
            from pptx.enum.text import MSO_AUTO_SIZE
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            warnings.append({"overflow_kind": "shrink",
                             "message": "text autofit enabled per overflow:shrink"})
        except Exception:
            pass

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = str(value)
    style = slot.get("style") or {}
    if style.get("size_pt"):
        run.font.size = Pt(style["size_pt"])
    if style.get("bold") is not None:
        run.font.bold = style["bold"]
    if style.get("italic") is not None:
        run.font.italic = style["italic"]
    font_name = _v5_resolve_font_name(slot, theme)
    if font_name:
        run.font.name = font_name
    color = _v5_resolve_color(slot, theme)
    if color is not None:
        run.font.color.rgb = color
    return warnings


def _v5_place_bullets(slide, slot: dict, items: list, slide_w, slide_h, theme, overflow) -> list[dict]:
    from pptx.util import Pt
    warnings: list[dict] = []
    left, top, w, h = _v5_emu_geometry(slot, slide_w, slide_h)
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    if overflow == "shrink":
        try:
            from pptx.enum.text import MSO_AUTO_SIZE
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            warnings.append({"overflow_kind": "shrink",
                             "message": "bullet autofit enabled per overflow:shrink"})
        except Exception:
            pass

    style = slot.get("style") or {}
    font_name = _v5_resolve_font_name(slot, theme)
    color = _v5_resolve_color(slot, theme)

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"• {item}"
        if style.get("size_pt"):
            run.font.size = Pt(style["size_pt"])
        if style.get("bold") is not None:
            run.font.bold = style["bold"]
        if font_name:
            run.font.name = font_name
        if color is not None:
            run.font.color.rgb = color
    return warnings


def _v5_draw_placeholder_box(slide, left: int, top: int, w: int, h: int, label: str) -> None:
    """Draw a labeled grey rectangle in place of a missing image asset.

    Used by the "placeholder" sentinel asset_id — agents emit this when
    find-asset returns empty for a required slot. The box is dashed +
    light grey so it reads as "fix me" in the final deck.
    """
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    from pptx.util import Pt
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
    line = shape.line
    line.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
    try:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        line.dash_style = MSO_LINE_DASH_STYLE.DASH
    except Exception:
        pass
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = label
    run.font.size = Pt(12)
    run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    run.font.italic = True


def _v5_place_image(slide, slot: dict, value, slide_w, slide_h, theme) -> list[dict]:
    warnings: list[dict] = []
    left, top, w, h = _v5_emu_geometry(slot, slide_w, slide_h)
    # Resolve asset id from value
    asset_id = None
    placeholder_label: str | None = None
    if isinstance(value, str):
        asset_id = value
    elif isinstance(value, dict):
        if value.get("placeholder"):
            placeholder_label = str(value.get("label") or "")
            asset_id = "placeholder"
        else:
            asset_id = value.get("asset") or value.get("asset_id")
    # Placeholder sentinel: emit a labeled grey box instead of a binary.
    # The agent uses this when find-asset returns empty for a required
    # slot and no online sourcing was possible. Build emits a warning so
    # the user sees in the sidecar which slots still need a real asset.
    if asset_id in ("placeholder", "asset_placeholder"):
        slot_id = slot.get("id") or "?"
        label = placeholder_label or f"image needed: {slot_id}"
        _v5_draw_placeholder_box(slide, left, top, w, h, label)
        warnings.append({
            "violation": "image_placeholder",
            "slot_id": slot_id,
            "message": f"placeholder rendered for slot {slot_id!r} — replace with a real asset",
        })
        return warnings
    if not asset_id:
        warnings.append({"violation": "no_asset", "message": "image value missing asset id"})
        return warnings
    # Find the asset binary
    assets_dir = _v5_assets_dir()
    if not assets_dir.exists():
        warnings.append({"violation": "no_assets_dir", "message": f"assets dir missing: {assets_dir}"})
        return warnings
    bin_path = None
    for cand in assets_dir.glob(f"{asset_id}.*"):
        if cand.suffix == ".yaml":
            continue
        bin_path = cand
        break
    if bin_path is None:
        warnings.append({"violation": "asset_not_found", "message": f"asset {asset_id} binary missing"})
        return warnings

    # Aspect-aware placement per the slot's auto_fit policy. Default
    # "cover" matches the agent contract — center-crop preserving
    # aspect to fill the slot. "contain" letterboxes; "stretch" is
    # the old distorting behaviour, kept for opt-in compatibility.
    fit = (slot.get("constraints") or {}).get("auto_fit") or "cover"
    asset_w, asset_h = _v5_image_dimensions(bin_path)
    if asset_w <= 0 or asset_h <= 0 or fit == "stretch":
        # Unknown dims or explicit stretch → fall back to direct fit
        # (matches the pre-aspect behaviour; cheaper than refusing).
        slide.shapes.add_picture(str(bin_path), left, top, w, h)
        if fit != "stretch" and (asset_w <= 0 or asset_h <= 0):
            warnings.append({
                "violation": "asset_dims_unknown",
                "message": f"could not read dimensions of {bin_path.name}; placed stretched",
            })
        return warnings

    asset_aspect = asset_w / asset_h
    slot_aspect = w / h if h > 0 else 1.0

    if fit == "contain":
        # Letterbox: shrink to fit inside slot, leave bands.
        if asset_aspect > slot_aspect:
            placed_w = w
            placed_h = int(w / asset_aspect)
        else:
            placed_h = h
            placed_w = int(h * asset_aspect)
        placed_left = left + (w - placed_w) // 2
        placed_top = top + (h - placed_h) // 2
        slide.shapes.add_picture(str(bin_path), placed_left, placed_top, placed_w, placed_h)
        return warnings

    # Default "cover": scale image larger than slot, crop overflow.
    # python-pptx exposes pic.crop_left/right/top/bottom as fractions
    # of the *displayed* image size (i.e. of placed_w / placed_h).
    if asset_aspect > slot_aspect:
        # Asset wider than slot → match height, crop sides
        placed_h = h
        placed_w = int(h * asset_aspect)
        crop_amount = (placed_w - w) / placed_w / 2
        pic = slide.shapes.add_picture(str(bin_path), left - int(placed_w - w) // 2, top, placed_w, placed_h)
        pic.crop_left = crop_amount
        pic.crop_right = crop_amount
    else:
        # Asset taller than slot → match width, crop top/bottom
        placed_w = w
        placed_h = int(w / asset_aspect)
        crop_amount = (placed_h - h) / placed_h / 2
        pic = slide.shapes.add_picture(str(bin_path), left, top - int(placed_h - h) // 2, placed_w, placed_h)
        pic.crop_top = crop_amount
        pic.crop_bottom = crop_amount
    return warnings


def _v5_image_dimensions(path: Path) -> tuple[int, int]:
    """Return (width, height) of a raster image in pixels, or (0, 0)
    on any failure. PIL is the existing dependency for v4 dominant-
    colour extraction, so it's already available.
    """
    try:
        from PIL import Image
        with Image.open(path) as img:
            return img.size
    except Exception:
        return 0, 0


def _v5_place_table(slide, slot: dict, value, slide_w, slide_h, theme) -> list[dict]:
    from pptx.util import Pt
    warnings: list[dict] = []
    left, top, w, h = _v5_emu_geometry(slot, slide_w, slide_h)
    if not isinstance(value, dict):
        warnings.append({"violation": "bad_table", "message": "table value must be dict"})
        return warnings
    data = value.get("data") or []
    rows = value.get("rows") or len(data)
    cols = value.get("cols") or (len(data[0]) if data else 0)
    if rows < 1 or cols < 1:
        warnings.append({"violation": "empty_table", "message": "table needs rows/cols"})
        return warnings
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, w, h)
    tbl = tbl_shape.table
    for r in range(min(rows, len(data))):
        row_data = data[r]
        for c in range(min(cols, len(row_data))):
            cell = tbl.cell(r, c)
            cell.text = str(row_data[c])
    return warnings


# Chart type strings → python-pptx XL_CHART_TYPE attribute names. Kept
# as a small whitelist; unknown types emit a warning and skip rather
# than crash. Scatter charts use XyChartData (different data shape) so
# they're not included here — add when there's a real need.
_V5_CHART_TYPE_MAP = {
    "bar": "BAR_CLUSTERED",
    "bar_clustered": "BAR_CLUSTERED",
    "bar_stacked": "BAR_STACKED",
    "column": "COLUMN_CLUSTERED",
    "column_clustered": "COLUMN_CLUSTERED",
    "column_stacked": "COLUMN_STACKED",
    "line": "LINE",
    "line_markers": "LINE_MARKERS",
    "pie": "PIE",
    "doughnut": "DOUGHNUT",
    "area": "AREA",
    "area_stacked": "AREA_STACKED",
}


def _v5_place_chart(slide, slot: dict, value, slide_w, slide_h, theme) -> list[dict]:
    """Build a category chart from primitives via python-pptx's
    add_chart. Replaces the old "chart_not_implemented" warning.

    Expected value shape:
      {
        "type": "bar|column|line|pie|doughnut|area" (+ variants),
        "categories": ["Q1", "Q2", "Q3"],
        "series": [{"name": "Revenue", "values": [10, 20, 30]}, ...]
      }

    Fail-soft: malformed value, unknown type, or any python-pptx
    exception → append a warning and leave the slot empty. Matches
    the previous fail-soft behaviour so a broken chart spec never
    crashes the whole compose run.
    """
    warnings: list[dict] = []
    left, top, w, h = _v5_emu_geometry(slot, slide_w, slide_h)

    if not isinstance(value, dict):
        warnings.append({
            "violation": "bad_chart",
            "message": "chart value must be a dict {type, categories, series}",
        })
        return warnings

    type_key = str(value.get("type") or "column").strip().lower()
    mapped = _V5_CHART_TYPE_MAP.get(type_key)
    if mapped is None:
        warnings.append({
            "violation": "unsupported_chart_type",
            "message": (f"chart type {type_key!r} not supported "
                        f"(use one of {sorted(_V5_CHART_TYPE_MAP)}); "
                        f"slot left empty"),
        })
        return warnings

    categories = value.get("categories") or []
    series = value.get("series") or []
    if not categories or not series:
        warnings.append({
            "violation": "empty_chart_data",
            "message": "chart needs at least one category and one series; slot left empty",
        })
        return warnings

    try:
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        cd = CategoryChartData()
        cd.categories = [str(c) for c in categories]
        n_cats = len(categories)
        added = 0
        for s in series:
            if not isinstance(s, dict):
                continue
            name = str(s.get("name") or "")
            raw_vals = s.get("values") or []
            vals = []
            for i in range(n_cats):
                v = raw_vals[i] if i < len(raw_vals) else None
                try:
                    vals.append(float(v) if v is not None else 0.0)
                except (TypeError, ValueError):
                    vals.append(0.0)
            cd.add_series(name, vals)
            added += 1
        if added == 0:
            warnings.append({
                "violation": "empty_chart_data",
                "message": "no usable series after filtering; slot left empty",
            })
            return warnings
        chart_type = getattr(XL_CHART_TYPE, mapped)
        slide.shapes.add_chart(chart_type, left, top, w, h, cd)
    except Exception as e:
        warnings.append({
            "violation": "chart_place_failed",
            "message": f"{type(e).__name__}: {e}; slot left empty",
        })
    return warnings


def cmd_v5_measure_text(args: argparse.Namespace) -> None:
    if args.array:
        try:
            items = json.loads(args.array)
        except json.JSONDecodeError as e:
            raise SystemExit(f"--array must be valid JSON: {e}")
        text = "\n".join(str(x) for x in items)
        n_items = len(items)
    else:
        text = args.text or ""
        n_items = 1
    chars = len(text)
    words = len(text.split())
    lines = text.count("\n") + 1
    out: dict = {"chars": chars, "words": words, "lines_est": lines}
    if args.array:
        out["items"] = n_items
    if args.against:
        try:
            sk_id, slot_id = args.against.split(".", 1)
        except ValueError:
            raise SystemExit("--against must be '<skeleton_id>.<slot_id>'")
        sk = _v5_load_skeleton(sk_id)
        if sk is None:
            raise SystemExit(f"skeleton not found: {sk_id}")
        slot = next((s for s in (sk.get("slots") or []) if s.get("id") == slot_id), None)
        if slot is None:
            raise SystemExit(f"slot {slot_id!r} not in {sk_id}")
        c = slot.get("constraints") or {}
        if args.array:
            fits, reason, hr = _v5_check_bullets_fit(items, c)
            out["fits"] = fits
            out["headroom"] = f"{hr['items']} items to spare" if fits else reason
        else:
            fits, reason, hr_chars = _v5_check_text_fit(text, c)
            out["fits"] = fits
            out["headroom"] = (f"{hr_chars} chars to spare" if fits
                               else f"{abs(hr_chars)} chars over (max {c.get('max_chars')})")
    json.dump(out, sys.stdout, indent=2, ensure_ascii=False)
    sys.stdout.write("\n")


# ===========================================================================
# Entry point
# ===========================================================================


def main(argv: list[str] | None = None) -> None:
    parser = argparse.ArgumentParser(prog="reader.py")
    sub = parser.add_subparsers(dest="cmd", required=True)

    p_lt = sub.add_parser("list-themes", help="List v5 themes.")
    p_lt.set_defaults(func=cmd_v5_list_themes)

    p_ls = sub.add_parser("list-skeletons", help="List v5 skeletons (filterable).")
    p_ls.add_argument("--category", action="append", default=None,
                      help="filter by category (repeatable; any-match)")
    p_ls.add_argument("--has-slot", action="append", default=None,
                      help="filter by slot kind present (repeatable; any-match)")
    p_ls.add_argument("--status", action="append", default=None,
                      help="status filter (pending/done/rejected; default excludes rejected)")
    p_ls.set_defaults(func=cmd_v5_list_skeletons)

    p_gs = sub.add_parser("get-skeleton", help="Get one v5 skeleton by id.")
    p_gs.add_argument("id")
    p_gs.set_defaults(func=cmd_v5_get_skeleton)

    p_gt = sub.add_parser("get-theme", help="Get one v5 theme by id.")
    p_gt.add_argument("id")
    p_gt.set_defaults(func=cmd_v5_get_theme)

    p_ms = sub.add_parser("match-skeletons",
                          help="Content-first ranked match. Returns matches or rephrase issues.")
    p_ms.add_argument("--content", required=True, help="JSON content dict")
    p_ms.add_argument("--category", action="append", default=None)
    p_ms.add_argument("--has-slot", action="append", default=None)
    p_ms.set_defaults(func=cmd_v5_match_skeletons)

    p_vp = sub.add_parser("validate-plan",
                          help="Pre-build constraint check on a full plan.")
    p_vp.add_argument("plan")
    p_vp.set_defaults(func=cmd_v5_validate_plan)

    p_cf = sub.add_parser("check-asset-fit",
                          help="Does this asset fit this skeleton slot?")
    p_cf.add_argument("asset_id")
    p_cf.add_argument("skeleton_id")
    p_cf.add_argument("slot_id")
    p_cf.set_defaults(func=cmd_v5_check_asset_fit)

    p_fa = sub.add_parser(
        "find-asset",
        help="Deterministic shortlist filter over the asset library "
             "(--kind required; --tags optional, AND-matched). Always "
             "call this BEFORE picking an asset_id by reading index.json.",
    )
    p_fa.add_argument(
        "--kind", required=True,
        help="photo|icon|logo|illustration|screenshot|vector|table|"
             "chart|callout|freeform|smartart",
    )
    p_fa.add_argument(
        "--tags", action="append", default=None,
        help="repeatable; AND-matched against the workspace tag vocab "
             "(see `tag_vocab` field on the find-asset response or in "
             "index.json).",
    )
    p_fa.add_argument(
        "--limit", type=int, default=5,
        help="cap on shortlist size (default 5)",
    )
    p_fa.set_defaults(func=cmd_v5_find_asset)

    p_mt = sub.add_parser("measure-text",
                          help="Char/word/line counts; optional fit against a slot.")
    p_mt.add_argument("text", nargs="?", default=None)
    p_mt.add_argument("--array", default=None, help="JSON array of items (for bullets)")
    p_mt.add_argument("--against", default=None, help="<skeleton_id>.<slot_id>")
    p_mt.set_defaults(func=cmd_v5_measure_text)

    p_cv = sub.add_parser("compose-v5",
                          help="Build a deck from a v5 plan on a chosen host theme.")
    p_cv.add_argument("plan")
    p_cv.add_argument("out")
    p_cv.add_argument("--theme", required=True, help="theme_id (see list-themes)")
    p_cv.set_defaults(func=cmd_v5_compose)

    args = parser.parse_args(argv)
    args.func(args)


if __name__ == "__main__":
    main()
