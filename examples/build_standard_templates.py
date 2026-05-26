"""Build a 10-slide standard templates deck for pptx-skill.

Each slide: header at top, text/bullets on the left, chart/image/table on
the right. Uses real python-pptx chart/table/picture objects so the
ingest pipeline detects the slot kinds correctly.
"""
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageDraw, ImageFont

# ---- Brand --------------------------------------------------------------
ACCENT = RGBColor(0x1F, 0x4E, 0x8C)        # deep blue
INK = RGBColor(0x22, 0x2A, 0x35)           # near-black
MUTED = RGBColor(0x6B, 0x73, 0x82)         # slate gray
LIGHT = RGBColor(0xF1, 0xF4, 0xF8)         # panel background
HAIR = RGBColor(0xD8, 0xDD, 0xE5)

TITLE_FONT = "Calibri"
BODY_FONT = "Calibri"

# ---- Helpers ------------------------------------------------------------
def make_placeholder_image(w_px: int, h_px: int, label: str) -> BytesIO:
    """Generate a soft-grey rectangle PNG with a centered label so the
    slide reads as 'image goes here' even before the agent fills it."""
    img = Image.new("RGB", (w_px, h_px), (235, 238, 243))
    d = ImageDraw.Draw(img)
    # subtle border
    d.rectangle([0, 0, w_px - 1, h_px - 1], outline=(200, 206, 216), width=2)
    # centered label
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 28)
    except Exception:
        font = ImageFont.load_default()
    bbox = d.textbbox((0, 0), label, font=font)
    tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
    d.text(((w_px - tw) / 2, (h_px - th) / 2 - 10), label,
           fill=(140, 148, 160), font=font)
    buf = BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def add_title(slide, text: str):
    """Header band across the top: bold title + thin accent rule."""
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.33), Inches(0.7))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = TITLE_FONT
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = INK
    # Accent rule under title
    from pptx.enum.shapes import MSO_SHAPE
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.5), Inches(1.08), Inches(0.7), Emu(28575))
    rule.fill.solid(); rule.fill.fore_color.rgb = ACCENT
    rule.line.fill.background()


def add_footer(slide, text: str):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.33), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.name = BODY_FONT; r.font.size = Pt(9); r.font.color.rgb = MUTED


def add_left_bullets(slide, items: list[str], left=Inches(0.5), top=Inches(1.45),
                     width=Inches(5.8), height=Inches(5.2)):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        r = p.add_run()
        r.text = "•  " + item
        r.font.name = BODY_FONT; r.font.size = Pt(16); r.font.color.rgb = INK


def add_left_paragraph(slide, text: str, left=Inches(0.5), top=Inches(1.45),
                       width=Inches(5.8), height=Inches(5.2)):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.name = BODY_FONT; r.font.size = Pt(15); r.font.color.rgb = INK


def add_left_lead_plus_bullets(slide, lead: str, bullets: list[str]):
    """A lead sentence (larger) followed by sub-bullets — heavier left col."""
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.45), Inches(5.8), Inches(5.2))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = lead
    r.font.name = BODY_FONT; r.font.size = Pt(18); r.font.bold = True
    r.font.color.rgb = INK
    p.space_after = Pt(14)
    for item in bullets:
        sp = tf.add_paragraph(); sp.space_after = Pt(6)
        sr = sp.add_run(); sr.text = "•  " + item
        sr.font.name = BODY_FONT; sr.font.size = Pt(14); sr.font.color.rgb = INK


def add_left_kpis(slide, stats: list[tuple[str, str]]):
    """Two big KPI stats (number + label)."""
    base_top = Inches(1.7)
    for i, (value, label) in enumerate(stats):
        block_top = base_top + Inches(2.1 * i)
        # value
        tb = slide.shapes.add_textbox(Inches(0.5), block_top, Inches(5.8), Inches(1.1))
        p = tb.text_frame.paragraphs[0]
        r = p.add_run(); r.text = value
        r.font.name = TITLE_FONT; r.font.size = Pt(56); r.font.bold = True
        r.font.color.rgb = ACCENT
        # label
        lb = slide.shapes.add_textbox(Inches(0.5), block_top + Inches(1.1),
                                      Inches(5.8), Inches(0.5))
        lp = lb.text_frame.paragraphs[0]
        lr = lp.add_run(); lr.text = label
        lr.font.name = BODY_FONT; lr.font.size = Pt(14); lr.font.color.rgb = MUTED


# Right-column geometry (consistent across all 10 slides)
R_LEFT, R_TOP, R_W, R_H = Inches(6.7), Inches(1.45), Inches(6.13), Inches(5.2)


def add_right_chart(slide, chart_type, categories, series, title=None):
    data = CategoryChartData()
    data.categories = categories
    for name, vals in series:
        data.add_series(name, vals)
    chart_shape = slide.shapes.add_chart(chart_type, R_LEFT, R_TOP, R_W, R_H, data)
    chart = chart_shape.chart
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
        for r in chart.chart_title.text_frame.paragraphs[0].runs:
            r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = INK
    else:
        chart.has_title = False
    if len(series) > 1:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    else:
        chart.has_legend = False


def add_right_image(slide, label: str = "Image placeholder"):
    # Use slot pixel size at ~96dpi for crisp placeholder text
    w_px, h_px = 760, 640
    buf = make_placeholder_image(w_px, h_px, label)
    slide.shapes.add_picture(buf, R_LEFT, R_TOP, R_W, R_H)


def add_right_table(slide, headers: list[str], rows: list[list[str]]):
    cols = len(headers); nrows = len(rows) + 1
    tbl_shape = slide.shapes.add_table(nrows, cols, R_LEFT, R_TOP, R_W, R_H)
    tbl = tbl_shape.table
    # Header row
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True; r.font.size = Pt(12); r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
    # Body rows
    for r_idx, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = tbl.cell(r_idx, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(12); run.font.color.rgb = INK


# ---- Build deck ---------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# Slide 1: Bullets + Bar chart
s = prs.slides.add_slide(blank)
add_title(s, "Quarterly Revenue Highlights")
add_left_bullets(s, [
    "Strong rebound vs. prior period",
    "EMEA leads regional growth",
    "Margin recovered to plan",
    "New logo motion fully ramped",
    "Pipeline coverage healthy into Q4",
])
add_right_chart(s, XL_CHART_TYPE.BAR_CLUSTERED,
                ["Americas", "EMEA", "APAC"],
                [("Revenue ($M)", (38, 52, 27))],
                title="Revenue by region")
add_footer(s, "Template 1 · bullets + bar chart")

# Slide 2: Bullets + Line chart
s = prs.slides.add_slide(blank)
add_title(s, "Adoption Trend — Last 12 Months")
add_left_bullets(s, [
    "Steady month-over-month gains",
    "Inflection at the May release",
    "Plateau in late summer, expected",
    "Q4 push tied to onboarding revamp",
])
add_right_chart(s, XL_CHART_TYPE.LINE,
                ["J","F","M","A","M","J","J","A","S","O","N","D"],
                [("Monthly active accounts", (12,14,15,18,24,28,30,31,30,33,38,42))],
                title="Monthly active accounts")
add_footer(s, "Template 2 · bullets + line chart")

# Slide 3: Bullets + Pie chart
s = prs.slides.add_slide(blank)
add_title(s, "Revenue Mix by Product Line")
add_left_bullets(s, [
    "Platform remains the anchor",
    "Add-ons crossed 20% mix",
    "Services held steady at 11%",
    "Targeting 25% add-on mix by EOY",
])
add_right_chart(s, XL_CHART_TYPE.PIE,
                ["Platform", "Add-ons", "Services", "Other"],
                [("Mix", (52, 22, 18, 8))],
                title="Mix")
add_footer(s, "Template 3 · bullets + pie chart")

# Slide 4: Bullets + Image
s = prs.slides.add_slide(blank)
add_title(s, "Product Walkthrough")
add_left_bullets(s, [
    "Updated home dashboard",
    "One-click templates panel",
    "New keyboard-driven nav",
    "Inline review and accept",
])
add_right_image(s, "Screenshot / hero image")
add_footer(s, "Template 4 · bullets + image")

# Slide 5: Bullets + Table
s = prs.slides.add_slide(blank)
add_title(s, "Top Accounts — Status")
add_left_bullets(s, [
    "All top-10 are renewal-ready",
    "Two upsell motions in flight",
    "One at-risk with mitigation",
    "Coverage 1.6x on the segment",
])
add_right_table(s,
    headers=["Account", "ARR", "Stage"],
    rows=[
        ["Acme Co.",      "$1.2M",  "Renewal"],
        ["Globex",        "$0.9M",  "Upsell"],
        ["Initech",       "$0.7M",  "At risk"],
        ["Umbrella",      "$0.6M",  "Renewal"],
        ["Stark Ind.",    "$0.5M",  "Renewal"],
    ],
)
add_footer(s, "Template 5 · bullets + table")

# Slide 6: Paragraph + Bar chart
s = prs.slides.add_slide(blank)
add_title(s, "Operating Leverage Improving")
add_left_paragraph(s,
    "Operating margin expanded for the fourth consecutive quarter as the "
    "platform-mix shift fed through to gross margin and headcount growth "
    "continued to lag revenue. Q3 marks the first period above the 20% "
    "operating-margin target set at the start of the year.")
add_right_chart(s, XL_CHART_TYPE.COLUMN_CLUSTERED,
                ["Q4 last", "Q1", "Q2", "Q3"],
                [("Op margin (%)", (12, 15, 18, 22))],
                title="Operating margin %")
add_footer(s, "Template 6 · paragraph + column chart")

# Slide 7: Paragraph + Image
s = prs.slides.add_slide(blank)
add_title(s, "New Onboarding Experience")
add_left_paragraph(s,
    "We rebuilt onboarding around a single guided checklist, with the "
    "team's most-used templates surfaced inline. Early users complete the "
    "first action in under three minutes and reach the activation event "
    "twice as fast as the previous flow.")
add_right_image(s, "Product screenshot")
add_footer(s, "Template 7 · paragraph + image")

# Slide 8: Paragraph + Table
s = prs.slides.add_slide(blank)
add_title(s, "Headcount Plan — Next Two Quarters")
add_left_paragraph(s,
    "The plan front-loads engineering hiring to support the Q1 roadmap, "
    "then opens GTM capacity once the new playbook is rolled out. All "
    "roles are open as of the start of the next quarter and recruiting "
    "capacity is in place to fill them at the cadence shown.")
add_right_table(s,
    headers=["Function", "Q4 plan", "Q1 plan"],
    rows=[
        ["Engineering", "+6", "+4"],
        ["Product",     "+2", "+1"],
        ["Sales",       "+3", "+5"],
        ["CS",          "+2", "+3"],
    ],
)
add_footer(s, "Template 8 · paragraph + table")

# Slide 9: Lead + sub-bullets + Image
s = prs.slides.add_slide(blank)
add_title(s, "Why This Launch Matters")
add_left_lead_plus_bullets(s,
    lead="The first release to unify ingest, review, and compose in one path.",
    bullets=[
        "Cuts setup time from an afternoon to ten minutes",
        "Unlocks templated decks for non-design users",
        "Sets the foundation for the v6 redesign",
        "Lands ahead of the annual user conference",
    ])
add_right_image(s, "Launch visual / mockup")
add_footer(s, "Template 9 · lead + sub-bullets + image")

# Slide 10: KPI stats + Column chart
s = prs.slides.add_slide(blank)
add_title(s, "By the Numbers")
add_left_kpis(s, [("+38%", "Year-over-year revenue growth"),
                  ("4.7 / 5", "Customer satisfaction score")])
add_right_chart(s, XL_CHART_TYPE.COLUMN_CLUSTERED,
                ["Q1", "Q2", "Q3", "Q4 (E)"],
                [("Revenue ($M)", (28, 33, 41, 48))],
                title="Quarterly revenue")
add_footer(s, "Template 10 · KPI stats + column chart")

# Save
out = Path(__file__).resolve().parent / "standard_templates.pptx"
out.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(out))
print(f"Wrote {out} ({out.stat().st_size // 1024} KB, {len(prs.slides)} slides)")
