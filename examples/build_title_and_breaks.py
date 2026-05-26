"""Build a 6-slide title/break/closing deck for pptx-skill.

Complements `build_standard_templates.py`. That deck covers content
shapes (header + body + chart/image/table). This one covers the
*structural* slides every real deck needs around them:

  1. Title cover           — deck title, subtitle, author/date
  2. Agenda / TOC          — numbered sections
  3. Section break (text)  — big phrase, accent rule
  4. Section break (number)— giant number with caption
  5. Pull quote            — large quote, attribution
  6. Closing / thank-you   — final headline + call to action

Same brand constants as the standard templates so the two decks
look like they belong to the same family. Ingest treats each slide
as its own skeleton, so the agent gets a vocabulary of opener,
divider, big-stat, quote, and closing shapes alongside the content
templates.
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# ---- Brand (matches build_standard_templates.py) ------------------------
ACCENT = RGBColor(0x1F, 0x4E, 0x8C)
INK = RGBColor(0x22, 0x2A, 0x35)
MUTED = RGBColor(0x6B, 0x73, 0x82)
LIGHT = RGBColor(0xF1, 0xF4, 0xF8)
HAIR = RGBColor(0xD8, 0xDD, 0xE5)

TITLE_FONT = "Calibri"
BODY_FONT = "Calibri"

# ---- Helpers ------------------------------------------------------------

def _run(p, text: str, *, font=TITLE_FONT, size=24, bold=False,
         color=INK, align=None):
    if align is not None:
        p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return r


def _text_box(slide, left, top, width, height, text: str, **kw):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    _run(tf.paragraphs[0], text, **kw)
    return tb


def _accent_rule(slide, left, top, width, height_emu=28575):
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Emu(height_emu)
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = ACCENT
    rule.line.fill.background()
    return rule


def add_footer(slide, text: str):
    tb = slide.shapes.add_textbox(
        Inches(0.5), Inches(7.0), Inches(12.33), Inches(0.3)
    )
    p = tb.text_frame.paragraphs[0]
    _run(p, text, font=BODY_FONT, size=9, color=MUTED)


# ---- Build deck ---------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
CENTER_LEFT_W = Inches(11)
CENTER_LEFT = (SLIDE_W - CENTER_LEFT_W) / 2

# Slide 1 — Title cover
s = prs.slides.add_slide(blank)
# Accent block on left edge as a hero detail
side = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), SLIDE_H)
side.fill.solid()
side.fill.fore_color.rgb = ACCENT
side.line.fill.background()
# Eyebrow label
_text_box(
    s, Inches(1.0), Inches(2.0), Inches(11.0), Inches(0.45),
    "ANNUAL REVIEW", font=BODY_FONT, size=12, bold=True, color=ACCENT,
)
# Title
_text_box(
    s, Inches(1.0), Inches(2.55), Inches(11.0), Inches(1.6),
    "FY26 Business Update",
    font=TITLE_FONT, size=54, bold=True, color=INK,
)
# Subtitle
_text_box(
    s, Inches(1.0), Inches(4.05), Inches(11.0), Inches(1.0),
    "Markets, momentum, and the path through next year",
    font=BODY_FONT, size=22, color=MUTED,
)
# Author / date block
_text_box(
    s, Inches(1.0), Inches(5.55), Inches(11.0), Inches(0.5),
    "Jakub Naskretski",
    font=BODY_FONT, size=14, bold=True, color=INK,
)
_text_box(
    s, Inches(1.0), Inches(6.05), Inches(11.0), Inches(0.45),
    "Q4 board review · May 2026",
    font=BODY_FONT, size=12, color=MUTED,
)
add_footer(s, "Template · title cover")

# Slide 2 — Agenda
s = prs.slides.add_slide(blank)
_text_box(
    s, Inches(0.5), Inches(0.55), Inches(12.33), Inches(0.7),
    "Agenda", font=TITLE_FONT, size=32, bold=True, color=INK,
)
_accent_rule(s, Inches(0.5), Inches(1.25), Inches(0.7))
agenda_items = [
    ("01", "Where we are"),
    ("02", "What changed this year"),
    ("03", "The product roadmap"),
    ("04", "Investment priorities"),
    ("05", "Next 90 days"),
]
top = Inches(1.95)
for idx, (num, label) in enumerate(agenda_items):
    row_top = top + Inches(0.95 * idx)
    # Numeral
    tb = s.shapes.add_textbox(Inches(0.5), row_top, Inches(1.6), Inches(0.85))
    _run(
        tb.text_frame.paragraphs[0], num,
        font=TITLE_FONT, size=40, bold=True, color=ACCENT,
    )
    # Label
    tb = s.shapes.add_textbox(Inches(2.0), row_top + Inches(0.18),
                              Inches(10.5), Inches(0.6))
    _run(
        tb.text_frame.paragraphs[0], label,
        font=BODY_FONT, size=22, color=INK,
    )
    # Divider hair-line between rows (skip after last)
    if idx < len(agenda_items) - 1:
        line = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), row_top + Inches(0.9),
            Inches(12.33), Emu(9525),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = HAIR
        line.line.fill.background()
add_footer(s, "Template · agenda / TOC")

# Slide 3 — Section break with big text
s = prs.slides.add_slide(blank)
# Eyebrow (section label)
_text_box(
    s, CENTER_LEFT, Inches(2.4), CENTER_LEFT_W, Inches(0.45),
    "PART TWO", font=BODY_FONT, size=14, bold=True, color=ACCENT,
    align=PP_ALIGN.CENTER,
)
# Big phrase
_text_box(
    s, CENTER_LEFT, Inches(3.05), CENTER_LEFT_W, Inches(1.8),
    "What changed this year",
    font=TITLE_FONT, size=60, bold=True, color=INK,
    align=PP_ALIGN.CENTER,
)
# Subtle accent under
rule = s.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    (SLIDE_W - Inches(1.4)) / 2, Inches(4.85),
    Inches(1.4), Emu(28575),
)
rule.fill.solid()
rule.fill.fore_color.rgb = ACCENT
rule.line.fill.background()
add_footer(s, "Template · section break (big text)")

# Slide 4 — Section break with big number
s = prs.slides.add_slide(blank)
# Giant numeral
_text_box(
    s, CENTER_LEFT, Inches(1.7), CENTER_LEFT_W, Inches(3.2),
    "03",
    font=TITLE_FONT, size=180, bold=True, color=ACCENT,
    align=PP_ALIGN.CENTER,
)
# Section title
_text_box(
    s, CENTER_LEFT, Inches(4.85), CENTER_LEFT_W, Inches(0.9),
    "The product roadmap",
    font=TITLE_FONT, size=36, bold=True, color=INK,
    align=PP_ALIGN.CENTER,
)
# Supporting caption
_text_box(
    s, CENTER_LEFT, Inches(5.75), CENTER_LEFT_W, Inches(0.6),
    "Where the next four releases take us",
    font=BODY_FONT, size=18, color=MUTED,
    align=PP_ALIGN.CENTER,
)
add_footer(s, "Template · section break (big number)")

# Slide 5 — Pull quote
s = prs.slides.add_slide(blank)
# Big opening quote glyph as a separate run for emphasis
tb = s.shapes.add_textbox(
    Inches(0.9), Inches(1.6), Inches(11.5), Inches(0.9)
)
_run(
    tb.text_frame.paragraphs[0], "“",
    font=TITLE_FONT, size=120, bold=True, color=ACCENT,
)
# The quote itself
_text_box(
    s, Inches(1.5), Inches(2.85), Inches(10.3), Inches(2.6),
    "We stopped building features to fill a roadmap and started "
    "building the three things customers asked for every quarter.",
    font=TITLE_FONT, size=28, color=INK,
)
# Attribution
_text_box(
    s, Inches(1.5), Inches(5.55), Inches(10.3), Inches(0.5),
    "— Internal product review, March 2026",
    font=BODY_FONT, size=14, color=MUTED,
)
add_footer(s, "Template · pull quote")

# Slide 6 — Closing / thank you
s = prs.slides.add_slide(blank)
# Eyebrow
_text_box(
    s, CENTER_LEFT, Inches(2.3), CENTER_LEFT_W, Inches(0.5),
    "QUESTIONS WELCOME", font=BODY_FONT, size=12, bold=True, color=ACCENT,
    align=PP_ALIGN.CENTER,
)
# Big headline
_text_box(
    s, CENTER_LEFT, Inches(2.85), CENTER_LEFT_W, Inches(1.6),
    "Thank you",
    font=TITLE_FONT, size=64, bold=True, color=INK,
    align=PP_ALIGN.CENTER,
)
# Call to action
_text_box(
    s, CENTER_LEFT, Inches(4.6), CENTER_LEFT_W, Inches(0.6),
    "Reach out — happy to dig deeper on any of this",
    font=BODY_FONT, size=20, color=MUTED,
    align=PP_ALIGN.CENTER,
)
# Contact block (single line, centred)
_text_box(
    s, CENTER_LEFT, Inches(5.55), CENTER_LEFT_W, Inches(0.5),
    "jakub@example.com · linkedin.com/in/jakubnaskretski",
    font=BODY_FONT, size=14, color=INK,
    align=PP_ALIGN.CENTER,
)
add_footer(s, "Template · closing / thank you")

# Save
out = Path(__file__).resolve().parent / "title_and_breaks.pptx"
out.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(out))
print(f"Wrote {out} ({out.stat().st_size // 1024} KB, {len(prs.slides)} slides)")
